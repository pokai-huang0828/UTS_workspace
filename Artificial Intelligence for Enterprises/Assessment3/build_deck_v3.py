# -*- coding: utf-8 -*-
"""A3 v3 投影片產生器 —— 依 notes/A3_提案內容_v3.md。
論證:先證明,再建置。復用 build_deck.py 的版面 helper(16:9 / 右下淨空區 / 三色制)。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from build_deck import (txt, header, picture, note, card,
                        NAVY, RED, GREY, DARK, WHITE, W, H, TOP, SAFE_R, FIG)

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Huang_26254793_421104_Assessment 3.pptx")

# v3 的五段導軌:介紹 / AI 機會 / 業務影響 / 效益與衡量 / 項目計畫


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    blank = prs.slide_layouts[6]
    S = lambda: prs.slides.add_slide(blank)
    gap = 0.16

    # ---------- 封面 ----------
    s = S()
    txt(s, 1.0, 2.0, 11.3, 1.5,
        "房貸固定利率到期留存:\n先證明,再建置", size=38, color=NAVY, bold=True, spacing=1.25)
    txt(s, 1.0, 3.62, 11.3, 0.6,
        "AI 路線圖提案  ·  提交技術審查委員會", size=20, color=GREY)
    txt(s, 1.0, 4.62, 11.3, 1.4,
        "Po-Kai Huang(26254793)\n421104 Artificial Intelligence for Enterprises\n"
        "Assessment 3  ·  2026 年 8 月", size=15, color=DARK, spacing=1.5)

    # ---------- P1 Executive Summary ----------
    s = S(); header(s, "Executive Summary — Decision Requested", 0)
    cw = 2.62
    cards = [
        ("批准什麼", "不是一套 AI 系統。\n\n是一個 8 週的\ndiscovery 與對照實驗,\n\n用來回答:\n房貸到期主動介入\n到底有沒有增量效果", RED),
        ("要多少", "約 3.5 人月\n(規劃值)\n\n無新增系統\n無新增平台授權\n無外部顧問", NAVY),
        ("為什麼先做這個", "到期日是確定事件,\n找出客戶不需要 AI。\n\n沒人知道的是\n介入有沒有用。\n\n在證明之前建模型,\n是把最貴的一步\n放在最前面", NAVY),
        ("何時要決定", "本次委員會\n\n8 週後帶著\n實驗結果回來,\n\n那時才決定\n要不要投第二階段", NAVY),
    ]
    for i, (t, b, c) in enumerate(cards):
        card(s, 0.55 + i * (cw + gap), TOP + 0.30, cw, 4.30, t, b, c)
    txt(s, 0.55, TOP + 4.85, W - 1.1 - SAFE_R, 0.9,
        "向一家在 Evident AI Index 2025 全球第 4、亞太第 1 的機構提案(CBA, 2025a)——\n"
        "它不缺 AI 能力。它缺的是這條業務線上的一個答案。",
        size=14, color=NAVY, bold=True, spacing=1.4)

    # ---------- P2 起點 ----------
    s = S(); header(s, "起點:CBA 已經在做的事", 0)
    picture(s, "fig_v3_01_baseline.png", top=TOP + 0.10, maxh=5.05)
    note(s, "主動說出對自己不利的事實,是投資提案裡最強的信任機制。", color=NAVY, size=12)

    # ---------- P3 缺口 ----------
    s = S(); header(s, "缺口:房貸到期不等於定存到期", 1)
    picture(s, "fig_v3_02_gap.png", top=TOP + 0.10, maxh=5.05)
    note(s, "CDR 降低了比價與資料取得的摩擦,但轉貸本身仍是一個完整的信貸流程。", size=11)

    # ---------- P4 六個未知 ----------
    s = S(); header(s, "六個未知:沒有一個能從公開資料回答", 1)
    picture(s, "fig_v3_03_unknowns.png", top=TOP + 0.05, maxh=5.15)
    note(s, "這六件事沒有一件我能從外面查到 —— 所以我請求的,是先把這六個答案買回來。",
         color=NAVY, size=12)

    # ---------- P5 三條路 ----------
    s = S(); header(s, "三條路,包含不用 AI 的那條", 1)
    picture(s, "fig_v3_04_options.png", top=TOP + 0.05, maxh=5.15)

    # ---------- P6 提案 ----------
    s = S(); header(s, "提案:8 週 discovery + 規則式對照實驗", 2)
    picture(s, "fig_v3_05_plan.png", top=TOP + 0.15, maxh=4.55)
    txt(s, 0.55, TOP + 4.80, W - 1.1 - SAFE_R, 0.9,
        "五個工作包的交付物全部可簽收:CEE 覆蓋現況報告、到期分佈實算、\n"
        "流失事件操作型定義(業務與資料雙簽)、合規意見書、實驗結果報告、回委員會的建議書。",
        size=12, color=DARK, spacing=1.4)

    # ---------- P7 實驗設計 ----------
    s = S(); header(s, "實驗設計:怎麼證明介入有效", 3)
    picture(s, "fig_v3_06_experiment.png", top=TOP + 0.05, maxh=5.15)

    # ---------- P8 風險 ----------
    s = S(); header(s, "限制與風險", 2)
    picture(s, "fig_v3_07_risk.png", top=TOP + 0.10, maxh=5.05)
    note(s, "「結論是不值得做」不是失敗 —— 是用 3.5 人月避免了 27.5 人月的錯誤投資。",
         color=NAVY, size=12)

    # ---------- P9 倫理 ----------
    s = S(); header(s, "社會文化與倫理影響", 2)
    picture(s, "fig_v3_08_ethics.png", top=TOP + 0.10, maxh=5.05)
    note(s, "這個系統最根本的倫理問題,不在它用了什麼特徵,在它決定先照顧誰。",
         color=RED, size=12)

    # ---------- P10 決策規則 ----------
    s = S(); header(s, "決策規則:什麼結果導向什麼下一步", 3)
    picture(s, "fig_v3_09_decision.png", top=TOP + 0.05, maxh=5.15)

    # ---------- P11 資源 ----------
    s = S(); header(s, "時程、資源與當責", 4)
    picture(s, "fig_v3_10_resource.png", top=TOP + 0.10, maxh=3.55, maxw=7.2, left=0.55)
    card(s, 8.05, TOP + 0.10, 3.15, 3.55, "關鍵利益相關者",
         "專案贊助人\n(零售銀行主管)\n\n業務負責人(房貸)\n\n資料負責人\n\n合規法務\n\n"
         "客戶經理代表\n\n財務(產品別 margin)", col=NAVY, bsize=10)
    txt(s, 0.55, TOP + 3.90, W - 1.1 - SAFE_R, 1.7,
        "關鍵路徑是「合規預審 → 對照實驗」—— 合規意見書未出,實驗不能開始。\n"
        "這是本案唯一的外部依賴,也是唯一需要簽核才能放行的關卡。",
        size=13, color=NAVY, spacing=1.45)

    # ---------- P12 結論 ----------
    s = S(); header(s, "結論與批准請求", None)
    cards2 = [
        ("批准什麼", "8 週 discovery\n+ 規則式對照實驗\n\n約 3.5 人月", NAVY),
        ("為什麼是這個", "到期日是確定事件,\n找出客戶不需要 AI。\n\n沒人知道介入有沒有用,\n而那是地基", NAVY),
        ("會拿到什麼", "六個未知的答案\n\n一個帶信賴區間的\n增量效果估計\n\n一份走或不走的建議書", NAVY),
        ("已知限制", "CEE 覆蓋範圍\n待內部確認\n\n實驗可能受執行率影響\n\n結論可能是\n「不值得做」", RED),
    ]
    for i, (t, b, c) in enumerate(cards2):
        card(s, 0.55 + i * (cw + gap), TOP + 0.35, cw, 3.65, t, b, c)
    txt(s, 0.55, TOP + 4.30, W - 1.1 - SAFE_R, 1.3,
        "我今天不是來要一套系統的預算。我是來要一個答案的預算。\n"
        "八週之後,我會帶著這個答案回來 —— 那時候各位要批的,才是一個真正的投資案。",
        size=16, color=NAVY, bold=True, spacing=1.45)

    # ---------- 參考文獻 ----------
    s = S(); header(s, "參考文獻", None)
    txt(s, 0.55, TOP + 0.25, W - 1.1, 5.7,
        "本頁依 notes/A3_參考文獻_查證版.md 填入 —— 該檔的每一筆都經 WebFetch 實地查證,\n"
        "查不到出處的一律標明,未採用。\n\n"
        "已查證可用(七筆):CBA 2019 分析師簡報(CEE 定存到期留存 85%)、CBA 2022 CEE 說明、\n"
        "CBA 2024 GenAI 與客服、CBA 2025 詐騙損失、CBA FY25 業績、\n"
        "Devriendt, Berrevoets & Verbeke (2021) Information Sciences(uplift modeling,有 DOI)、\n"
        "澳洲 CDR 官方資料。\n\n"
        "⚠️ 已排除:A$611.5bn(查無精確出處,改用 A$634bn);市占 25.4% 標明對應 1H25 非 FY25。\n\n"
        "三條機械判準:每筆可點擊深層連結 / 零重複條目 / 零孤兒條目。APA 7,字級不小於 12pt。",
        size=13, color=DARK, spacing=1.5)

    prs.save(OUT)
    return prs, OUT


if __name__ == "__main__":
    prs, out = build()
    n = len(prs.slides._sldIdLst)
    print("投影片頁數:", n)
    print("輸出:", os.path.basename(out))
    print("大小: %.1f KB" % (os.path.getsize(out) / 1024))
