# -*- coding: utf-8 -*-
"""掃描全片是否殘留作廢值 —— 規格、講稿骨架、交付檔(含講者備註)三處都掃。

🔴 為什麼要有這支:數字改過好幾輪,而**改一個地方就會在別處留下矛盾**。
   2026-08-16 實測抓到:第 3 頁的 why 欄還寫著 43%–78%(已作廢,現值 50%–87.5%),
   而 why 欄**會進講者備註**(【這一頁的作用】)—— 也就是錄影時眼睛看得到。
   靠人記是記不住的,所以做成一條指令。

用法:python scripts/sweep_stale.py
"""
import json, re, sys
P = json.load(open('notes/deck_spec/pages.json', encoding='utf-8'))
STALE = [("43%–78%", r"43%[–-]78%"), ("2.5–5.5 天", r"2\.5[–-]5\.5"),
         ("22%–57%", r"22%[–-]57%"), ("NT$(匯率宣告除外)", r"NT\$(?!22\.6)"),
         ("60%–64%", r"60%[–-]64%"), ("16–23 人", r"16[–-]23 人"),
         ("4.4 分鐘", r"4\.4 分"), ("1,144–2,106", r"1,144[–-]2,106"),
         ("56.9–113.8", r"56\.9[–-]113\.8"), ("2.2–3.5 小時", r"2\.2[–-]3\.5 小時")]

# 🔑 白名單 —— **刻意提到作廢值**的地方,而且理由必須寫在這裡。
#    純字串黑名單分不出「拿它當主張」跟「解釋我們為什麼不拿它當主張」,
#    但後者在答辯裡是必要的誠實。所以放行採**逐字全句比對**,不是關鍵字放行:
#    句子只要改一個字就會重新被抓到,不會變成一張擴散的免死金牌。
ALLOW = [
    # 第 3 頁備註:坦承 0.25/0.25 是比營運端原口述(0.5/0.5)收緊過的值。
    # 不寫出 43%–78%,委員會就不知道收緊幅度有多大 —— 那就不叫坦承。
    "(用 0.5/0.5 算是 43%–78%)",
]


def strip_allowed(t):
    for a in ALLOW:
        t = t.replace(a, "")
    return t


def shape_text(sl):
    """🔴 表格文字不在 text_frame 裡 —— 只走 text_frame 會漏掉整張表。
       同一個 bug 讓答辯題庫的 agent 斷言「片上沒有 X」,而 X 就印在表格上。
       見 scripts/dump_deck.py 的事故紀錄。"""
    out = []
    for sh in sl.shapes:
        if sh.has_table:
            out += [c.text for r in sh.table.rows for c in r.cells]
        elif sh.has_text_frame:
            out.append(sh.text_frame.text)
    return " ".join(out)
bad = []
for pg in P:
    s = strip_allowed(json.dumps(pg, ensure_ascii=False))
    for name, pat in STALE:
        for m in re.finditer(pat, s):
            bad.append((pg['n'], name, s[max(0, m.start()-40):m.start()+30]))
print("=== pages.json 殘留作廢值 ===")
if bad:
    for n, name, ctx in bad:
        print(f"  🔴 第 {n} 頁 [{name}] …{ctx}…")
else:
    print("  ✅ 無")

sys.path.insert(0, 'scripts')
import pitch
print("\n=== pitch.py(講稿骨架)殘留 ===")
b2 = []
for i, sk in enumerate(pitch.SKELETON):
    blob = strip_allowed(json.dumps(sk, ensure_ascii=False))
    for name, pat in STALE:
        if re.search(pat, blob):
            b2.append((i+1, name))
print("  ✅ 無" if not b2 else "\n".join(f"  🔴 第 {n} 頁 [{x}]" for n, x in b2))

from pptx import Presentation
prs = Presentation('Huang_26254793_421104_Assessment 3.pptx')
print("\n=== 交付檔(含講者備註)殘留 ===")
b3 = []
for i, sl in enumerate(prs.slides, 1):
    txt = shape_text(sl) + " " + (sl.notes_slide.notes_text_frame.text or "")
    txt = strip_allowed(txt)
    for name, pat in STALE:
        if re.search(pat, txt):
            b3.append((i, name))
print("  ✅ 無" if not b3 else "\n".join(f"  🔴 第 {n} 頁 [{x}]" for n, x in b3))
