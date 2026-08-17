# -*- coding: utf-8 -*-
"""把 notes/deck_spec/pages.json 渲染成全原生形狀的 pptx。

與前兩版最大的不同:**沒有任何一張 PNG**。
每個方塊、每條線、每段文字都是 PowerPoint 原生物件,Kenny 可以直接點進去改。

用法:python scripts/render_deck.py
輸出:Huang_26254793_421104_Assessment 3.pptx
"""
import json
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches

import native as N
import pitch as PITCH

HERE = os.path.dirname(os.path.abspath(__file__))
A3 = os.path.dirname(HERE)
SPEC = os.path.join(A3, "notes", "deck_spec", "pages.json")
OUT = os.path.join(A3, "Huang_26254793_421104_Assessment 3.pptx")

M = 0.72                       # 版心左右
TITLE_Y = 0.42
RULE_Y = 1.52
BODY_TOP = 1.76
BODY_BOT = N.H - 0.38          # 內容底線(0.42→0.38:第 9 頁差 0.12" 才裝得下八條 KPI)
FOOT_Y = N.H - 0.62

COVER = dict(
    # 🔴 2026-08-16 Kenny 逐頁審第 1 頁的定案。
    title="讓機器讀懂事件的完整脈絡",
    # 🔴 「三道決策門」是白話化之前的代號;第 6、10 頁片上已全面換成「三個喊停點」。
    #    封面副標是觀眾看到的第一行說明文字,卻留著舊詞 —— 同字數直接換。
    subtitle="地端 AI 事件判讀 —— 六個月、三個喊停點的提案計畫",
    # 「光是一個車隊」比原本的「今天」有力 —— 它強調 3,000 筆只是**一個**客戶;
    # 但後半句「而客戶要的答案,兩天內」不能省:張力來自量 vs 期限,少了它只剩量大。
    hook="光是一個車隊,五個人要看完三千支影片;而客戶要的答案,兩天內。",
    scale=["我們客戶是車隊,包含物流 · 客運 · 礦業;車上裝我們的鏡頭與車機",
           "鏡頭偵測到一次危險駕駛 → 記一筆「事件」(影片 + 時間地點 + 系統分類)",
           "客戶回頭說「這筆 AI 判錯了」= 一次「爭議」—— 由人重看影片才判得出來",
           "單一大型車隊約 3,000 筆事件/天 · 部門僅 5 位專職分析人員逐筆看,而客戶還在持續增長"],
    who=["Po-Kai Huang(學號 26254793)",
         "421104 Artificial Intelligence for Enterprises · Assessment 3",
         "2026 年 8 月"],
    confid="The company name has been changed for commercial in confidence reasons.　"
           "內部數字均標示為假設或推算量級,並於各頁附推算方式。",
)

# 🔴 2026-08-17 全面重整:12 筆 → 7 筆。
#    Kenny:「第 12 頁針對我前面的內容重金校對,不要有錯標。」
#    投影片改過非常多輪(自動結案整條拿掉、風險頁重寫、倫理頁改成「沒有一條變差」),
#    但參考文獻沒有跟著更新 —— 逐條核對後,**七筆是為了已經不存在的內容而引的**:
#      · NIST AI RMF —— 全片「NIST」只出現兩次,一次是第 7 頁備註**叫自己不要講**
#      · Rudin (2019) —— **方向相反**:它反對高風險決策用黑箱 + 事後解釋,
#        而本案用的就是兩顆黑箱 VLM + 事後說明
#      · Ensign (2018) —— 唯一錨點是第 7 頁「以前判成誤報的紀錄算不算證據」那一列,
#        而那一列 Kenny 已經下令刪掉
#      · Hopp & Spearman —— 引排隊論去撐一段片上明說「不是瓶頸、不在本案範圍」的區段;
#        指定 chaps. 7–9 更是自己遞刀(那你的 Kingman 分析呢)
#      · Northcutt (2021) —— **會反咬**:它證明測試集本身會標錯,
#        而第 7 頁第 2 列的把握寫「高 · 答案是確定的」
#      · Guo (2017) —— 講信心校準,不是幻覺;片上不顯示信心分數的理由是「人會照抄」
#      · AHRC (2020) —— 講受保護群體的差別待遇,片上通篇沒有族群/歧視
#    留下的五筆**書目全部有錯**(篇名被砍掉後半、機關名年份錯置、缺 DOI),已逐筆查證修正。
#    另補兩筆填補零支撐的主張(流暢型幻覺、量化精度損失)—— 那兩個正是評分者的主場。
#    每一筆後面掛頁碼標註,讓評分者一眼看到都真的被用到。
REFS = [
    # 🔴 APA 7 依作者姓氏字母序 —— 不是依重要性排。Department < Dettmers(Dep < Det)。
    # 第 8 頁全頁的骨架。🔴 機關名年份錯置:八原則 2019-11-07 由 DIIS 發布,
    #    DISR 這個名字 2022 年 7 月才存在(等於讓 2022 年成立的機關掛 2019 年的著作)。
    #    篇名也不可擅自展開成 artificial intelligence —— 官方標題就是 AI。
    "Department of Industry, Innovation and Science. (2019). Australia's AI ethics "
    "principles. https://www.industry.gov.au/publications/australias-ai-ethics-principles"
    "　(第 8 頁)",
    # 🔴 新增:第 5 頁「8-bit 損失比 4-bit 小」原本零文獻支撐,而它撐著「為什麼要兩台」。
    #    ⚠️ 這篇一般被稱作 LLM.int8()(arXiv 2208.07339),但 **NeurIPS 論文集裡的標題是
    #    GPT3.int8()** —— 片上照版本紀錄寫,被問到就說明,那反而是個加分的細節。
    "Dettmers, T., Lewis, M., Belkada, Y., & Zettlemoyer, L. (2022). GPT3.int8(): 8-bit "
    "matrix multiplication for transformers at scale. In Advances in neural information "
    "processing systems (Vol. 35).　(第 5 頁)",
    # 第 7 頁第 3 列「機器給了判定結果,人可能照抄」+ 介面層緩解。
    #    原篇名被砍掉後半,拿那個標題去 PubMed 搜不到這篇。
    "Goddard, K., Roudsari, A., & Wyatt, J. C. (2012). Automation bias: A systematic "
    "review of frequency, effect mediators, and mitigators. Journal of the American "
    "Medical Informatics Association, 19(1), 121–127. "
    "https://doi.org/10.1136/amiajnl-2011-000089　(第 7 頁)",
    # 第 9 頁⑨ 盲測「比對照組差就不放行」。這篇獨有的貢獻是
    #    「有人監督本身不構成防護」—— 正是把它列成喊停條件的依據。
    "Parasuraman, R., & Manzey, D. H. (2010). Complacency and bias in human use of "
    "automation: An attentional integration. Human Factors, 52(3), 381–410. "
    "https://doi.org/10.1177/0018720810376055　(第 9 頁)",
    # 🔴 新增:第 7 頁第 1 列「機器會判錯,但講得很順」原本零文獻支撐 ——
    #    那是風險表第一列、整套「標出證據」設計的存在理由,也是評分者的主場。
    "Rohrbach, A., Hendricks, L. A., Burns, K., Darrell, T., & Saenko, K. (2018). "
    "Object hallucination in image captioning. In Proceedings of the 2018 Conference on "
    "Empirical Methods in Natural Language Processing (pp. 4035–4045). "
    "https://doi.org/10.18653/v1/D18-1437　(第 7 頁)",
    # 第 10 頁「順序不能換」+ 第 4 頁「判準不一致的標籤餵進去」。
    "Sambasivan, N., Kapania, S., Highfill, H., Akrong, D., Paritosh, P., & Aroyo, L. M. "
    "(2021). “Everyone wants to do the model work, not the data work”: Data cascades "
    "in high-stakes AI. In Proceedings of the 2021 CHI Conference on Human Factors in "
    "Computing Systems. https://doi.org/10.1145/3411764.3445518　(第 4、10 頁)",
    # 第 11 頁「近半數是顧它的人力」+「五項還沒估價」= 模型以外的成本。
    #    🔴 原本用 APA 省略號吃掉三位作者;這篇共十位,APA 7 要列前 19 位。
    "Sculley, D., Holt, G., Golovin, D., Davydov, E., Phillips, T., Ebner, D., "
    "Chaudhary, V., Young, M., Crespo, J.-F., & Dennison, D. (2015). Hidden technical "
    "debt in machine learning systems. In Advances in neural information processing "
    "systems (Vol. 28).　(第 11 頁)",
]

# 每種 block 的相對高度權重 —— 用來分配版心的垂直空間
WEIGHT = {"cards": 1.35, "matrix": 1.25, "compare": 1.20,
          "bar": 1.15, "rows": 1.00, "callout": 0.34}


def narration():
    """從 notes/_v2_parts 取出十二個單元的口白逐字。

    順序:鉤子 → P1…P10 → 參考文獻,正好對上投影片 1…12。
    舞台指示(整段被括號包住的錄影標記)剝掉 —— 那是給錄影用的,不唸出聲。
    """
    import glob
    H1 = re.compile(r"(?m)^#\s+(.+)$")
    NARR = re.compile(r"(?m)^##\s*口白逐字\s*$")
    NXT = re.compile(r"(?m)^#{1,2}\s+")
    STAGE = re.compile(r"(?ms)^[ \t]*[*_]{0,2}[(（].*?[)）][*_]{0,2}[ \t]*$")
    out = []
    for f in sorted(glob.glob(os.path.join(A3, "notes", "_v2_parts", "*.md"))):
        t = open(f, encoding="utf-8").read()
        hs = list(H1.finditer(t))
        for i, m in enumerate(hs):
            if i == 0:
                continue                      # 分段標題(甲/乙/丙段),不是單元
            e = hs[i + 1].start() if i + 1 < len(hs) else len(t)
            u = t[m.start():e]
            nm = NARR.search(u)
            if not nm:
                out.append((m.group(1).strip(), "", []))
                continue
            rest = u[nm.end():]
            nx = NXT.search(rest)
            raw = rest[:nx.start()] if nx else rest
            # 舞台指示不唸出聲,但它是**錄影提示**(哪裡停、哪裡切全螢幕人像),
            # 所以要留下來放在講稿下方,不是丟掉。
            cues = [re.sub(r"[*_`~]", "", c).strip("()（） \n")
                    for c in STAGE.findall(raw)]
            blk = STAGE.sub("", raw)
            body = "\n".join(l for l in blk.split("\n")
                             if not l.lstrip().startswith(">")).strip()
            out.append((m.group(1).strip(), re.sub(r"[*_`~]", "", body), cues))
    return out


NARR_CACHE = None


def unit(idx):
    """第 idx 頁的 (單元名, 逐字, 錄影提示)。

    🔑 pitch.py 的 say 欄可以覆寫逐字 —— 第 1 頁的開場口白只寫在那裡,
       不在 notes/_v2_parts(那批檔案是十二個單元的原始口白,不含新補的開場)。

    🔴 2026-08-16:投影片變 13 頁,但 notes/_v2_parts 仍是 12 個單元 ——
       新加的「選型與部署」沒有對應單元。它在 pitch.py 標 narr=None,
       而它**後面每一頁的口白索引都要往前挪一格**,否則整份講稿會錯位一頁。
    """
    global NARR_CACHE
    if NARR_CACHE is None:
        NARR_CACHE = narration()
    sk = PITCH.SKELETON[idx] if idx < len(PITCH.SKELETON) else None
    if sk is not None and "narr" in sk and sk["narr"] is None:
        return sk.get("title", ""), sk.get("say", ""), []
    # 🔴 頁序改過好幾輪(插入選型頁、刪掉漏放頁),索引不能再用「猜的」——
    #    骨架裡標了 narr=N 就用 N,那是唯一可靠的對應。
    if sk is not None and isinstance(sk.get("narr"), int):
        ni = sk["narr"]
    else:
        skip = sum(1 for j in range(idx)
                   if j < len(PITCH.SKELETON)
                   and PITCH.SKELETON[j].get("narr", "?") is None)
        ni = idx - skip
    title, body, cues = NARR_CACHE[ni] if ni < len(NARR_CACHE) else ("", "", [])
    if sk and sk.get("say"):
        body = sk["say"]
    # 🔴 單元名來自 notes/_v2_parts,那批檔案停在舊版結構。
    #    第 6 頁片上已改成「三條路」,備註抬頭卻還印著「P5 · 五條路,含三條不用 AI」——
    #    而那一行是他上台前看的第一行字。骨架寫了 title 就以骨架為準。
    if sk and sk.get("title"):
        title = sk["title"]
    return title, body, cues


RATE = 4.4          # 字/秒(全片配時基準)
CJK = re.compile(r"[一-鿿A-Za-z0-9]")


def put_notes(slide, idx, extra_lines):
    """講者備註 = **提詞稿**,兩層。

    🔴 2026-08-15 改版。Jiwei Guan 的作業說明寫著:
       「视频不是简单地"念 PPT"，而是一个 Selling Pitch」
       「如果完全照读 AI 生成的稿件，往往会显得不自然」
       —— 而先前的備註就是一份 2,464 字的逐字稿,照唸正好踩中他點名的那條。
       逐字稿還有個更實際的問題:眼睛在稿子裡找位置的時候,人就不看鏡頭了。

    所以第一層改成**骨架**(自己的話講),逐字稿降為第二層(忘詞才看):
        【第 N 頁 · 目標 XX 秒 · 累計 X:XX】
        ★ 一句話  ← 這頁只有一句會被記住,講完它才換頁
        · 關鍵詞  ← 撐起那句話的三個支點
        # 數字    ← 講錯就毀掉可信度的,照著唸
        ⚠️ 紅線   ← 這頁特有的禁語與必附限定語
        ── 錄影提示 ──
        → 換下一頁
        ── 逐字備援(忘詞才看,不是拿來唸的)──
        ── 備查(答辯用)──
    """
    global NARR_CACHE
    if NARR_CACHE is None:
        NARR_CACHE = narration()
    parts = []
    if unit(idx)[1]:      # 頁數已與 NARR_CACHE 脫鉤,不能再用長度判斷
        title, body, cues = unit(idx)
        n = len(CJK.findall(body))
        secs = n / RATE
        cum = sum(len(CJK.findall(unit(k)[1])) for k in range(idx + 1)) / RATE
        parts.append(f"【第 {idx + 1} 頁 · {title}】"
                     f"　目標 {secs:.0f} 秒 · 累計 {int(cum // 60)}:{int(cum % 60):02d}")

        sk = PITCH.SKELETON[idx] if idx < len(PITCH.SKELETON) else None
        if sk:
            parts.append("")
            parts.append(f"★ {sk['one']}")
            if sk.get("keys"):
                parts.append("")
                parts.extend(f"· {k}" for k in sk["keys"])
            if sk.get("nums"):
                parts.append("")
                parts.extend(f"# {x}" for x in sk["nums"])
            if sk.get("red"):
                parts.append("")
                parts.extend(f"⚠️ {x}" for x in sk["red"])
            if sk.get("cut"):
                parts.append("")
                parts.append("✂️ " + sk["cut"])
            if sk.get("cue") and not any(sk["cue"][:12] in c for c in cues):
                cues = list(cues) + [sk["cue"]]   # 口白檔裡可能已經寫過同一條,別重複
        if cues:
            parts.append("")
            parts.append("── 錄影提示(不唸)──")
            parts.extend(f"· {c}" for c in cues)
        parts.append("")
        parts.append("→ 換下一頁")
        parts.append("")
        parts.append(f"── 逐字備援 · {n} 字(忘詞才看,不是拿來唸的)──")
        parts.append(body)
        if extra_lines:
            parts.append("")
            parts.append("── 備查(不唸,答辯用)──")
    parts.extend(extra_lines)
    if parts:
        slide.notes_slide.notes_text_frame.text = chr(10).join(parts)


def need_h(kind, sp, w):
    """這個版塊**真正需要**多少高度(吋)。

    🔴 為什麼要有這個函式:
       版心高度原本按 WEIGHT 這張跟內容無關的固定權重分配,
       於是「三張卡片、每張三行長句」跟「三張卡片、每張一個詞」拿到一樣的高度。
       實際後果是第 6 頁那張〈這 85% 到底是什麼(三層,一層都不可省)〉——
       標題印出來了,**三層本身被 clamp 吃掉**,片上只剩一句半。
       那是全案最重要的一句話,而它是被一個猜出來的高度砍掉的。

    這裡的估算刻意跟各 blk_* 的實際幾何對齊(卡片標題 0.62 + 徽章 0.42、
    表格表頭 0.36、rows 的 PAD 0.10 …),數值不必精準到吋,
    但**必須隨內容單調成長** —— 分配是按比例正規化的,重要的是相對大小。
    """
    if kind == "callout":
        return N.text_h(sp.get("text", ""), w - 0.40, 17) + 0.30

    if kind == "matrix":
        rows, cols = sp.get("rows", []), sp.get("cols", [])
        name_w = w * sp.get("name_frac", 0.30)
        cw = (w - name_w) / max(len(cols), 1)
        # 每格允許到 2 行時的高度 —— 壓到 1 行雖然放得下,但那是降級不是需求
        return 0.36 + sum(
            max([min(N.text_h(str(r["name"]), name_w - 0.16, 12,
                              safety=N.PINNED_SAFETY), 2 * 0.345)]
                + [min(N.text_h(str(c), cw - 0.16, 12,
                                safety=N.PINNED_SAFETY), 2 * 0.345)
                   for c in r["cells"]]) + 0.06
            for r in rows)

    if kind == "cards":
        items = sp["items"]
        cw = (w - 0.16 * (len(items) - 1)) / max(len(items), 1)
        return max(0.82 + (0.42 if it.get("badge") else 0) + 0.14
                   + sum(N.text_h(ln, cw - 0.24, 12, 1.35)
                         for ln in it.get("lines", []))
                   for it in items)

    if kind == "compare":
        cw = (w - 0.30) / 2
        return 0.90 + (0.48 if sp.get("verdict") else 0) + max(
            sum(N.text_h(ln, cw - 0.32, 13, 1.5) for ln in sp[s]["lines"])
            for s in ("left", "right"))

    if kind == "rows":
        # 🔴 欄寬必須跟 blk_rows **算得一模一樣**,否則估出來的高度是假的。
        #    舊版寫死 VW = w*0.30、NW = w*0.40,但 blk_rows 的值欄是**跟著內容走**的
        #    (只放「第 20–26 週」就只吃 1.58",不是 3.57")——
        #    於是說明欄實際有 7.46" 卻被當成 4.76",行數估多、整塊高度虛胖 0.5"+。
        #    後果是第 10 頁 rows 分到 3.74" 只畫得出 3.20",中間空 0.67",
        #    而「用不完就還給表格」那條規則因為看的是這個假高度,永遠不會觸發。
        items = sp["items"]
        LW = w * 0.24
        vmax = max([N.text_w(str(it["value"]), 15) for it in items] or [0]) + 0.24
        VW = max(min(w * 0.30, vmax), w * 0.10)
        NW = w - LW - VW
        return sum(max(N.text_h(it["label"], LW - 0.24, 14),
                       N.text_h(str(it["value"]), VW - 0.14, 15),
                       min(N.text_h(it.get("note", ""), NW - 0.14, 12),
                           3 * 0.30)) + 0.20
                   for it in items)

    if kind == "bar":
        h = (0.48 if sp.get("label") else 0) + 0.72
        if sp.get("bracket"):
            h += 1.24 if isinstance(sp["bracket"].get("big"), str) else 0.62
        return h + (0.30 if sp.get("note") else 0)

    return 2.0


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    N.box(s, 0, 0, 0.34, N.H, fill=N.NAVY, line=None, radius=False, who="cover-bar")
    N.text(s, 1.05, 1.30, N.W - 2.6, 1.15, COVER["title"], 40, N.NAVY, bold=True,
           who="cover-title")
    N.text(s, 1.05, 2.58, N.W - 2.6, 0.55, COVER["subtitle"], 20, N.INK,
           who="cover-sub")
    N.rule(s, 1.05, 3.30, 3.2, color=N.ORANGE)
    N.text(s, 1.05, 3.60, N.W - 2.6, 0.6, COVER["hook"], 18, N.ORANGE, bold=True,
           who="cover-hook")
    N.text(s, 1.05, 4.26, N.W - 2.6, 1.16, "\n".join(COVER["scale"]), 13, N.GREY,
           spacing=1.30, who="cover-scale")
    N.text(s, 1.05, 5.56, N.W - 2.6 - N.SAFE_R, 0.92, "\n".join(COVER["who"]),
           14, N.INK, spacing=1.30, who="cover-who")
    N.text(s, 1.05, N.H - 0.70, N.SAFE_X - 1.10, 0.5, COVER["confid"], 12, N.GREY,
           spacing=1.2, who="cover-conf")
    put_notes(s, 0, [])
    return s


def refs_page(prs, idx=11):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    N.text(s, M, TITLE_Y + 0.28, 6.0, 0.6, "參考文獻", 28, N.NAVY, bold=True,
           who="ref-title")
    N.rule(s, M, RULE_Y, N.W - 2 * M)
    # 文獻列會排到版面底部,兩欄都必須收在人像淨空區左緣以內
    half = (N.SAFE_X - M - 0.45) / 2
    # 🔴 行高必須**量出來**,不能寫死。
    #    舊版用「每筆 0.80"、每欄 6 筆」的固定格子 —— 那是在條目只有兩三行時成立的。
    #    2026-08-17 補上 DOI 之後每筆變成 4–6 行,固定格子讓七筆**整片疊在一起**,
    #    最後一筆還壓到頁尾那行 APA 聲明上。COM 開得起來、python 也不會叫,
    #    但那一頁的字是糊的。
    # 🔑 改成:先量每一筆要多高,再由高度決定換欄,而不是由筆數。
    avail = (N.H - 0.62) - BODY_TOP          # 版心底到頁尾聲明之間
    # 🔴 行距 1.20 是實測出來的:1.25 時第四筆會卡在兩欄之間(左欄 5.20" > 5.12")。
    REF_SP = 1.20
    hs = [N.text_h(r, half, 12, REF_SP) + 0.10 for r in REFS]
    col, y = 0, BODY_TOP
    for i, (r, h) in enumerate(zip(REFS, hs)):
        if y + h > BODY_TOP + avail + 1e-6 and col == 0:
            col, y = 1, BODY_TOP             # 第一欄裝滿 → 換第二欄
        N.text(s, M + col * (half + 0.45), y, half, h - 0.10,
               r, 12, N.INK, spacing=REF_SP, who=f"ref{i}")
        y += h
        if y > BODY_TOP + avail + 1e-6 and col == 1:
            print(f"   🔴 參考文獻第 {i + 1} 筆之後放不下 —— 需要縮字級或減筆數")
    N.text(s, M, N.H - 0.52, N.SAFE_X - M, 0.34,
           "APA 7th　·　The company name has been changed for commercial in "
           "confidence reasons.", 12, N.GREY, who="ref-foot")
    put_notes(s, idx, [])
    return s


def content(prs, pg):
    N.PAGE = pg["n"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if pg.get("eyebrow"):
        N.text(s, M, TITLE_Y, 6.0, 0.30, pg["eyebrow"], 13, N.ORANGE, bold=True,
               who="eyebrow")
    # 🔴 標題框只有 0.62" 高 = 一行 27pt。放不下就會**折到第二行,壓在副標上** ——
    #    第 6 頁實測撞出「…能不能反／悔」蓋住副標,而所有既有檢查(overflow、鎖定、
    #    作廢值、COM 開檔)沒有一支會叫。標題是全片字最大的一行,撞版一眼就看得到。
    _tw = N.text_w(pg["title"], 27)
    _tmax = N.W - M - 0.55
    if _tw > _tmax:
        print(f'   🔴 第 {pg["n"]} 頁標題放不下(需 {_tw:.2f}" · 只有 {_tmax:.2f}")'
              f' —— 會折行壓到副標:{pg["title"]}')
    N.text(s, M, TITLE_Y + 0.32, _tmax, 0.62, pg["title"], 27, N.NAVY,
           bold=True, who="title")
    N.text(s, M, TITLE_Y + 1.00, N.W - M - 0.55, 0.40, pg["subtitle"], 15,
           N.GREY, who="subtitle")
    N.rule(s, M, RULE_Y + 0.14, N.W - 2 * M)

    blocks = pg["blocks"]
    foot = pg.get("foot")
    bot = (FOOT_Y - 0.14) if foot else BODY_BOT
    avail = bot - (BODY_TOP + 0.16)
    gap = 0.14
    usable = avail - gap * (len(blocks) - 1)

    # 🔑 高度按**實際需要**分配,不用固定權重。
    #    固定權重的問題:某個 block 被分層縮短後,多出來的空間就空在那裡,
    #    而隔壁塞不下的 block 卻還在把內容往備註丟。
    #    callout 有天花板(它是一句話,給再多空間也只用一行高)。
    specs = [json.loads(b["spec"]) if isinstance(b["spec"], str) else b["spec"]
             for b in blocks]
    CAP = {"callout": 1.35}
    needs = [min(need_h(b["kind"], sp, N.W - 2 * M), CAP.get(b["kind"], 1e9))
             for b, sp in zip(blocks, specs)]

    # 🔴 表格先拿到**保證下限**,剩下的空間才按需求比例分。
    #    理由是兩種版塊掉東西的後果不一樣:
    #      · cards / rows 少印一行,是說明變短 —— 難看,但沒有自相矛盾;
    #      · matrix 少印一列,是**標題承諾「十條風險」而片上只有六條** ——
    #        那是投影片自己跟自己打架,而且看的人一數就知道。
    #    所以表格的最小高度(每格壓到 1 行)不參與競爭,直接扣下來。
    #    callout 也給一條下限 —— 第 9 頁那條被分到 **0.00"**,等於這個版塊不存在。
    #    🔴 2026-08-16:**rows 也要有下限**,理由跟表格一模一樣。
    #    第 10 頁副標寫「六包只有一包是 AI」而片上只印得出兩包、
    #    第 11 頁標題寫「成本主體是人」而「人力」那一列根本沒印 ——
    #    那不是「說明變短」,是投影片自己承諾的條數對不上自己印出來的條數。
    #    0.34" 是 blk_rows 的最小列高,0.10" 是列距。
    def _pre_min(kind):
        return {"callout": 0.36, "cards": 2.32, "rows": 0.44,
                "compare": 2.15, "bar": 2.15, "matrix": 0.70}.get(kind, 0.3)

    def _floor(kind, sp):
        if kind == "matrix":
            return 0.36 + len(sp.get("rows", [])) * 0.32
        if kind == "rows":
            n = len(sp.get("items", []))
            return n * 0.34 + 0.10 * max(n - 1, 0)
        return {"callout": 0.62}.get(kind, 0.0)

    floors = [_floor(b["kind"], sp) for b, sp in zip(blocks, specs)]
    if sum(max(floors[i], _pre_min(b["kind"]))
           for i, b in enumerate(blocks)) > usable + 1e-6:
        floors = [0.0 if b["kind"] == "rows" else f
                  for b, f in zip(blocks, floors)]

    # 下限加總可能就已經超過版心(第 9 頁:兩張表 + callout)。
    # 這時要縮下限,否則 slack 變負數、版塊會被排到版面外(實測掉出底部 0.74")。
    # 🔑 縮的順序有先後:**先縮 callout,表格的下限最後才動** ——
    #    callout 短一句只是話變短,表格少一列是投影片自打嘴巴。
    if sum(floors) > usable:
        over = sum(floors) - usable
        for i, b in enumerate(blocks):
            if over <= 1e-9:
                break
            if b["kind"] not in ("matrix", "rows"):
                cut = min(floors[i], over)
                floors[i] -= cut
                over -= cut
        if over > 1e-9:                       # 光表格 + rows 的下限就塞不下,只好等比縮
            k = usable / sum(floors)
            floors = [f * k for f in floors]
    slack = max(usable - sum(floors), 0.0)
    extra = [max(n - f, 0.0) for n, f in zip(needs, floors)]
    tot = sum(extra) or 1.0
    heights = [f + slack * e / tot for f, e in zip(floors, extra)]

    # 🔴 有些版塊**拿到再多高度也用不完**:rows 把每一列畫成自然高度就停了、
    #    callout 就是一句話。多給的那一塊不會消失,它會變成**版面正中央的一個洞**。
    #    第 4 頁刪掉一列之後就是這樣:rows 分到 1.86" 卻只畫得出 1.31",
    #    於是表格和底部那條 callout 中間空了 1.2",看起來像沒做完。
    # 🔑 表格不一樣 —— 它永遠用得完:給多少就把列高平均加上去,字距鬆開反而更好讀。
    #    所以把用不完的收回來交給表格;沒有表格的頁面就讓空白落在**底部**,
    #    那比落在中間好看得多。
    # ⚠️ 門檻 0.30":小於這個的空隙**肉眼看不出來**,而重排會動到每一個版塊的
    #    y 座標。Kenny 逐頁核可過的頁面不該為了 0.08" 的空隙整頁下移 ——
    #    實測加這條之前,第 5 頁十四個形狀全部被推低 0.08",只為了填一個看不見的縫。
    # 🔴 上限不能低於該版塊的**硬下限**(cards 2.32 / compare · bar 2.15 …)。
    #    第一版只用 need 當上限,把 cards 壓到 1.99" —— 低於它的 2.32" 下限,
    #    於是後面那道「補足下限」的借用又從 rows 借回去,rows 反而從 2.74 掉到 2.43,
    #    多丟一列。上限與下限必須是同一組數字,否則兩道規則會互相拉扯。
    cap = [1e9 if b["kind"] == "matrix" else n
           for b, n in zip(blocks, needs)]
    spare = sum(max(h_ - c, 0.0) for h_, c in zip(heights, cap))
    if spare > 1e-6:
        heights = [min(h_, c) for h_, c in zip(heights, cap)]
        # 🔴 先補**還沒吃飽的**版塊,再談填空白。
        #    舊版只把剩餘還給表格,沒有表格的頁面就讓它空著 ——
        #    第 11 頁實測:整頁還剩 0.31" 沒人用,而 rows 差 **0.02"** 就放得下第五列,
        #    於是「不做這件事會怎樣(要再請 12–19 個人)」被丟進備註。
        #    一邊有空位、一邊在丟內容,那不是取捨,是分配錯了。
        short = [max(n - h_, 0.0) for n, h_ in zip(needs, heights)]
        tot_short = sum(short)
        if tot_short > 1e-6:
            give = min(spare, tot_short)
            heights = [h_ + give * sh / tot_short
                       for h_, sh in zip(heights, short)]
            spare -= give
        # 還有剩才輪到「填空白」—— 表格是唯一給多少都用得完的版塊
        mx = [i for i, b in enumerate(blocks) if b["kind"] == "matrix"]
        if mx and spare > 0.30:
            heights[mx[0]] += spare

    # 🔴 每種版塊都有一個「再少就畫不出來」的高度。低於它會產生**零高或負高的形狀**,
    #    python-pptx 照寫、zip 也合法,但 PowerPoint 會拒絕開啟整個檔案。
    #    不足的部分一律向最高的那個版塊借。
    # cards 的下限 2.20" 是算出來的:標題 0.62 + 徽章 0.42 + 內距 0.14 = 1.18" 的固定開銷,
    # 再加三行 12pt(3 x 0.304)。低於這個數,卡片就只剩標題 + 一行殘句。
    # 每個數字都是從對應 blk_* 的固定開銷反推的,不是猜的:
    #   cards   標題區 0.82 + 徽章 0.42 + 內距 0.14 + 三行 12pt(3×0.3038)= 2.29
    #           ⚠️ 是 0.82 不是 0.62 —— blk_cards 的內文起點是 y+0.82(標題框 0.62
    #           上面還有 0.14 的內距,下面還有留白)。用 0.62 算會少 0.20",
    #           結果第 6 頁那張〈三層,一層都不可省〉只印得出兩層,少的正好是放行層。
    #   compare 標題 0.42 + 裁決 0.48 + 內距 + 一行 13pt
    #   rows    fit_items 保底兩列 × 0.34 + 列距 0.10
    #   bar     標題 0.48 + 長條 0.72 + 括號 0.62 + 註記 0.30 = 2.12 → 取 2.15
    #           (低於這個數,註記會被壓到括號說明那一行上面 —— 第 3 頁踩過)
    _MIN = {"callout": 0.36, "cards": 2.32, "rows": 0.80,
            "compare": 2.15, "bar": 2.15, "matrix": 0.70}

    # 🔴 rows 的下限得看**列數**,不能是一個常數。第 2 頁收成一列之後,
    #    常數 0.80" 反而去跟 cards / compare 搶 0.46",兩邊都被擠到印不全。
    # 🔴 下限不能大於這個版塊**實際需要**的高度。
    #    cards 的 2.32" 是照「三行說明」算的,但第 11 頁那三張卡只有兩行(只需 1.99")——
    #    一條「至少要這麼高」的規則把 0.33" 鎖在用不到的地方,
    #    而隔壁的 rows 正好差 0.10" 放不下第五列,被迫把
    #    「不做這件事會怎樣(要再請 12–19 個人)」丟進備註。
    #    「最少要多高才畫得出來」只在它真的有那麼多內容時才成立。
    hard = [min(_MIN.get(b["kind"], 0.3), needs[i])
            for i, b in enumerate(blocks)]

    class MIN:
        @staticmethod
        def get(kind, default=0.3, _i=[0]):
            return _MIN.get(kind, default)

    mins = [max(floors[i], hard[i]) if b["kind"] != "rows"
            else max(floors[i], 0.44)
            for i, b in enumerate(blocks)]
    for _ in range(len(heights)):
        i = min(range(len(heights)),
                key=lambda k: heights[k] - mins[k])
        gap_i = hard[i] - heights[i]
        if gap_i <= 1e-6:
            break
        # 向**餘裕最多**的版塊借,不是向最高的借 —— 最高的那個可能自己就剛好卡在下限
        # (第 10 頁 cards 1.69" 既是最矮也是最高,迴圈當場空轉)。
        j = max(range(len(heights)),
                key=lambda k: heights[k] - max(mins[k], floors[k]))
        if j == i:
            break
        # 借的時候不能借到低於對方的 floor —— 否則剛保住的表格列又被借走
        # (實測第 10 頁四臂表就是這樣從 1.64" 被扣到 1.58",少印一臂)。
        take = min(gap_i, heights[j] - max(mins[j], floors[j]))
        if take <= 1e-6:
            break
        heights[j] -= take
        heights[i] += take
    tight = [f'{b["kind"]}={h:.2f}"(需 {m:.2f}")'
             for b, h, m in zip(blocks, heights, mins) if h < m - 1e-6]
    if tight:
        # 🔴 這一頁的內容超過版心裝得下的量。不讓 build 掛掉(那會擋住其他頁的產出),
        #    但一定要吵 —— 這是「該砍內容了」的訊號,不是排版還能再擠一擠。
        print(f'   🔴 第 {pg["n"]} 頁版塊塞不下:' + "、".join(tight))

    y = BODY_TOP + 0.16
    for b, h in zip(blocks, heights):
        spec = json.loads(b["spec"]) if isinstance(b["spec"], str) else b["spec"]
        fn = N.RENDER[b["kind"]]
        width = N.W - 2 * M
        # 底部的 block 若會伸進人像淨空區,寬度收到 SAFE_X
        if y + h > N.SAFE_Y:
            width = min(width, N.SAFE_X - M)
        if b["kind"] == "cards":
            fn(s, M, y, width, h, spec["items"])
        elif b["kind"] == "rows":
            fn(s, M, y, width, h, spec["items"])
        else:
            fn(s, M, y, width, h, spec)
        y += h + gap

    if foot:
        N.text(s, M, FOOT_Y, N.SAFE_X - M, 0.5, foot, 13, N.GREY, who="foot")

    # 版塊放不下而被截斷的完整內容 → 講者備註(內容不丟)
    extra = [f"【{k}】{v}" for k, v in N.OVERFLOW]
    N.OVERFLOW.clear()
    if pg.get("why"):
        extra.append(f"【這一頁的作用】{pg['why']}")
    put_notes(s, pg["n"] - 1, extra)      # 投影片 n → 單元索引 n-1
    return s


def main():
    pages = json.load(open(SPEC, encoding="utf-8"))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(N.W), Inches(N.H)

    cover(prs)
    for pg in pages:
        content(prs, pg)
    refs_page(prs, len(pages) + 1)   # 封面 0、內容 1..N、文獻 N+1

    prs.save(OUT)

    pic = sum(1 for sl in prs.slides for sh in sl.shapes if sh.shape_type == 13)
    shp = sum(len(sl.shapes) for sl in prs.slides)
    try:
        import lock_pages
        lock_pages.check(OUT)
    except Exception as e:
        print(f"   ⚠️ 鎖定比對跳過:{e}")
    print(f"✅ {os.path.basename(OUT)}")
    print(f"   {len(prs.slides)} 頁 · {os.path.getsize(OUT)/2**20:.2f} MB")
    print(f"   形狀 {shp} 個 · 圖片 {pic} 個 {'✅ 全原生可編輯' if pic == 0 else '🔴 仍有貼圖'}")
    return OUT


if __name__ == "__main__":
    main()
