# -*- coding: utf-8 -*-
"""A3 v4 投影片產生器 —— **組版框架**(12 張:封面 + P1-P10 + 參考文獻)。

本檔只負責「框架」:版面常數、色票、共用 helper、頁序、輸出。
**十二頁的實際內容不寫在這裡**,分三個模組並行填寫(避免互相覆蓋):

    deck/slides_01_04.py   封面 + P1-P4   (page_idx 0-4)
    deck/slides_05_07.py   P5-P7          (page_idx 5-7)
    deck/slides_08_12.py   P8-P10 + 參考文獻 (page_idx 8-11)

每個模組只需實作 `def build(prs, ctx)`,用 `ctx["new_slide"](prs, idx)` 開頁,
其餘一律走 ctx 裡的 helper。尚未填的頁由本框架自動補成帶正確標題與導軌的空頁,
所以任何時候 `python build_deck_v4.py` 都跑得完、都輸出 12 張。

版面硬規則
----------
- 16:9,13.333 x 7.5 吋;字體 Microsoft JhengHei
- 右下淨空區(SAFE_R x SAFE_B)錄影時放人像小窗:**文字框一律不得進入**
  (`place_slide_text()` 會 assert 擋下來)
- 投影片上的字級 **不得低於 12pt**(MIN_PT)
- 每頁只有一個標題,由 `page_header()` 統一放;圖裡不再畫標題
- 圖一律 **100% 原尺寸置入,絕不縮放**(`place_fig()`);300dpi PNG 縮了字就不足 9pt
"""
import os
import sys
import importlib

# ------------------------------------------------------------------
# 路徑:本檔在 Assessment3/v4/scripts/,共用版面庫在 Assessment3/build_deck.py
# ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))   # .../v4/scripts
V4 = os.path.dirname(HERE)                          # .../v4
A3 = os.path.dirname(V4)                            # .../Assessment3
for _p in (A3, HERE):                               # HERE 讓 deck/ 與 figs/ 可被 import
    if _p not in sys.path:
        sys.path.insert(0, _p)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

import build_deck
from build_deck import (txt, header, picture, note, card,
                        RED, GREY, DARK, WHITE, LIGHT, FONT,
                        W, H, TOP, BOT, SAFE_R, SAFE_B, SECTIONS)

# ------------------------------------------------------------------
# v4 色票(裁決 E)。build_deck.py 是 v3 也在用的共用庫,**不改那個檔**,
# 只在執行期覆寫它的模組屬性 —— header()/card() 內部查的是 build_deck 的
# 模組層變數,只改本檔的區域變數沒有用。
# ------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x4E, 0x79)      # 深藍:自動化 / 已具備能力
ORANGE = RGBColor(0xE8, 0x83, 0x3A)    # 橘:人力成本
DEEP_ORANGE = RGBColor(0xB4, 0x5F, 0x1F)   # 橘字(小字用,對比足)
build_deck.NAVY = NAVY
build_deck.ORANGE = ORANGE

# 圖檔目錄改指 v4/figures(picture() 讀的是 build_deck.FIG)
FIG = os.path.join(V4, "figures")
build_deck.FIG = FIG

# 交件檔輸出到 v4/ 之下,不覆蓋上一層 Assessment3/ 的 v3 同名檔
OUT = os.path.join(V4, "Huang_26254793_421104_Assessment 3.pptx")

N_SLIDES = 12
MIN_PT = 12.0                          # 投影片字級下限(比圖內的 9pt 嚴)

# 右下淨空區的左上角座標
SAFE_X = W - SAFE_R                    # 11.283
SAFE_Y = H - SAFE_B                    # 6.35
_EPS = 1e-6

# ------------------------------------------------------------------
# 頁序:0 = 封面,1-10 = P1-P10,11 = 參考文獻
# 標題不含「P1 ·」前綴 —— header() 會自己補上段名。
# ------------------------------------------------------------------
PAGE_TITLES = [
    None,                                        # 0  封面(不放 header)
    "用例:兩階段推論,與它下游那道人工防線",      # 1  P1
    "為什麼是現在:工作量沒變多,客戶要的時效變了",  # 2  P2
    "發現過程:兩種病,兩個缺口",                  # 3  P3
    "我們看不見的那一半:有資料,沒有習慣",        # 4  P4
    "五條路,含三條不用 AI",                      # 5  P5
    "風險:七類八條,緩解位移與殘餘風險",          # 6  P6
    "倫理與可申訴性:對照澳洲 AI 倫理八原則",      # 7  P7
    "怎麼算成功:三層 KPI,一條回圈",              # 8  P8
    "六個月 · 六個工作包 · 三道決策門",           # 9  P9
    "資源、當責與請求",                           # 10 P10
    "參考文獻",                                   # 11
]

# 五段導軌:SECTIONS = 介紹 / AI 機會 / 業務影響 / 效益與衡量 / 項目計畫
# (依 notes/04 各頁「進度指示」段的指定格)
PAGE_SECTION = [
    None,   # 0  封面
    0,      # 1  P1  介紹
    1,      # 2  P2  AI 機會
    1,      # 3  P3  AI 機會
    1,      # 4  P4  AI 機會
    1,      # 5  P5  AI 機會
    2,      # 6  P6  業務影響
    2,      # 7  P7  業務影響
    3,      # 8  P8  效益與衡量
    4,      # 9  P9  項目計畫
    4,      # 10 P10 項目計畫
    None,   # 11 參考文獻(不佔十頁額度,不加導軌)
]

# 各頁配圖(None = 該頁不放圖)。實體尺寸見 FIG_INCHES。
PAGE_FIG = [None] + ["fig_v4_%02d.png" % i for i in range(1, 11)] + [None]

# 300dpi PNG 的實體尺寸(吋),置入時就是這個尺寸,不得縮放
FIG_INCHES = {
    "fig_v4_01.png": (12.20, 5.70), "fig_v4_02.png": (11.95, 2.81),
    "fig_v4_03.png": (9.69, 4.93),  "fig_v4_04.png": (10.07, 5.62),
    "fig_v4_05.png": (12.20, 5.70), "fig_v4_06.png": (12.20, 5.70),
    "fig_v4_07.png": (12.20, 5.70), "fig_v4_08.png": (12.20, 5.70),
    "fig_v4_09.png": (12.20, 5.70), "fig_v4_10.png": (12.20, 5.70),
}

# 摘要 / 結論頁的四卡橫排參數
CARD_W, CARD_GAP, CARD_X0 = 2.62, 0.16, 0.55

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
           "bottom": MSO_ANCHOR.BOTTOM}


# ==================================================================
# 小工具
# ==================================================================
def rgb(c):
    """接受 RGBColor / "#1F4E79" / "1F4E79" / (31, 78, 121)。"""
    if isinstance(c, RGBColor):
        return c
    if isinstance(c, str):
        return RGBColor.from_string(c.lstrip("#").upper())
    if isinstance(c, (tuple, list)) and len(c) == 3:
        return RGBColor(*c)
    raise TypeError("看不懂的顏色:%r" % (c,))


def assert_clear_of_safe(x, y, w, h, who="(未命名)"):
    """右下淨空區要放錄影的人像小窗,任何文字框都不得與它相交。"""
    if x + w > SAFE_X + _EPS and y + h > SAFE_Y + _EPS:
        raise AssertionError(
            "%s 落進右下淨空區:框 x %.2f-%.2f / y %.2f-%.2f,"
            "淨空區自 x %.2f 且 y %.2f 起。"
            "若要保留這個 y,寬度最多 %.2f 吋;若要保留這個寬度,底緣最多到 y %.2f。"
            % (who, x, x + w, y, y + h, SAFE_X, SAFE_Y,
               max(0.0, SAFE_X - x), SAFE_Y))


def assert_on_canvas(x, y, w, h, who="(未命名)"):
    if x < -_EPS or y < -_EPS or x + w > W + _EPS or y + h > H + _EPS:
        raise AssertionError(
            "%s 超出版面:框 x %.2f-%.2f / y %.2f-%.2f,版面 %.3f x %.3f"
            % (who, x, x + w, y, y + h, W, H))


# ==================================================================
# 🔑 三段共用:批次放原生文字框
# ==================================================================
def place_slide_text(slide, blocks, default_color=DARK):
    """把「移出圖的文字」批次放成投影片原生文字框。

    blocks = list[dict],每個 dict 就是一塊文字:

        必填  x, y, w, h   吋,文字框的左 / 上 / 寬 / 高
        文字  text=str(可含 \\n)  或  text=list[str]  或  lines=list[str]
              也接受 rows=[(左欄, 右欄), ...] —— 會併成「左欄  右欄」單欄輸出
        選填  id      塊代號(如 "P5-T01"),assert 出錯時會印出來,強烈建議帶
              size    字級,預設 12;**低於 12 直接 assert 擋下**
              color   RGBColor / "#RRGGBB" / (r,g,b),預設 DARK
              bold    True / False
              align   "left"(預設)/ "center" / "right" / "justify"
              anchor  "top"(預設)/ "middle" / "bottom"
              spacing 行距倍數,預設 1.20
              gap     rows 模式的左右欄間距(全形空白數),預設 2

    放之前逐塊檢查:字級 >= 12pt、框在版面內、框不與右下淨空區相交。
    回傳放好的 textbox list(順序同 blocks)。
    """
    made = []
    for i, b in enumerate(blocks):
        who = str(b.get("id") or b.get("label") or "blocks[%d]" % i)

        if "rows" in b and b.get("text") is None and "lines" not in b:
            gap = "　" * int(b.get("gap", 2))
            body = "\n".join(gap.join(str(c) for c in r) for r in b["rows"])
        else:
            raw = b.get("text", b.get("lines"))
            if raw is None:
                raise AssertionError("%s 沒有任何文字(text / lines / rows 三選一)" % who)
            body = "\n".join(raw) if isinstance(raw, (list, tuple)) else str(raw)
        if not body.strip():
            raise AssertionError("%s 的文字是空的 —— 移出圖的字一個都不能少" % who)

        size = float(b.get("size", MIN_PT))
        if size < MIN_PT - _EPS:
            raise AssertionError(
                "%s 字級 %.1fpt 低於投影片下限 %.0fpt(投影用,不是圖內的 9pt)"
                % (who, size, MIN_PT))

        try:
            x, y, w, h = (float(b[k]) for k in ("x", "y", "w", "h"))
        except KeyError as e:
            raise AssertionError("%s 缺少位置參數 %s" % (who, e))
        assert_on_canvas(x, y, w, h, who)
        assert_clear_of_safe(x, y, w, h, who)

        made.append(txt(
            slide, x, y, w, h, body,
            size=size,
            color=rgb(b.get("color", default_color)),
            bold=bool(b.get("bold", False)),
            align=_ALIGN[b.get("align", "left")],
            anchor=_ANCHOR[b.get("anchor", "top")],
            spacing=float(b.get("spacing", 1.20)),
        ))
    return made


# ==================================================================
# 圖:100% 原尺寸置入,絕不縮放
# ==================================================================
def fig_inches(name):
    """回傳 (寬, 高) 吋 —— 直接由 PNG 的畫素與 dpi 算,不查表。"""
    path = os.path.join(FIG, name)
    with Image.open(path) as im:
        dpi = im.info.get("dpi", (300.0, 300.0))
        return im.size[0] / dpi[0], im.size[1] / dpi[1]


def place_fig(slide, name, left=None, top=TOP, center=True):
    """以 100% 原尺寸置入圖片(不縮放),回傳它佔的矩形。

    left=None 時:center=True 於整個版面置中,center=False 則靠左 0.55。
    ⚠️ 滿版圖(12.2 x 5.7)本來就會蓋到右下淨空區的範圍,那是可以的
       —— 圖在該處刻意留白。**不可以進去的是文字框**,place_slide_text() 會擋。
       置入前務必先讀該頁 scripts/figs/fig_v4_XX.py 的檔頭,確認留白區座標,
       文字框只能放在留白區裡。
    """
    path = os.path.join(FIG, name)
    if not os.path.exists(path):
        raise AssertionError("缺圖:%s" % path)
    fw, fh = fig_inches(name)
    if left is None:
        left = (W - fw) / 2.0 if center else 0.55
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(fw), Inches(fh))
    return dict(left=left, top=top, w=fw, h=fh,
                right=left + fw, bottom=top + fh)


# ==================================================================
# 開頁 / 標題
# ==================================================================
_HEADED = set()


def page_header(slide, page_idx, sub=None, title=None):
    """放本頁唯一的標題與五段導軌。同一頁呼叫第二次會直接 assert。"""
    if id(slide) in _HEADED:
        raise AssertionError("page %s 已經有標題了 —— 每頁只能出現一次標題" % page_idx)
    _HEADED.add(id(slide))
    t = PAGE_TITLES[page_idx] if title is None else title
    if t:
        header(slide, t, PAGE_SECTION[page_idx], sub=sub)
    return slide


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def new_slide(prs, page_idx, sub=None, title=None, ctx=None):
    """開一頁並放好標題 / 導軌。三段模組一律用這個開頁。"""
    s = blank_slide(prs)
    page_header(s, page_idx, sub=sub, title=title)
    if ctx is not None:
        ctx["emitted"].append(page_idx)
    return s


def speaker_notes(slide, text):
    """講者備註。放不下的長文可以擺這裡,但正面該有的內容不得用它頂替。"""
    slide.notes_slide.notes_text_frame.text = text
    return slide


# ==================================================================
# ctx —— 三段模組共用的東西全在這一包裡
# ==================================================================
def make_ctx():
    return {
        # helper
        "new_slide": new_slide, "blank_slide": blank_slide,
        "page_header": page_header, "place_slide_text": place_slide_text,
        "place_fig": place_fig, "fig_inches": fig_inches,
        "speaker_notes": speaker_notes,
        "txt": txt, "header": header, "picture": picture,
        "note": note, "card": card, "rgb": rgb,
        "assert_clear_of_safe": assert_clear_of_safe,
        # 色票
        "NAVY": NAVY, "ORANGE": ORANGE, "DEEP_ORANGE": DEEP_ORANGE,
        "RED": RED, "GREY": GREY, "DARK": DARK, "WHITE": WHITE, "LIGHT": LIGHT,
        "FONT": FONT,
        # 版面常數
        "W": W, "H": H, "TOP": TOP, "BOT": BOT,
        "SAFE_R": SAFE_R, "SAFE_B": SAFE_B, "SAFE_X": SAFE_X, "SAFE_Y": SAFE_Y,
        "MIN_PT": MIN_PT,
        "CARD_W": CARD_W, "CARD_GAP": CARD_GAP, "CARD_X0": CARD_X0,
        # 頁序與資源
        "SECTIONS": SECTIONS, "PAGE_TITLES": PAGE_TITLES,
        "PAGE_SECTION": PAGE_SECTION, "PAGE_FIG": PAGE_FIG,
        "FIG_INCHES": FIG_INCHES, "FIG": FIG, "V4": V4,
        # 本次 build 已產出的頁號(框架用來補未填頁,模組請勿自行改寫)
        "emitted": [],
    }


# ==================================================================
# 分段模組:模組名 -> 負責的頁號
# ==================================================================
SEGMENTS = [
    ("deck.slides_01_04", list(range(0, 5))),    # 封面 + P1-P4
    ("deck.slides_05_07", list(range(5, 8))),    # P5-P7
    ("deck.slides_08_12", list(range(8, 12))),   # P8-P10 + 參考文獻
]


def _placeholder(prs, idx, ctx):
    """該頁尚未填內容 —— 補一張有正確標題與導軌的空頁,讓整份還是 12 張。"""
    s = new_slide(prs, idx, ctx=ctx)
    if idx == 0:
        txt(s, 1.0, 3.2, 11.3, 0.8, "(封面待填)", size=20, color=GREY)
    else:
        txt(s, 0.55, TOP + 0.40, W - 1.1 - SAFE_R, 0.8,
            "(本頁內容待填)", size=16, color=GREY)
    return s


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    ctx = make_ctx()

    for mod_name, pages in SEGMENTS:
        mod = importlib.import_module(mod_name)
        before = list(ctx["emitted"])
        mod.build(prs, ctx)
        got = ctx["emitted"][len(before):]

        bad = [i for i in got if i not in pages]
        if bad:
            raise AssertionError("%s 產出了不屬於它的頁:%s(它只負責 %s)"
                                 % (mod_name, bad, pages))
        if got != sorted(got):
            raise AssertionError("%s 的頁序不是遞增:%s" % (mod_name, got))
        if len(set(got)) != len(got):
            raise AssertionError("%s 有重複的頁:%s" % (mod_name, got))

        for idx in pages:                      # 未填的頁補空殼,順序照 pages
            if idx not in got:
                _placeholder(prs, idx, ctx)

    n = len(prs.slides._sldIdLst)
    assert n == N_SLIDES, "v4 必須是 %d 張,目前 %d 張" % (N_SLIDES, n)
    assert sorted(ctx["emitted"]) == list(range(N_SLIDES)), \
        "頁號不完整:%s" % sorted(ctx["emitted"])

    prs.save(OUT)
    return prs, OUT


if __name__ == "__main__":
    prs, out = build()
    print("投影片頁數:", len(prs.slides._sldIdLst))
    print("版面尺寸: %.3f x %.3f 吋" % (prs.slide_width.inches, prs.slide_height.inches))
    print("圖檔目錄:", build_deck.FIG)
    print("深藍 / 橘: #%s / #%s" % (NAVY, ORANGE))
    print("輸出:", out)
    print("大小: %.1f KB" % (os.path.getsize(out) / 1024))
