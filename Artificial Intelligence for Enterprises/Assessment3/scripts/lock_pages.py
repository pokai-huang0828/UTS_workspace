# -*- coding: utf-8 -*-
"""把已定案的頁鎖起來 —— 之後改引擎不能動到它們。

用法:
    python scripts/lock_pages.py 1 2      # 把第 1、2 頁的現況存成指紋
    python scripts/lock_pages.py --check  # 比對(render_deck 會自動呼叫)

🔴 為什麼要有這支:Kenny 逐頁審,審過的頁就定案了。
   但引擎是全域的 —— 為了第 10 頁調一個下限,第 1 頁的字可能就跟著跑。
   指紋比對讓「不要再動」變成一條**機器守得住**的規則,而不是我記得就好。
"""
import json, os, sys
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
A3 = os.path.dirname(HERE)
DECK = os.path.join(A3, "Huang_26254793_421104_Assessment 3.pptx")
LOCK = os.path.join(A3, "notes", "deck_spec", "_locked_pages.json")


def fingerprint(deck=DECK, pages=None):
    prs = Presentation(deck)
    out = {}
    for i, s in enumerate(prs.slides, 1):
        if pages and i not in pages:
            continue
        items = []
        for sh in s.shapes:
            geo = (round(sh.left / 914400, 2), round(sh.top / 914400, 2),
                   round(sh.width / 914400, 2), round(sh.height / 914400, 2))
            if sh.has_table:
                body = ["｜".join(c.text for c in r.cells) for r in sh.table.rows]
            elif sh.has_text_frame:
                body = sh.text_frame.text.split("\n")
            else:
                body = []
            items.append({"geo": geo, "text": body})
        out[str(i)] = items
    return out


def save(pages):
    data = json.load(open(LOCK, encoding="utf-8")) if os.path.exists(LOCK) else {}
    data.update(fingerprint(pages=set(pages)))
    json.dump(data, open(LOCK, "w", encoding="utf-8"), ensure_ascii=False, indent=1)
    print(f"🔒 已鎖定第 {'、'.join(map(str, pages))} 頁")


def check(deck=DECK, quiet=False):
    if not os.path.exists(LOCK):
        return True
    want = json.load(open(LOCK, encoding="utf-8"))
    got = fingerprint(deck, pages={int(k) for k in want})
    ok = True
    for n, items in want.items():
        cur = got.get(n, [])
        if len(cur) != len(items):
            print(f"   🔒🔴 第 {n} 頁被動到了:形狀數 {len(items)} → {len(cur)}")
            ok = False
            continue
        for a, b in zip(items, cur):
            # ⚠️ JSON 存出來是 list、剛渲染的是 tuple —— 不轉型就永遠不相等,
            #    整個比對會變成每一行都在叫「被動到了」。
            if list(a["geo"]) != list(b["geo"]):
                print(f'   🔒🔴 第 {n} 頁被動到了(位置/大小):{a["text"][:1]} '
                      f'{a["geo"]} → {b["geo"]}')
                ok = False
            if a["text"] != b["text"]:
                print(f'   🔒🔴 第 {n} 頁被動到了(文字):{a["text"][:1]} → {b["text"][:1]}')
                ok = False
    if ok and not quiet:
        print(f"   🔒 已鎖定的第 {'、'.join(sorted(want))} 頁:未被更動 ✅")
    return ok


if __name__ == "__main__":
    if "--check" in sys.argv:
        sys.exit(0 if check() else 1)
    save([int(a) for a in sys.argv[1:] if a.isdigit()])
