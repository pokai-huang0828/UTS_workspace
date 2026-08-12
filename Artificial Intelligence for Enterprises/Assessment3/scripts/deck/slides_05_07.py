# -*- coding: utf-8 -*-
"""第 2 段:P5-P7(page_idx 5, 6, 7)。

三頁的配圖都是 12.20 x 5.70 吋的滿版圖,100% 原尺寸置入後幾乎佔滿整個內容區,
所以原生文字框只能放在兩種地方:

  (1) 圖內實測全白的帶(用 PNG 逐列掃描出來的,不是估算);
  (2) 圖片下緣到版面底之間的頁腳帶(且越過 y 6.35 之後右緣要收到 x 11.28)。

每頁的落點常數寫在該頁的 build 函式裡,座標一律「投影片絕對英吋」。
標題只由 new_slide 放一次;圖裡沒有標題,本檔也不再加標題文字框。
"""
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

from figs.fig_v4_05 import SLIDE_TEXT as P5_TEXT

# ------------------------------------------------------------------ 版面常數
FIG_W, FIG_H = 12.20, 5.70
FIG_L = (13.333 - FIG_W) / 2.0          # 0.5665,置中置圖的左緣
BODY_L, BODY_R = 0.55, 11.28            # 頁腳可用的左右界(右界 = 淨空區左緣)
BODY_W = BODY_R - BODY_L                # 10.73

LINE = 0.176                            # 12pt + spacing 0.88 的實際行高(吋)
SP = 0.88                               # 密排帶用的行距倍數
SP_L = 0.88                             # 頁腳單行用的行距倍數


def _flush(boxes):
    """把文字框的四邊內距歸零 —— 座標算的是文字的實際落點,不是框加內距後的落點。"""
    for tb in boxes:
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = 0
        tf.margin_left = tf.margin_right = Inches(0.02)
    return boxes


def _bar(slide, x, y, w, h, color):
    """實心橫條(反白帶的底)。文字框稍後疊在上面,所以先畫。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ==================================================================
# P5 · 五條路,含三條不用 AI
# ==================================================================
# fig_v4_05 逐列掃描結果:整張圖只有最上 0.12 吋與最下 0.103 吋是全白,
# 中間五欄卡片 / G1 判定門帶 / 判定分岔區完全沒有可疊字的留白。
# 置圖 top = 1.05 -> 圖佔 y 1.05-6.75,圖內下緣白帶對到投影片 y 6.65-6.75。
# 因此本頁的原生文字帶只有 y 6.65-7.48 這一條(高 0.83,約四行 12pt)。
# 規格把最下 12% 指定給「三條策略窄帶的量化落點 -> 競爭地位帶 -> 頁腳引文條」,
# 這條頁腳帶就照那個次序放;其餘各欄的長文放講者備註(見 _p5_notes)。
P5_FOOT_Y = 6.65

# 這三塊放在投影片正面,其餘塊放講者備註
P5_ON_SLIDE = ("P5-T12", "P5-T13", "P5-T14")

# 純施工用的括註 / 版面警語,不是給觀眾看的內容,講者備註不收(紅線 5)
P5_STRIP = (
    "(這一行不得省,它是本頁最重要的邏輯)",
    "〔紅字警語〕",
    "本頁不得把這兩個數字放在同一張圖或同一行上。",
)


def _p5_block(b):
    """把 SLIDE_TEXT 的一塊還原成純文字(rows 併成「左欄  右欄」單欄)。"""
    if "rows" in b and b.get("text") is None:
        return "\n".join("　　".join(str(c) for c in r) for r in b["rows"])
    return "\n".join(b["text"])


def _p5_notes():
    """沒放上正面的塊,逐字進講者備註 —— 移出圖的字不能在交件物裡消失。"""
    out = []
    for b in P5_TEXT:
        if b["id"] in P5_ON_SLIDE:
            continue
        t = _p5_block(b)
        for s in P5_STRIP:
            t = t.replace(s, "")
        out.append(t.strip())
    return "\n\n".join(out)


def slide_05(prs, ctx):
    s = ctx["new_slide"](prs, 5, ctx=ctx)
    ctx["place_fig"](s, "fig_v4_05.png", top=ctx["TOP"])

    by_id = {b["id"]: b for b in P5_TEXT}

    # (1) 三條策略窄帶裡唯一能整塊放進頁腳的量化落點
    _flush(ctx["place_slide_text"](s, [
        dict(id="P5-T12", x=BODY_L, y=P5_FOOT_Y, w=BODY_W, h=0.19,
             text="　".join(by_id["P5-T12"]["text"]),
             size=12, color=ctx["DEEP_ORANGE"], spacing=SP_L),
    ]))

    # (2) 競爭地位帶:深藍實心橫條 + 白字(規格指定的反白帶,先畫底再疊字)
    _bar(s, BODY_L, 6.86, BODY_W, 0.40, ctx["NAVY"])
    _flush(ctx["place_slide_text"](s, [
        dict(id="P5-T13", x=BODY_L + 0.08, y=6.86, w=BODY_W - 0.16, h=0.40,
             text=_p5_block(by_id["P5-T13"]),
             size=12, color=ctx["WHITE"], spacing=SP, anchor="middle"),
    ]))

    # (3) 頁腳引文條:左句 + 右引文
    t14 = by_id["P5-T14"]["text"]
    _flush(ctx["place_slide_text"](s, [
        dict(id="P5-T14a", x=BODY_L, y=7.29, w=BODY_W - 2.10, h=0.19,
             text=t14[0], size=12, color=ctx["GREY"], spacing=SP_L),
        dict(id="P5-T14b", x=BODY_R - 2.10, y=7.29, w=2.10, h=0.19,
             text=t14[1], size=12, color=ctx["GREY"], spacing=SP_L,
             align="right"),
    ]))

    ctx["speaker_notes"](s, _p5_notes())
    return s


# ==================================================================
# P6 · 風險:七類八條,緩解位移與殘餘風險
# ==================================================================
# fig_v4_06 逐列掃描:最上 0.147 吋、最下 0.25 吋全白(整張寬)。
# 置圖 top = 1.05 -> 圖佔 y 1.05-6.75,圖內下緣白帶對到投影片 y 6.50-6.75,
# 接上圖外的 6.75-7.50,得到一條 y 6.50-7.48 的通欄文字帶(高 0.98,約五行)。
P6_SUB = "技術局限 · 數據(2 條)· 倫理 · 社會性 · 經濟 · 組織與管理 · 政治法律政策"

# 「緩解自身的風險」九條(順序同圖上表列),合成一段通欄放進頁腳帶
P6_SELF_RISK = [
    "R1 技術局限:規則層擋不住的情境仍回到黑箱",
    "R2 數據:雙人複核使複核成本近倍增,需以抽樣比例控制",
    "R2 附 數據:地圖速限資料本身有誤時整批傳遞,且比模型錯更難察覺",
    "R3 數據:不顯示信心使複核變慢,與時效缺口相衝",
    "R4 倫理:限制用途會降低客戶端採用意願",
    "R5 社會性:參與式做法拉長判準定案時程",
    "R6 經濟:分段釋出使長週期工作包難以排人",
    "R7 組織與管理:決策門增加治理負擔與會議成本",
    "R8 政治法律政策:無:目前沒有可執行的緩解措施",
]

# R2 那一列的偵測訊號寫「基準見頁腳」,這條頁腳註腳沒放,圖上就會指向空處
P6_R2_BASE = ("R2 偵測訊號的基準 = 某小型車隊單月樣本的全部誤報中 "
              "103 / 124 = 83% 屬交通標誌辨識根因;此為該樣本內的比例,不是全公司數字。")


def slide_06(prs, ctx):
    s = ctx["new_slide"](prs, 6, ctx=ctx)
    ctx["place_fig"](s, "fig_v4_06.png", top=ctx["TOP"])

    _flush(ctx["place_slide_text"](s, [
        dict(id="P6-SUB", x=BODY_L, y=0.90, w=10.70, h=0.19,
             text=P6_SUB, size=12, color=ctx["GREY"], spacing=SP),
        dict(id="P6-SELF", x=BODY_L, y=6.52, w=BODY_W, h=0.71,
             text="緩解自身的風險(九條):" + " · ".join(P6_SELF_RISK),
             size=12, color=ctx["DARK"], spacing=SP),
        dict(id="P6-R2BASE", x=BODY_L, y=7.27, w=BODY_W, h=0.19,
             text=P6_R2_BASE, size=12, color=ctx["GREY"], spacing=SP_L),
    ]))
    return s


# ==================================================================
# P7 · 倫理與可申訴性
# ==================================================================
# 標題只寫「倫理與可申訴性」:圖的左欄小標已經是「澳洲 AI 倫理八原則 · 逐條對照」,
# 標題再寫全就會在同一頁出現兩次同義標題。
P7_TITLE = "倫理與可申訴性"

# 置圖 top = 1.08 -> 圖佔 y 1.08-6.78。逐列掃描出的右半四條全白帶(圖內座標,
# 由上而下),加上 1.08 換算成投影片 y:
#   白帶 0  圖 1.087-1.730 -> 投影片 2.167-2.810
#   白帶 1  圖 2.350-3.000 -> 投影片 3.430-4.080
#   白帶 2  圖 3.620-4.273 -> 投影片 4.700-5.353
#   白帶 3  圖 4.890-5.700 -> 投影片 5.970-6.780(y > 6.35 之後右緣要收到 11.28)
# 全白帶的圖內 x 為 7.27-12.14 -> 投影片 7.837-12.707,文字框取 7.90 / 寬 4.74。
P7_FIG_TOP = 1.08
P7_CARD_X, P7_CARD_W = 7.90, 4.74

P7_A = ("本頁對照的是 2019 年版八原則;主管機關已於 2025 年發布後續的 AI 採用指引,"
        "官方頁面現以過去式描述這八條。本案仍以 2019 年八原則作為對照骨架,"
        "理由是它逐條可對、可公開查證;移交前由提案人依後續指引重跑一次對照。")

P7_B1 = "我決定的:錯誤率量測一旦建立,複核優先序改以風險排序而非客戶大小,新的排序規則提交委員會核定"
P7_B2 = "我不能決定的:商業上的客戶分級本身"
P7_B3 = ("當責 · 時點 · 會改變哪個決策:複核營運主管於 W6 交出「這道防線目前覆蓋了誰」的第一份紀錄;"
         "新的風險排序規則於 G1 提交委員會;委員會核定與否,決定這一條在移交時是紅燈還是琥珀。")
P7_B4 = "被產能擋住的不是簽約,是我們能對多少客戶維持同一個品質水準。"

# 🔴 2026-08-12 裁決 X:舊句「判定仍由人做…設計本身要求的一環」在讀法(i) 下是假的,全片撤除。
#    改為把「專案期」與「移交後」兩個時點分開講(P05-07.md v2.5 §P7 追加行)。
P7_C1 = ("這不是說法,是設計上的必要 —— 採用案的分流率是推估的 85%〔非實測〕,"
         "也就是設計上預期約 15% 要人自己調影片看;模型做的是把證據整理好,"
         "專案期內判定仍全數由人做(影子模式),所以人的覆核在期內不是過渡安排;"
         "但移交後開到天花板時,那 850 筆將由機器單方面結案 —— 兩個時點必須分開講。")
P7_C2 = "我決定的:轉任與再訓練規劃列為移交前置條件"
P7_C3 = ("我不能決定的:人事承諾本身 —— 需贊助人於 G1(第 6 週)核定;"
         "在核定之前,本案不對任何人的職位或編制做出承諾。")
P7_C4 = "當責:項目贊助人 · 時點:G1 · 會改變哪個決策:核定與否決定移交條件能不能簽"

P7_D1 = "處置:判準手冊明文化 + 抽樣雙人複核 + 判定可回溯"
P7_D2 = "我決定的:申訴機制到位前不供個人考核使用"
P7_D3 = ("我決定的:回抽採最小必要抽樣,以事件片段入庫而非整日影片,樣本用畢即銷毀;"
         "資料授權範圍的確認列為 G1 前置事項")
P7_D4 = "我不能決定的:客戶端的申訴流程本身"

P7_E1 = "我決定的:內部稽核報表常態化"
P7_E2 = "我不能決定的:要不要對客戶揭露、揭露到哪"
P7_E3 = ("當責 · 時點 · 會改變哪個決策:提案人於 G1 把揭露層級的選項與各自後果整理成一頁提交;"
         "委員會在 G1 選一個層級,這個選擇決定 WP-F 移交文件裡對外報表的欄位。")

P7_G = "表上的狀態不是自我打分,是移交前的待辦清單。"

# 四條白帶各 0.65 吋,12pt 只裝得下三行;裝不進正面的四塊改進講者備註,逐字保留。
P7_NOTES = "\n\n".join([P7_B4, P7_B3, P7_D3, P7_E3])


def slide_07(prs, ctx):
    s = ctx["new_slide"](prs, 7, title=P7_TITLE, ctx=ctx)

    # 副標段(2019 / 2025 口徑):規格禁止縮成頁腳。圖自 y 1.13 起就滿版有墨,
    # 標題下方放不下兩行,故放標題右側的空白(標題只佔到 x 5.2,導軌自 x 11.05 起)。
    _flush(ctx["place_slide_text"](s, [
        dict(id="P7-A", x=5.50, y=0.32, w=5.50, h=0.74,
             text=P7_A, size=12, color=ctx["GREY"], spacing=SP),
    ]))

    ctx["place_fig"](s, "fig_v4_07.png", top=P7_FIG_TOP)

    x, w = P7_CARD_X, P7_CARD_W
    _flush(ctx["place_slide_text"](s, [
        # ---- 白帶 0(卡片零 · 品質防線是配給制)----
        dict(id="P7-B1", x=x, y=2.19, w=w, h=0.36,
             text=P7_B1, size=12, color=ctx["NAVY"], spacing=SP),
        dict(id="P7-B2", x=x, y=2.56, w=w, h=0.19,
             text=P7_B2, size=12, color=ctx["RED"], spacing=SP),
        # ---- 白帶 1(卡片一 · 工作性質改變)----
        dict(id="P7-C2", x=x, y=3.45, w=w, h=0.19,
             text=P7_C2, size=12, color=ctx["NAVY"], spacing=SP),
        dict(id="P7-C3", x=x, y=3.65, w=w, h=0.36,
             text=P7_C3, size=12, color=ctx["DEEP_ORANGE"], spacing=SP),
        # ---- 白帶 2(卡片二 · 判定會影響駕駛)----
        dict(id="P7-D12", x=x, y=4.72, w=w, h=0.36,
             text=[P7_D1, P7_D2], size=12, color=ctx["NAVY"], spacing=SP),
        dict(id="P7-D4", x=x, y=5.09, w=w, h=0.19,
             text=P7_D4, size=12, color=ctx["RED"], spacing=SP),
        # ---- 白帶 3 上半(卡片三 · 透明度的界線);下半越過淨空區,讓給通欄帶 ----
        dict(id="P7-E1", x=x, y=5.98, w=w, h=0.19,
             text=P7_E1, size=12, color=ctx["NAVY"], spacing=SP),
        dict(id="P7-E2", x=x, y=6.16, w=w, h=0.18,
             text=P7_E2, size=12, color=ctx["DEEP_ORANGE"], spacing=SP),
    ]))

    # ---- 底部窄帶(規格:通欄一行、深底白字)----
    _bar(s, BODY_L, 6.52, BODY_W, 0.26, ctx["NAVY"])
    _flush(ctx["place_slide_text"](s, [
        dict(id="P7-G", x=BODY_L, y=6.52, w=BODY_W, h=0.26,
             text=P7_G, size=12, color=ctx["WHITE"], bold=True,
             align="center", anchor="middle", spacing=SP),
        # ---- 設計上的必要 · 當責與時點 ----
        dict(id="P7-C1", x=BODY_L, y=6.86, w=BODY_W, h=0.36,
             text=P7_C1, size=12, color=ctx["DARK"], spacing=SP),
        dict(id="P7-C4", x=BODY_L, y=7.26, w=BODY_W, h=0.19,
             text=P7_C4, size=12, color=ctx["DARK"], spacing=SP),
    ]))

    ctx["speaker_notes"](s, P7_NOTES)
    return s


# ==================================================================
def build(prs, ctx):
    slide_05(prs, ctx)
    slide_06(prs, ctx)
    slide_07(prs, ctx)
