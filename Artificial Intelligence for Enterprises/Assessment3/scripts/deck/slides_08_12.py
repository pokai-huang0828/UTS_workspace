# -*- coding: utf-8 -*-
"""第 3 段:P8-P10 + 參考文獻(page_idx 8, 9, 10, 11)。

只准動這個檔,不要動 build_deck_v4.py,也不要動另外兩個分段模組。

版位推導(每一頁都先量過圖的墨跡,文字框只放留白區)
----------------------------------------------------
  版面 13.333 x 7.5;圖 12.20 x 5.70 以 100% 原尺寸置中置入 top=TOP=1.05
  -> 圖佔 x 0.567-12.767 / y 1.05-6.75
  標題 26pt 佔 y 0.30-0.80  ->  **上緣條** y 0.79-1.05 可放一行 12pt
  右下淨空區自 x 11.283 且 y 6.35 起  ->  圖下方的頁腳條寬度最多到 x 11.25
  圖內留白(逐頁量測 PNG 墨跡得出,見各頁註解)才是文字框的落點

ctx 有什麼(三段共用)
--------------------
  helper   new_slide(prs, idx, sub=None, title=None, ctx=ctx) · blank_slide ·
           page_header · place_slide_text(slide, blocks) · place_fig(slide, name,
           left=None, top=TOP, center=True) · fig_inches(name) ·
           speaker_notes(slide, text) · txt / header / picture / note / card ·
           rgb("#1F4E79") · assert_clear_of_safe(x, y, w, h, who)
  色票     NAVY · ORANGE · DEEP_ORANGE · RED · GREY · DARK · WHITE · LIGHT · FONT
  版面     W H TOP BOT SAFE_R SAFE_B SAFE_X SAFE_Y MIN_PT CARD_W CARD_GAP CARD_X0
  頁序     SECTIONS · PAGE_TITLES · PAGE_SECTION · PAGE_FIG · FIG_INCHES · FIG · V4

規則
----
  1. 標題只由 new_slide 放,本檔不另加標題文字框。
  2. 文字框一律走 place_slide_text(擋 <12pt 與右下淨空區)。
  3. 圖 100% 原尺寸,不縮放。
  4. 移出圖的文字放不進版位的,改放講者備註(仍在交件物內),並於回報逐條列出。
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from figs.fig_v4_09 import MOVED_TO_SLIDE as P9_MOVED

# ------------------------------------------------------------------ 共用版位
# Microsoft JhengHei 單行高 = 字級 x 1.342;12pt 一行 = 0.224 吋、20pt 一行 = 0.373 吋。
# 文字框內距全部歸零(_zero_margins),所以「框高 >= 行數 x 行高」就是排得下。
# ⚠️ 標題 26pt 行距 1.25 的行框吃到 y 0.956,圖上緣在 1.05 —— 中間只剩 0.09 吋,
#    放不下一行 12pt(0.224 吋),所以本段三頁都沒有「標題與圖之間」的文字條,
#    移出圖的帶狀文字一律走頁腳條或圖內留白區。
FOOT_X, FOOT_W = 0.55, 10.70               # 圖下方頁腳條(避開右下淨空區)
FOOT_Y, FOOT_H = 6.73, 0.77                # 圖下緣 6.75 之下,最多三行 12pt


def _zero_margins(boxes):
    """把文字框內距歸零 —— 密排頁靠這 0.2 吋寬 / 0.1 吋高過活。"""
    for tb in boxes:
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
    return boxes


def _shape(slide, x, y, w, h, line, fill):
    """畫一個圓角容器框(只畫框,字另外用 place_slide_text 放上去)。"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    box.adjustments[0] = 0.06
    return box


# ==================================================================
# P8 · 怎麼算成功:三層 KPI,一條回圈
# ==================================================================
# 圖內墨跡量測:fig_v4_08 由 y 0.10 一路滿到 y 5.60(圖內座標),沒有帶狀留白,
# 所以移出的文字只能放「上緣條」與「圖下方頁腳條」。
P8_SUB = "最上層第一個指標,是客戶等多久拿到分析。下面兩層的指標,都是它的支撐。"

P8_FOOT = (
    P8_SUB + "　"
    "當責三件套(誰負責 · 什麼時候交 · 它會改變哪一個決策)｜"
    "② 獲得人工驗證判定的客戶涵蓋率:業務負責人 · W6 交出基線 · "
    "它決定 G1 要不要把複核優先序從客戶大小改成風險大小｜"
    "⑩ 漏放率:技術負責人 · W6 交出點估計與 95% 信賴區間 · "
    "它決定 G1 放不放行後段工作"
)

P8_NOTES = "\n\n".join([
    "⑤ 每筆批次複核的人工工時:複核營運主管 · W6 交出計時量測 · "
    "它決定全客戶人力推算要不要在 G2 整組重算。",

    "⑪ 複核者一致性:複核營運主管 · W10 交出基線 · 它決定 G2 准不准開放對照重訓｜"
    "錯誤攔截比例(期內不可得):需漏放率先建立才算得出來,期內不可得。"
    "當責:技術負責人;待漏放率基線建立後進入移交後的常態儀表板。",

    "限定語｜① 時效改善的一半來自證據包自動化,另一半來自誤報量下降(脈絡特徵 → 對照重訓);"
    "「2 天內;單純案件當天回覆」是客戶期望,非已簽訂之服務水準協議｜"
    "② 現況為配給制,依購買量與問題嚴重性排序｜"
    "③ 支撐全客戶 13,000+ 筆/天〔所需人力由 G2 後依實測單位工時重估〕;"
    "此項的驅動因子不在我方控制內 —— 由客戶的開通時點決定。",

    "④ 每年約 NT$1,144-2,106 萬:以初階工程師的公開薪資量級推算,非核定預算,"
    "非本案承諾的節省額｜"
    "⑤ 約 20 秒/筆 = 例行批次複核 · 全量;營運端的實務估計,尚未以計時量測驗證｜"
    "⑥ 約 5.25 分鐘/筆:由一次實際判讀反推(5 人 × 3.6 個分析日),不含報告與回信｜"
    "⑦ 降至 7% 以下:以該樣本為基準,G1 後依重新校準的量測重訂。",

    "⑧ 客戶回頭標記率同時是「駕駛端無效警示率」的可觀測代理｜"
    "⑨ 尖峰徵調 7 位研發工程師:離開本職、非常態編制、非專職;"
    "欄 2 全文為「期內建立峰值徵調人時的量測基準,目標值於 G2 訂定;歸零列為移交後目標」,"
    "量測相依:脈絡特徵那一包 W5-12｜"
    "底部閉環同時回答:省下的人去哪了 / 訓練資料哪裡來。"
    "移交後的成熟度階段語彙參照 (IBM, 2020)。",

    "凡今天訂不出目標值的,一律寫明它由哪一道門、在第幾週被訂出來,"
    "並標明是誰負責交、它會改變哪一個決策。",
])


def _p8(prs, ctx):
    s = ctx["new_slide"](prs, 8, ctx=ctx)
    ctx["place_fig"](s, "fig_v4_08.png", top=ctx["TOP"])
    _zero_margins(ctx["place_slide_text"](s, [
        dict(id="P8-F1", x=FOOT_X, y=FOOT_Y, w=FOOT_W, h=FOOT_H,
             text=P8_FOOT, size=12, color=ctx["GREY"], spacing=1.0),
    ]))
    ctx["speaker_notes"](s, P8_NOTES)
    return s


# ==================================================================
# P9 · 六個月 · 六個工作包 · 三道決策門
# ==================================================================
# 圖內墨跡量測:fig_v4_09 甘特圖主體滿到 y 4.30(軸線),W2-W26 刻度字在 y 4.30-4.55,
# 中央圖例在 x 3.58-8.65 / y 4.75-4.95,底部色塊圖例在 y 5.25-5.50。
# -> 留白區 A(左下):圖內 x 0.33-3.53 / y 4.32-5.25  = 投影片 x 0.90-4.10 / y 5.37-6.30
#    留白區 B(右下):圖內 x 8.68-12.10 / y 4.55-5.25 = 投影片 x 9.25-12.67 / y 5.60-6.30
_P9 = {code: lines for code, _pos, lines in P9_MOVED}


def _p9(prs, ctx):
    s = ctx["new_slide"](prs, 9, ctx=ctx)
    ctx["place_fig"](s, "fig_v4_09.png", top=ctx["TOP"])
    t4 = _P9["P9-T4"]
    _zero_margins(ctx["place_slide_text"](s, [
        # 留白區 A / B:T5 兩行當責註記
        dict(id="P9-T5a", x=0.90, y=5.37, w=3.20, h=0.92,
             text=_P9["P9-T5"][0], size=12, color=ctx["NAVY"], spacing=1.0),
        dict(id="P9-T5b", x=9.25, y=5.60, w=3.42, h=0.69,
             text=_P9["P9-T5"][1], size=12, color=ctx["GREY"], spacing=1.0),
        # 頁腳條:T1 + T4「今天 -> 第 26 週」對照(欄首 + 四列,連續排成三行)
        dict(id="P9-T4", x=FOOT_X, y=FOOT_Y, w=FOOT_W, h=FOOT_H, size=12,
             color=ctx["DARK"], spacing=1.0,
             text=_P9["P9-T1"][0] + "　" + t4[0] + ":" + "｜".join(t4[1:])),
    ]))
    ctx["speaker_notes"](s, "\n\n".join(
        ["【三道決策門】"] + _p9_gates() + ["", "【六個工作包 · 副標與交付】"] + _P9["P9-T3"]))
    return s


def _p9_gates():
    """三道門說明卡:去掉組版流程用的內部註記,只留給聽眾看的內容。"""
    out = []
    for line in _P9["P9-T2"]:
        line = line.replace(
            "三方結構、凍結評估集、預註冊門檻與不可取消實驗四項一字不動(裁決 K);"
            "本輪只升級第三組的技術主體,規則式脈絡基線正式成為地端 VLM 的對照基線。",
            "三方結構、凍結評估集、預註冊門檻與不可取消實驗四項一字不動;"
            "規則式脈絡基線是地端 VLM 的對照基線。")
        out.append(line)
    return out


# ==================================================================
# P10 · 資源、當責與請求
# ==================================================================
# 圖內墨跡量測:fig_v4_10 於 y >= 4.35(圖內座標)刻意留白,墨跡最後一列止於 4.35。
# -> 留白帶 = 投影片 x 0.57-12.77 / y 5.40-6.75,再加圖下方 y 6.75-7.50。
#    右下淨空區自 x 11.283 且 y 6.35 起,故文字框一律收在 x <= 11.25。
P10_BOX1 = [
    "算式框一 · 不做本案時的年度增聘成本暴露",
    "　　單人年薪 NT$80-90 萬(初階工程師的公開薪資量級,未使用公司內部薪資資料)",
    "×　間接費率 1.3(提案人設定的假設)",
    "=　NT$104-117 萬/人/年 × 需增聘 11-18 人",
    "=　每年約 NT$1,144-2,106 萬,約 A$54-100 萬(1 A$ 約 NT$21,概算)",
]

P10_BOX2 = [
    "算式框二 · 人力成本估算",
    "　　18.2 人月 × 月薪 NT$6.7-7.5 萬(年薪 80-90 萬 / 12)× 間接費率 1.3",
    "=　約 NT$158-177 萬",
    "=　約 1.4-1.7 個 loaded FTE-year",
]

P10_GATE = [
    "放行條件 · 本案的損益兩平門檻",
    "① 若公司選擇以現行流程覆蓋全客戶,年度擴編暴露約 NT$1,144-2,106 萬;這不是本案承諾節省。",
    "② 完整專案的人力成本約等於 1.4-1.7 個 loaded FTE-year。",
    "③ 只有 G2 實測證明可避免或延後至少此容量,且業務確認會轉成少招人或少徵調研發,才放行後續預算。",
]

P10_ASK_MAIN = "請求核准:六個月整體額度輪廓 · 分段釋出"
P10_ASK_SUB = ("第一段 W1-6(WP-A + WP-B/WP-C 限額探索)約 4.85 人月 / "
               "量級 NT$42-47 萬 · G1 由贊助人決定是否續行")

P10_NOTES = "\n\n".join([
    "算式框二 · 人力成本估算:18.2 人月 × 月薪 NT$6.7-7.5 萬(年薪 80-90 萬 / 12)"
    "× 間接費率 1.3 = 約 NT$158-177 萬 = 約 1.4-1.7 個 loaded FTE-year。"
    "這個金額全片唯一叫法是「人力成本估算」。",

    "框一與框二不得相減、不得畫成同一條式子:框二是本案要花的錢;"
    "框一是不做本案、維持現行做法的年度增聘成本暴露。"
    "反事實成本暴露 → 只有被實測證明可避免的那一部分才算數。",

    "放行條件框的當責:業務負責人提出容量轉換的確認 · 時點 G2(W12)· 它決定第二段預算放不放行。"
    "G1 決定的是是否續行剩餘工作,不是是否開始。",

    "金額為推算量級,非核定預算;薪資、間接費率、匯率三項均為提案人設定的假設值。"
    "本估算涵蓋人力與地端設備,不涵蓋下列五項非人力成本(已辨識、尚未估價):"
    "① 影片儲存 ② 雲端運算 ③ 圖資更新 ④ 測試環境建置 ⑤ 訓練資料蒐集;"
    "量級由技術負責人與數據負責人於 G1 前補齊,它決定第二段額度要不要上修。"
    "因這五項仍未估價,人力與設備相加後仍不得稱為「總成本」,只能稱「已估成本」。",

    "薪資尺說明:NT$80-90 萬為初階工程師的公開薪資量級;要增聘的複核分析人員與"
    "專案內五種角色同屬工程職、共用同一把尺,故此尺可用;但它仍是公開量級,"
    "非公司內部薪資資料。",

    "地端設備:128G 設備 NT$16-20 萬 / 台(一次性),先寫一台起步。"
    "採購時點在 G1 之後、W8 之前,不在第一段 W1-6 的額度內。"
    "設備等於一個人年成本的 14%-19%;硬體只佔本專案已估成本的 8%-12%"
    "—— 這個專案的成本主體是人,不是設備。"
    "每日批次的設備佔用率僅 9%-14%(估算,非實測),餘裕充足;"
    "要不要再加一台由 W1-6 的量測結果評估,列為 G1 的加購判準。"
    "同一套基礎設施的第二個應用,是把內部規格文件做成客戶技術問答的線上回覆"
    "—— 不在本次請求範圍內,但它讓這筆硬體投資的攤提對象不只一個。",

    "WP-E:W8-12 最小離線模型試驗 + W12-24 對照重訓;總投入 4.50 人月不變,"
    "係內部人力前移;逐格 FTE 由技術負責人於 G1 前重排,待確認。",

    "第一段(W1-6)= 評估底座 2.08 + 標註品質前 4 週 1.39 + 脈絡特徵前 2 週 0.69 "
    "+ 專案負責人 0.69 = 約 4.85 人月。W1-6 的範圍是 WP-A 全部 + WP-B / WP-C 的限額探索,"
    "不是只做 WP-A。",

    "六方當責:項目贊助人 / 業務負責人(流程與 KPI)/ 數據負責人(品質與權限)/ "
    "技術負責人(模型 · 系統 · 安全)/ 複核營運主管 / 專案負責人(提案人,PM)。"
    "複核營運主管是評估底座那一包的當責人,因為那件事發生在他的團隊裡"
    "(標註品質那一包該列給 C)。每個工作包只有一個 A;三道決策門的 A 一律落在項目贊助人。",

    "受影響但不在核准鏈上的利害關係人:"
    "① 5 位專職分析人員 —— 工作從逐筆看片變成處理難題與治理;"
    "轉任與再訓練規劃列為移交前置條件,由複核營運主管代表。"
    "② 車隊客戶的安全主管 —— 判定結果的實際使用者,卻不在核准鏈上;由業務負責人代表。"
    "③ 駕駛 —— 被判定的對象,目前無申訴管道;在申訴機制到位前,"
    "本案產出的判準與一致性量測不對外供個人考核使用。",

    "相依關係(待共同確認):本案 W1-26(評估底座 · 標註品質 · 脈絡特徵 · "
    "爭議回覆自動化 · 對照重訓 · 治理移交)→ 公司既有的雲端判讀規劃。"
    "請求委員會指派一次共同評估,列為 G1 前置事項;"
    "若評估結果是兩案應合併或改序,我接受。",

    "組織障礙:政策可改 · 優先序照購買量 · 尖峰徵調研發,"
    "所以我要的不只是預算,是一個贊助人跟三道門。",

    "六個月之後,委員會手上會多三樣東西:① 一份可稽核的錯誤率 "
    "② 一層把畫面外脈絡帶進事件的判斷層 ③ 一份證明或推翻「脈絡到底有沒有用」的對照實驗結果。"
    "就算第二十四週證不出增量,本案仍留下公司第一套可重複的模型評估與標註品質流程,"
    "那個東西不隨第三道門撤掉。",

    "不做這件事的代價不是多花錢,是三件事同時發生:客戶繼續等兩到五天、"
    "小客戶繼續拿到沒有人看過的判定,而我們的研發工程師在尖峰時繼續被拉去看影片。",
])


def _p10(prs, ctx):
    s = ctx["new_slide"](prs, 10, ctx=ctx)
    ctx["place_fig"](s, "fig_v4_10.png", top=ctx["TOP"])

    orange, red, navy = ctx["ORANGE"], ctx["RED"], ctx["NAVY"]
    _shape(s, 0.58, 5.42, 5.62, 1.42, orange, ctx["rgb"]("#FDF1E7"))   # 算式框一
    _shape(s, 6.32, 5.42, 4.93, 1.42, red, ctx["rgb"]("#FBECEC"))      # 放行條件
    _shape(s, 0.58, 6.86, 10.67, 0.64, navy, navy)                     # 結論條

    _zero_margins(ctx["place_slide_text"](s, [
        dict(id="P10-B1H", x=0.66, y=5.47, w=5.46, h=0.23,
             text=P10_BOX1[0], size=12, color=ctx["DEEP_ORANGE"], bold=True,
             spacing=1.0),
        dict(id="P10-B1", x=0.66, y=5.71, w=5.46, h=1.13,
             text=P10_BOX1[1:], size=12, color=ctx["DARK"], spacing=1.0),

        dict(id="P10-GH", x=6.40, y=5.47, w=4.77, h=0.23,
             text=P10_GATE[0], size=12, color=red, bold=True, spacing=1.0),
        dict(id="P10-G", x=6.40, y=5.71, w=4.77, h=1.13,
             text=P10_GATE[1:], size=12, color=ctx["DARK"], spacing=1.0),

        dict(id="P10-ASK", x=0.58, y=6.89, w=10.67, h=0.38,
             text=P10_ASK_MAIN, size=20, color=ctx["WHITE"], bold=True,
             align="center", spacing=1.0),
        dict(id="P10-ASK2", x=0.58, y=7.27, w=10.67, h=0.23,
             text=P10_ASK_SUB, size=12, color=ctx["rgb"]("#F6C9A0"),
             bold=True, align="center", spacing=1.0),
    ]))
    ctx["speaker_notes"](s, P10_NOTES)
    return s


# ==================================================================
# 參考文獻(page_idx 11)—— 無圖、無導軌,雙欄:左 8 條、右 7 條
# ==================================================================
REF_ENTRIES = [
    "Australian Human Rights Commission. (2020). Using artificial intelligence to "
    "make decisions: Addressing the problem of algorithmic bias (Technical paper). "
    "https://humanrights.gov.au/sites/default/files/document/publication/"
    "ahrc_technical_paper_algorithmic_bias_2020.pdf",

    "Department of Industry, Science and Resources. (2019). Australia's AI ethics "
    "principles. Australian Government. https://www.industry.gov.au/publications/"
    "australias-ai-ethics-principles",

    "Ensign, D., Friedler, S. A., Neville, S., Scheidegger, C., & "
    "Venkatasubramanian, S. (2018). Runaway feedback loops in predictive policing. "
    "Proceedings of Machine Learning Research, 81, 160-171. "
    "https://proceedings.mlr.press/v81/ensign18a.html",

    "Goddard, K., Roudsari, A., & Wyatt, J. C. (2012). Automation bias: A systematic "
    "review of frequency, effect mediators, and mitigators. Journal of the American "
    "Medical Informatics Association, 19(1), 121-127. "
    "https://doi.org/10.1136/amiajnl-2011-000089",

    "Guo, C., Pleiss, G., Sun, Y., & Weinberger, K. Q. (2017). On calibration of "
    "modern neural networks. Proceedings of Machine Learning Research, 70, 1321-1330. "
    "https://proceedings.mlr.press/v70/guo17a.html",

    "Hopp, W. J., & Spearman, M. L. (2011). Factory physics (3rd ed., chaps. 7-9). "
    "Waveland Press.",

    "IBM. (2020). AI maturity framework for enterprise applications [White paper]. "
    "IBM Corporation. https://www.ibm.com/downloads/cas/OB8M18WR",

    "IBM. (n.d.). University Hospitals Coventry and Warwickshire (UHCW) NHS Trust: "
    "Patient-centered care from AI-centered efficiencies [Case study]. "
    "https://www.ibm.com/case-studies/uhcw-nhs-trust",

    "National Institute of Standards and Technology. (2023). Artificial intelligence "
    "risk management framework (AI RMF 1.0) (NIST AI 100-1). U.S. Department of "
    "Commerce. https://doi.org/10.6028/NIST.AI.100-1",

    "Northcutt, C. G., Athalye, A., & Mueller, J. (2021). Pervasive label errors in "
    "test sets destabilize machine learning benchmarks. In Proceedings of the Neural "
    "Information Processing Systems Track on Datasets and Benchmarks (Vol. 1). "
    "https://arxiv.org/abs/2103.14749",

    "Parasuraman, R., & Manzey, D. H. (2010). Complacency and bias in human use of "
    "automation: An attentional integration. Human Factors, 52(3), 381-410. "
    "https://doi.org/10.1177/0018720810376055",

    "Press, G. (2016, March 23). Cleaning big data: Most time-consuming, least "
    "enjoyable data science task, survey says. Forbes. "
    "https://www.forbes.com/sites/gilpress/2016/03/23/data-preparation-most-"
    "time-consuming-least-enjoyable-data-science-task-survey-says/",

    "Rudin, C. (2019). Stop explaining black box machine learning models for high "
    "stakes decisions and use interpretable models instead. Nature Machine "
    "Intelligence, 1(5), 206-215. https://doi.org/10.1038/s42256-019-0048-x",

    "Sambasivan, N., Kapania, S., Highfill, H., Akrong, D., Paritosh, P., & Aroyo, "
    "L. M. (2021). \"Everyone wants to do the model work, not the data work\": Data "
    "cascades in high-stakes AI. In Proceedings of the 2021 CHI Conference on Human "
    "Factors in Computing Systems (Article 39, pp. 1-15). ACM. "
    "https://doi.org/10.1145/3411764.3445518",

    "Sculley, D., Holt, G., Golovin, D., Davydov, E., Phillips, T., Ebner, D., "
    "Chaudhary, V., Young, M., Crespo, J.-F., & Dennison, D. (2015). Hidden technical "
    "debt in machine learning systems. In Advances in Neural Information Processing "
    "Systems 28 (pp. 2503-2511). Curran Associates. "
    "https://proceedings.neurips.cc/paper/2015/hash/"
    "86df7dcfd896fcaf2674f757a2463eba-Abstract.html",
]

REF_FOOT = ("The company name has been changed for commercial in confidence reasons."
            "　內部數字均標示為假設或推算量級,並於各頁附推算方式。")

# ------------------------------------------------------------------
# 🔴 換行量測 —— 這一頁出過疊印,原因是「行數用猜的」
# ------------------------------------------------------------------
# 舊版把 15 條的行數寫死成常數,其中 4 條少算一行,於是下一條直接壓上去
# (Department of Industry 的 principles 行、Parasuraman 的 DOI 行、
#  Press 的 task-survey-says/ 行、Sculley 的 Abstract.html 行)。
# 改成用 PIL 載 Microsoft JhengHei、逐 token 模擬 PowerPoint 換行,建版時現算。
#
# 模擬規則是拿 PowerPoint 匯出的 PNG 反推、再逐條對回去校準出來的:
#   · 拉丁字以空白斷詞;CJK 每字可斷
#   · **連字號 "-" 之後可斷行,斜線 "/" 不可** —— 所以長網址不是在 / 斷,
#     而是整串移到新的一行後再硬切(humanrights 那條就是這樣變成 5 行的)
#   · 單一 token 比整行長時,填滿一行硬切,剩下的必定進下一行
# 校準結果:15 條 + 頁尾在 5.30 吋欄寬下,模擬行數與 PowerPoint 實測**完全一致**。
_MSJH = r"C:\Windows\Fonts\msjh.ttc"
_BREAK_AFTER = "-"
_FONT_CACHE = {}


def _font(size_pt):
    """以 1px = 1pt 載入 JhengHei,getlength() 回傳的就是點值。"""
    key = int(round(size_pt))
    if key not in _FONT_CACHE:
        from PIL import ImageFont
        _FONT_CACHE[key] = ImageFont.truetype(_MSJH, key, index=0)
    return _FONT_CACHE[key]


def _is_cjk(ch):
    o = ord(ch)
    return (0x2E80 <= o <= 0x9FFF or 0xF900 <= o <= 0xFAFF
            or 0xFF00 <= o <= 0xFF60 or 0x3000 <= o <= 0x303F)


def _tokens(text):
    """切成 PowerPoint 的斷行單位。"""
    out, buf = [], ""
    for ch in text:
        if ch == " ":
            if buf:
                out.append(buf)
                buf = ""
            out.append(" ")
        elif _is_cjk(ch):
            if buf:
                out.append(buf)
                buf = ""
            out.append(ch)
        else:
            buf += ch
            if ch in _BREAK_AFTER:
                out.append(buf)
                buf = ""
    if buf:
        out.append(buf)
    return out


def line_count(text, width_in, size_pt=12.0):
    """這段文字在 width_in 吋欄寬(內距已歸零)、size_pt 字級下會排成幾行。"""
    f = _font(size_pt)
    limit = width_in * 72.0
    lines, cur = 1, 0.0
    for tk in _tokens(text):
        w = f.getlength(tk)
        if tk == " ":
            if cur > 0:                       # 行首的空白吃掉
                cur += w
            continue
        if cur + w <= limit + 1e-6:
            cur += w
            continue
        if cur > 0:                           # 這個 token 整串移到下一行
            lines += 1
            cur = 0.0
        while w > limit + 1e-6:               # 比整行還長 -> 填滿一行硬切
            acc, i = 0.0, 0
            while i < len(tk):
                cw = f.getlength(tk[i])
                if acc + cw > limit + 1e-6:
                    break
                acc += cw
                i += 1
            i = max(i, 1)
            tk = tk[i:]
            w = f.getlength(tk)
            lines += 1
        cur = w
    return lines


# ------------------------------------------------------------------ 參考文獻版位
# 行高:PowerPoint 的行節距 = 字級 x 1.20 x 行距倍數(匯出 PNG 量到 0.2092 吋
# @ 12pt / 1.05,對上 12 x 1.20 x 1.05 / 72 = 0.2100)。單行的墨跡帶約 0.176 吋,
# 所以節距**必須 > 0.176** 才不會上下行的字互相碰到 —— 舊版的 0.88 行距剛好壓在
# 這條線上(0.176),那也是「條目黏成一片」的來源。這裡取 0.95 -> 節距 0.190 吋。
REF_SP = 0.95
REF_LH = 12.0 * 1.20 * REF_SP / 72.0        # 一行 = 0.1900 吋
# 段距 0.08 吋:框高 n x 節距 只涵蓋到末行的基線帶,末行的降部(p / g / y / _)
# 還會再往下約 0.013 吋,所以 0.08 的段距實得約 0.067 吋淨空 —— 匯出 PNG 實測
# 相鄰兩條的墨跡淨空 0.093-0.140 吋,肉眼分得開,這一頁的「黏成一片」由此解掉。
REF_GAP = 0.08
REF_TOP = 0.80                              # 標題 26pt 的墨跡收在 y 0.74

# 🔑 欄寬從 5.30 加寬到 5.90:總行數由 63 掉到 56,這是騰出段距的關鍵。
# 右欄 x+w = 12.50 越過右下淨空區左緣(11.283),所以**右欄底緣必須收在 6.35 以上**;
# 左欄不越界,可以一路排到頁尾條之前。兩欄等寬,右欄短、左欄長。
REF_W = 5.90
REF_X_L, REF_X_R = 0.55, 6.60
REF_SPLIT = 9                    # 前 9 條走左欄 —— 依字母序切,不打亂 APA 排序
REF_BOT_L, REF_BOT_R = 7.18, 6.35

REF_FOOT_X, REF_FOOT_W = 0.55, 10.70
REF_FOOT_Y, REF_FOOT_H = 7.22, 0.20


def _ref_blocks(color, foot_color):
    """由上往下累加 y,逐條現算行數 —— 不再有寫死的行數常數。"""
    blocks = []
    for col, (x, entries, first) in enumerate((
            (REF_X_L, REF_ENTRIES[:REF_SPLIT], 1),
            (REF_X_R, REF_ENTRIES[REF_SPLIT:], REF_SPLIT + 1))):
        y = REF_TOP
        for i, e in enumerate(entries):
            n = line_count(e, REF_W)
            h = n * REF_LH
            blocks.append(dict(id="REF-%s%02d" % ("LR"[col], first + i),
                               x=x, y=y, w=REF_W, h=h, text=e,
                               size=12, color=color, spacing=REF_SP, lines=n))
            y += h + REF_GAP
    blocks.append(dict(id="REF-FOOT", x=REF_FOOT_X, y=REF_FOOT_Y,
                       w=REF_FOOT_W, h=REF_FOOT_H, text=REF_FOOT,
                       size=12, color=foot_color, spacing=1.0,
                       lines=line_count(REF_FOOT, REF_FOOT_W)))
    return blocks


def assert_ref_layout(blocks):
    """把 place_slide_text 不檢查的那一項補上:框高夠不夠、框之間會不會疊到。"""
    # 1. 宣告高度必須裝得下實際行數
    for b in blocks:
        need = b["lines"] * (12.0 * 1.20 * b["spacing"] / 72.0)
        if b["h"] + 1e-6 < need:
            raise AssertionError(
                "%s 框高不足:宣告 %.4f 吋,%d 行實際需要 %.4f 吋"
                % (b["id"], b["h"], b["lines"], need))
    # 2. 兩兩不得相交(同欄比 y,跨欄 x 本來就分開,用矩形交集一次判掉)
    for i in range(len(blocks)):
        a = blocks[i]
        for j in range(i + 1, len(blocks)):
            c = blocks[j]
            if (a["x"] < c["x"] + c["w"] - 1e-6 and c["x"] < a["x"] + a["w"] - 1e-6
                    and a["y"] < c["y"] + c["h"] - 1e-6
                    and c["y"] < a["y"] + a["h"] - 1e-6):
                raise AssertionError(
                    "%s 與 %s 疊印:y %.3f-%.3f vs %.3f-%.3f"
                    % (a["id"], c["id"], a["y"], a["y"] + a["h"],
                       c["y"], c["y"] + c["h"]))
    # 3. 底緣:左欄收在頁尾條之前,右欄收在淨空區之上,頁尾條收在安全邊之內
    for b in blocks:
        if b["id"] == "REF-FOOT":
            lim = 7.44
        else:
            lim = REF_BOT_R if b["x"] > 3.0 else REF_BOT_L
        if b["y"] + b["h"] > lim + 1e-6:
            raise AssertionError("%s 底緣 %.3f 超過上限 %.3f"
                                 % (b["id"], b["y"] + b["h"], lim))
    return blocks


def ref_layout(ctx=None):
    """給驗證腳本用:不開 PowerPoint 就拿得到這一頁的完整版位。"""
    from pptx.dml.color import RGBColor
    dark = ctx["DARK"] if ctx else RGBColor(0x33, 0x33, 0x33)
    grey = ctx["GREY"] if ctx else RGBColor(0x80, 0x80, 0x80)
    return assert_ref_layout(_ref_blocks(dark, grey))


def _refs(prs, ctx):
    s = ctx["new_slide"](prs, 11, ctx=ctx)
    blocks = assert_ref_layout(_ref_blocks(ctx["DARK"], ctx["GREY"]))
    _zero_margins(ctx["place_slide_text"](s, blocks))
    ctx["speaker_notes"](s, "這些是本提案引用的公開文獻,完整清單在這一頁。")
    return s


# ==================================================================
def build(prs, ctx):
    """建立 P8-P10 與參考文獻頁(page_idx 8-11,遞增)。"""
    _p8(prs, ctx)
    _p9(prs, ctx)
    _p10(prs, ctx)
    _refs(prs, ctx)
