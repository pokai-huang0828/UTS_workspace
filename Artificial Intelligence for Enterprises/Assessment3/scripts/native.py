# -*- coding: utf-8 -*-
"""原生形狀引擎 —— 用 python-pptx 的真形狀與真文字框畫圖,不出 PNG。

為什麼要有這支:
    先前兩版投影片的圖都是 matplotlib 出 PNG 再貼進去。
    在 PowerPoint 裡那就是**一張死圖** —— Kenny 一個字都改不了,
    而他要的是「可以編輯的內容,包括裡面的圖字」。

    這裡的每一個方塊、每一條線、每一段文字都是原生物件:
    點下去可以改字、拖位置、換顏色,不需要回頭找 Python。

座標一律用「吋」,原點在左上角(與 PowerPoint 一致)。
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 色票(全片鎖定,不得自行取色)────────────────────────────────
NAVY = "#1F4E79"
ORANGE = "#E8833A"
RED = "#C00000"
# 🔴 2026-08-16 新增:第 8 頁的「改善」要跟「不變」在視覺上分得開。
#    深綠(不是螢光綠)—— 跟 NAVY 同一個飽和度層級,投影與錄影都讀得出來。
GREEN = "#1E7145"
GREY = "#6B7885"
LIGHT = "#DCE3EA"
PALE = "#EAF1F8"
INK = "#16212E"
WHITE = "#FFFFFF"
FONT = "Microsoft JhengHei"

TONE = {"navy": NAVY, "orange": ORANGE, "red": RED, "grey": GREY,
        "chosen": NAVY, "light": LIGHT, "ink": INK}

W, H = 13.333, 7.5
SAFE_R, SAFE_B = 2.05, 1.15          # 右下人像淨空區
SAFE_X, SAFE_Y = W - SAFE_R, H - SAFE_B
MIN_PT = 12.0                        # 投影片字級下限


def rgb(h):
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def _guard(x, y, w, h, who):
    # 🔴 尺寸必須為正。負高度的形狀 python-pptx 照寫不誤,zip 也是合法的,
    #    但 **PowerPoint 會直接拒絕開啟整個檔案**(「PowerPoint could not open the file」)。
    #    2026-08-14 我改版心分配時就做出過這種檔:第 9 頁 callout 高 0、
    #    第 10 頁三個文字框高 −0.17"。python-pptx 讀得開,所以驗證全過,
    #    是拿 PowerPoint COM 去導 PNG 才炸出來。
    #    → 這條 assert 就是那個缺掉的守門員:寧可 build 失敗,不可交出開不了的檔。
    assert w > 0.01 and h > 0.01, f"{who} 尺寸非正:w={w:.3f} h={h:.3f}(PowerPoint 會拒絕開檔)"
    assert -0.01 <= x and x + w <= W + 0.01, f"{who} 超出版面寬:{x:.2f}+{w:.2f}"
    assert -0.01 <= y and y + h <= H + 0.01, f"{who} 超出版面高:{y:.2f}+{h:.2f}"
    if x + w > SAFE_X + 1e-6 and y + h > SAFE_Y + 1e-6:
        raise AssertionError(
            f"{who} 侵入右下人像淨空區(x>{SAFE_X:.2f} 且 y>{SAFE_Y:.2f})")


# ══════════════════════════════════════════════════════════════════
# 基本元件
# ══════════════════════════════════════════════════════════════════
def text(slide, x, y, w, h, s, size=14, color=INK, bold=False,
         align="left", anchor="top", spacing=1.25, who="text"):
    """真文字框。s 可含 \\n;每一行都是一個可編輯段落。"""
    assert size >= MIN_PT - 1e-6, f"{who} 字級 {size}pt 低於下限 {MIN_PT}pt"
    _guard(x, y, w, h, who)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for i, line in enumerate(str(s).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = rgb(color)
    return tb


def box(slide, x, y, w, h, fill=None, line=None, lw=1.6, radius=True,
        who="box", dash=False):
    """真矩形(圓角可選)。fill / line 給 None 表示不填 / 不描邊。"""
    _guard(x, y, w, h, who)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        # 圓角半徑改小,預設太圓像按鈕
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(lw)
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rule(slide, x, y, w, color=NAVY, thick=0.035):
    """細分隔線(用扁矩形做,PowerPoint 裡好選好改)。"""
    return box(slide, x, y, w, thick, fill=color, line=None, radius=False,
               who="rule")


def arrow(slide, x, y, w, h=0.16, color=GREY):
    _guard(x, y, w, h, "arrow")
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def badge(slide, x, y, w, h, s, color, size=12):
    """狀態標籤(圓角小條 + 白字)。

    🔴 標籤是**一個詞**,不是一句話。規格常常塞進「提案推算,非核定編制」
       這種長句,兩行就撐爆膠囊 —— 一律 clamp 成一行,全文進備註。
    """
    box(slide, x, y, w, h, fill=color, line=None, who="badge")
    text(slide, x, y, w, h, clamp(str(s), w - 0.10, size, 1, "badge"),
         size=size, color=WHITE, bold=True,
         align="center", anchor="middle", who="badge-text")


# ══════════════════════════════════════════════════════════════════
# 文字量測 —— 版塊要按內容決定高度,不能用固定行高
# ══════════════════════════════════════════════════════════════════
def _wide(ch):
    """全形(CJK / 全形標點)算 1.0 個字寬,半形算 0.55。"""
    o = ord(ch)
    return 1.0 if (0x2E80 <= o <= 0x9FFF or 0xFF00 <= o <= 0xFF60
                   or 0x3000 <= o <= 0x303F) else 0.55


def lines_needed(s, w_in, size_pt):
    """在寬 w_in 吋、字級 size_pt 下,這段文字要幾行。"""
    if not s:
        return 0
    per_line = max(w_in * 72.0 / size_pt, 1.0)      # 一行放得下幾個全形字
    n = 0
    for para in str(s).split(chr(10)):
        u = sum(_wide(c) for c in para)
        n += max(1, int(u / per_line) + (1 if u % per_line else 0))
    return n


def text_w(s, size_pt):
    """這段文字排成**一行**要多寬(吋)。用來把欄寬收到剛好夠。"""
    return max(sum(_wide(c) for c in str(s).replace(chr(10), "")), 1) * size_pt / 72.0


SAFETY = 1.35          # 量測安全係數


PINNED_SAFETY = 1.15   # 行數已被 clamp 釘死時用的係數(見下)


def text_h(s, w_in, size_pt, spacing=1.25, safety=None):
    """這段文字需要的高度(吋)。

    🔴 乘上 SAFETY。lines_needed 是估算(CJK 1.0 / 半形 0.55 的近似字寬),
       單一區塊誤差很小,但一頁疊三四個版塊後會累積成溢出 ——
       實測 1.18 仍會爆,調到 1.35 才乾淨。
       **寧可留白,不可重疊**:重疊在投影片上是致命的,留白只是不好看。

    🔑 但 SAFETY 保的是「**這段文字會佔幾行**」這個估計的誤差。
       文字已經被 clamp 釘死在 N 行之後,行數不再是估計值,
       這時再乘 1.35 就是重複計價 —— 表格因此每列虛胖兩成,
       10 條風險量成 6.59" 塞不進 4.19",然後默默丟掉 4 條。
       所以呼叫端可以傳 safety=PINNED_SAFETY,只留字體度量的餘裕。
    """
    return (lines_needed(s, w_in, size_pt) * size_pt * spacing / 72.0
            * (SAFETY if safety is None else safety))


OVERFLOW = []          # 被截斷或未上版的完整內容,由渲染器收進講者備註
PAGE = 0               # 目前渲染到第幾頁(只為了讓警告訊息指得出頁碼)


def fit_items(items, avail_h, measure, min_keep=2, who=""):
    """在 avail_h 吋內,能放幾項就放幾項;放不下的交給講者備註。

    🔑 這是本引擎的核心設計判斷。
       先前的做法是二選一 —— 內容全塞(爆版)或整批砍掉(空洞)。
       兩者都是用猜的。正確做法是**讓量測決定那條線在哪**:
       同一份規格,換版面、換字級都會自動重算,不需要人再判斷一次。

    🔴 **順序不等於重要性。** 第一版只留「前 N 項」,結果 P7 把 R9 / R10
       (這個方案自己造成的兩條新風險,而且底部紅條正在指名它們)砍掉,
       留下了 R1–R5。所以要挑的是**重要的**,不是**排前面的** ——
       標了 hot 的優先保留,渲染時再依原順序排回去。

    items    : list
    avail_h  : 可用高度(吋)
    measure  : callable(item) -> 該項需要的高度(吋)
    min_keep : 至少保留幾項(即使超高)—— 一項都不放的版面沒有意義
    回傳     : (shown, hidden),shown 已還原成原始順序
    """
    def _hot(it):
        return bool(it.get("hot")) if isinstance(it, dict) else False

    order = sorted(range(len(items)), key=lambda i: (not _hot(items[i]), i))
    keep, used = set(), 0.0
    for i in order:
        h = measure(items[i])
        if keep and used + h > avail_h + 1e-6 and len(keep) >= min_keep:
            continue                      # 這一項放不下,但後面較小的也許放得下
        keep.add(i)
        used += h
    shown = [items[i] for i in sorted(keep)]
    hidden = [items[i] for i in range(len(items)) if i not in keep]
    for it in hidden:
        OVERFLOW.append((who, str(it)))
    if hidden:
        # 🔴 大聲報。matrix 已經有這條,但 rows / compare / cards 一直是**默默**丟 ——
        #    第 2 頁的 compare 規格三行只印一行(而漏掉的是「營運端的實務估計,
        #    尚未以計時量測驗證」這條必附限定語),rows 三列只印兩列,
        #    我完全不知道,是逐頁審到第 2 頁才發現。
        def _lbl(it):
            if isinstance(it, dict):
                return str(it.get("label") or it.get("name") or it)[:20]
            return str(it)[:20]
        print(f'   ⚠️ 第 {PAGE} 頁 {who} 丟了 {len(hidden)} / {len(items)} 項'
              f'(可用 {avail_h:.2f}"):' + '、'.join(_lbl(it) for it in hidden))
    return shown, hidden


def more_marker(slide, x, y, w, n, who="more"):
    """「另 N 項見備註」——讓觀眾知道有東西被收起來,而不是憑空消失。"""
    if n <= 0:
        return
    text(slide, x, y, w, 0.28, f"另 {n} 項見講者備註", 12, GREY,
         align="right", who=who)


def clamp(s, w_in, size_pt, max_lines, who=""):
    """超過 max_lines 就截斷,完整內容存進 OVERFLOW 交給講者備註。

    🔑 這樣做的理由:表格欄位塞進整段句子會overlap,
       但那些句子有內容不能丟 —— 投影片留可掃描的版本,全文進備註。
    """
    if lines_needed(s, w_in, size_pt) <= max_lines:
        return s
    s = str(s)
    # 🔴 2026-08-14:硬換行必須先攤平,否則 clamp 根本壓不到 max_lines。
    #    「A\nB」壓成 1 行,舊版只按字寬截斷,\n 原封不動留著 → 還是 2 行。
    #    於是 blk_matrix 的 3→2→1 降階跑到 1 也沒用,表格照樣塞不下、照樣丟列。
    #    (第 7 頁 10 條風險需 6.59" 只有 4.19" —— 每列實測仍是 2.1 行。)
    if s.count(chr(10)) + 1 > max_lines:
        s = "　".join(p.strip() for p in s.split(chr(10)) if p.strip())
    per_line = max(w_in * 72.0 / size_pt, 1.0)
    budget = int(per_line * max_lines) - 1
    acc, out = 0.0, []
    for c in str(s):
        acc += _wide(c)
        if acc > budget:
            break
        out.append(c)
    OVERFLOW.append((who, str(s)))
    return "".join(out).rstrip("　,、;·") + "…"


# ══════════════════════════════════════════════════════════════════
# 版塊 —— 六種,對應設計稿的 block kind
# ══════════════════════════════════════════════════════════════════
def blk_cards(slide, x, y, w, h, items):
    """横排卡片。items = [{title, badge, tone, lines[]}]"""
    n = len(items)
    gap = 0.16
    cw = (w - gap * (n - 1)) / n
    for i, it in enumerate(items):
        cx = x + i * (cw + gap)
        tone = TONE.get(it.get("tone", "navy"), NAVY)
        chosen = it.get("tone") == "chosen"
        box(slide, cx, y, cw, h, fill=(PALE if chosen else WHITE),
            line=tone, lw=(3.0 if chosen else 1.6), who=f"card{i}")
        text(slide, cx + 0.12, y + 0.14, cw - 0.24, 0.62, it["title"],
             size=15, color=INK, bold=True, align="center", who=f"card{i}-t")
        by = y + 0.82
        if it.get("badge"):
            # 🔴 膠囊寬度要**跟著字長走**,不能寫死 1.35"。
            #    1.35" 在 12pt 只放得下 8 個全形字,「提案推算,非核定編制」是 10 個
            #    → clamp 成「提案推算,非…」印在片上,那是一句沒有意義的話。
            need_w = lines_needed(it["badge"], 99, 12) and \
                sum(_wide(c) for c in str(it["badge"])) * 12 / 72.0 + 0.28
            bw = max(0.9, min(cw - 0.20, need_w))
            badge(slide, cx + (cw - bw) / 2, by, bw, 0.30, it["badge"], tone)
            by += 0.42
        if it.get("lines"):
            room = h - (by - y) - 0.14
            if room < 0.24:
                # 🔴 同 blk_compare:高度不足就整批進備註,絕不畫負高度的框
                #    (第 10 頁實測 room = −0.64",那會讓 PowerPoint 拒絕開檔)
                for ln in it["lines"]:
                    OVERFLOW.append((f"card:{it['title']}", str(ln)))
                print(f'   ⚠️ 第 {PAGE} 頁 card:{it["title"][:16]} 高度不足,'
                      f'{len(it["lines"])} 行全部進備註')
                continue
            # 🔑 卡片內的說明行走同一條規則:放得下幾行就放幾行,其餘進備註
            ml = max(1, int(room * 72.0 / (12 * 1.35 * SAFETY)))
            lines = [clamp(ln, cw - 0.24, 12, ml, f"card:{it['title']}")
                     for ln in it["lines"]]
            shown, hidden = fit_items(
                lines, room,
                lambda s, _w=cw: text_h(s, _w - 0.24, 12, 1.35),
                min_keep=1, who=f"card:{it['title']}")
            text(slide, cx + 0.12, by, cw - 0.24, room, chr(10).join(shown),
                 size=12, color=GREY, align="center",
                 spacing=1.35, who=f"card{i}-l")
            # 🔴 不在頁面上放「另 N 項」標記 —— 它一直壓到內文。
            #    被收起來的內容在講者備註裡,錄影時你看得到。


def blk_bar(slide, x, y, w, h, spec):
    """分段長條 + 括弧標註。spec = {label, segments[{name,value,weight,hot}], bracket, note}"""
    if spec.get("label"):
        text(slide, x, y, w, 0.38, spec["label"], size=17, color=INK, bold=True,
             who="bar-label")
    by = y + (0.48 if spec.get("label") else 0)
    bh = 0.62
    tot = sum(s.get("weight", 1) for s in spec["segments"])
    cx = x
    hot_x0 = hot_x1 = None
    for s in spec["segments"]:
        sw = w * s.get("weight", 1) / tot
        hot = s.get("hot")
        box(slide, cx, by, sw, bh, fill=(NAVY if hot else LIGHT), line=WHITE,
            lw=1.5, radius=False, who="seg")
        text(slide, cx, by + 0.05, sw, 0.26, s["name"], size=12,
             color=(WHITE if hot else INK), bold=True, align="center", who="seg-n")
        if s.get("value"):
            text(slide, cx, by + 0.31, sw, 0.26, s["value"], size=12,
                 color=(WHITE if hot else GREY), align="center", who="seg-v")
        if hot:
            hot_x0 = cx if hot_x0 is None else hot_x0
            hot_x1 = cx + sw
        cx += sw
    yy = by + bh + 0.10
    if spec.get("bracket") and hot_x0 is not None:
        b = spec["bracket"]
        rule(slide, hot_x0, yy, hot_x1 - hot_x0, color=NAVY, thick=0.03)
        # 🔴 big 可能是**布林旗標**(「這一行要放大」)或字串(要另外放大的那句)。
        #    先前一律當字串印,結果 "big": true 就在投影片中央印出一個大大的 True。
        big = b.get("big")
        big_str = big if isinstance(big, str) else None
        emph = bool(big)
        text(slide, hot_x0, yy + 0.10, hot_x1 - hot_x0, 0.44, b["text"],
             size=(17 if emph and not big_str else 14), color=NAVY, bold=True,
             align="center", who="brk")
        if big_str:
            text(slide, hot_x0, yy + 0.58, hot_x1 - hot_x0, 0.60, big_str,
                 size=30, color=NAVY, bold=True, align="center", who="brk-big")
            yy += 1.24
        else:
            yy += 0.62
    if spec.get("note"):
        # 🔴 2026-08-16:舊版寫 min(yy, y + h - 0.30) —— 版塊高度不夠時,
        #    這個 min 會把註記**往上夾**,直接壓在括號說明那一行上面。
        #    第 3 頁就是這樣:「合計 2.5–5.5 天」印在「AI 可處理段…」正上方,
        #    兩行疊在一起。Kenny 在 PowerPoint 裡手動把它往下拖 0.34" 才看得清楚。
        #    正解是**不要夾**:放在該放的位置,高度不夠就讓 _guard 擋下來,
        #    再由 need_h / MIN 去把版塊撐開(bar 的下限已同步調到 2.15")。
        text(slide, x, yy, w, 0.30, spec["note"], size=12,
             color=GREY, who="bar-note")


def blk_rows(slide, x, y, w, h, items):
    """左標右值的資料列。items = [{label, value, note, hot}]

    🔴 行高由**量測內容**決定,不用固定值。註記欄放得下多少就放多少,
       放不下的整段進講者備註(clamp),不讓它壓到下一列。
    """
    # 🔴 2026-08-16:值欄寬度改成**跟著內容走**,不再固定 30%。
    #    Kenny 在 PowerPoint 裡手動把第 10 頁的說明欄從 4.64" 拉到 7.92" ——
    #    他是對的:那一頁的值欄固定吃掉 3.45",裡面只放「W1–6」四個字,
    #    而旁邊的說明欄擠到必須截斷。固定比例在值短的時候就是在浪費版面。
    LW = w * 0.24
    vmax = max([text_w(str(it["value"]), 15) for it in items] or [0]) + 0.24
    VW = max(min(w * 0.30, vmax), w * 0.10)
    NW = w - LW - VW
    NX = x + LW + VW
    PAD = 0.10
    max_note_lines = 3

    def _need(it):
        return max(text_h(it["label"], LW - 0.24, 14),
                   text_h(it["value"], VW - 0.14, 15),
                   text_h(clamp(it.get("note", ""), NW - 0.14, 12, max_note_lines),
                          NW - 0.14, 12)) + 0.10 + PAD

    # 🔑 放得下幾列就放幾列,其餘進備註 —— 不猜、不整批砍
    items, _hidden = fit_items(items, h, _need, min_keep=2, who="rows")
    n = max(len(items), 1)
    budget = (h - PAD * (n - 1)) / n
    ry = y
    for i, it in enumerate(items):
        c = RED if it.get("hot") else NAVY
        note = it.get("note", "")
        if note:
            note = clamp(note, NW - 0.14, 12, max_note_lines, f"rows[{i}]{it['label']}")
        need = max(
            text_h(it["label"], LW - 0.24, 14),
            text_h(it["value"], VW - 0.14, 15),
            text_h(note, NW - 0.14, 12) if note else 0,
        ) + 0.10
        rh = max(min(need, budget), 0.34)
        box(slide, x, ry + 0.03, 0.05, rh - 0.10, fill=c, line=None,
            radius=False, who=f"row{i}-tick")
        text(slide, x + 0.16, ry, LW - 0.22, rh, it["label"],
             size=14, color=INK, bold=True, anchor="middle", who=f"row{i}-l")
        text(slide, x + LW, ry, VW - 0.12, rh,
             clamp(it["value"], VW - 0.14, 15, 2, f"rows-value[{it['label']}]"),
             size=15, color=c, bold=True, anchor="middle", who=f"row{i}-v")
        if note:
            text(slide, NX, ry, NW - 0.12, rh, note,
                 size=12, color=GREY, anchor="middle", spacing=1.2,
                 who=f"row{i}-n")
        ry += rh + PAD


def blk_compare(slide, x, y, w, h, spec):
    """左右兩欄對照 + 一句裁決。"""
    gap = 0.30
    cw = (w - gap) / 2
    for i, side in enumerate(("left", "right")):
        s = spec[side]
        cx = x + i * (cw + gap)
        tone = TONE.get(s.get("tone", "navy"), NAVY)
        vh = h - (0.48 if spec.get("verdict") else 0)
        box(slide, cx, y, cw, vh, fill=WHITE, line=tone, lw=2.0, who=f"cmp{i}")
        text(slide, cx + 0.16, y + 0.16, cw - 0.32, 0.42, s["title"], size=17,
             color=tone, bold=True, align="center", who=f"cmp{i}-t")
        room = vh - 0.90
        if room < 0.30:
            # 🔴 高度不足以放任何一行 —— 只印標題與裁決,絕不產出負高度的文字框
            #    (負高度會讓 PowerPoint 拒絕開啟整個檔案)。
            for ln in s["lines"]:
                OVERFLOW.append((f"compare:{s['title']}", str(ln)))
            print(f'   ⚠️ 第 {PAGE} 頁 compare:{s["title"][:16]} 高度不足,'
                  f'{len(s["lines"])} 行全部進備註')
            continue
        shown, hidden = fit_items(
            s["lines"], room,
            lambda ln, _w=cw: text_h(ln, _w - 0.32, 13, 1.5),
            min_keep=1, who=f"compare:{s['title']}")
        text(slide, cx + 0.16, y + 0.70, cw - 0.32, room, chr(10).join(shown),
             size=13, color=INK, spacing=1.5, align="center", who=f"cmp{i}-l")
    if spec.get("verdict"):
        text(slide, x, y + h - 0.40, w, 0.38, spec["verdict"], size=15,
             color=INK, bold=True, align="center", who="cmp-v")


def blk_matrix(slide, x, y, w, h, spec):
    """真 PowerPoint 表格(GraphicFrame),不是用文字框排出來的假表。

    🔑 為什麼改用真表格:
       ① Kenny 可以在 PowerPoint 裡增刪列、拉欄寬 —— 假表做不到
       ② **列高由 PowerPoint 自己算**,我不必再手動量測換行
          (先前 matrix 的重疊,根因就是我算的列高跟實際渲染對不上)
       ③ 表格樣式一致,不會每一格各自為政
    """
    cols = list(spec["cols"])
    rows = list(spec["rows"])
    name_w = w * spec.get("name_frac", 0.30)
    cw = (w - name_w) / max(len(cols), 1)

    # 🔴 欄數對齊:規格常常 cols 有 4 個但 rows 只給 3 格 → 右邊多一個空欄。
    #    少的補空字串,多的截掉。
    for r in rows:
        cs = list(r["cells"])[:len(cols)]
        r["cells"] = cs + [""] * (len(cols) - len(cs))

    # 🔴 整欄都空的直接丟掉。規格給 4 個欄名卻只有 3 格資料時,
    #    補空字串會留下一整條空白欄 —— 那比沒有這一欄更糟,
    #    委員會會以為那裡的東西被漏掉了。
    keep_c = [j for j in range(len(cols))
              if any(str(r["cells"][j]).strip() for r in rows)]
    if len(keep_c) < len(cols):
        cols = [cols[j] for j in keep_c]
        for r in rows:
            r["cells"] = [r["cells"][j] for j in keep_c]
        cw = (w - name_w) / max(len(cols), 1)

    # 🔴 PowerPoint 表格**會自己長高** —— add_table 給的列高是最小值不是上限,
    #    內容放不下它就撐開,於是壓到下一個版塊。
    #    所以先把每一格 clamp 到固定行數,列高就可預測。
    #    行數已由 clamp 釘死,所以這裡用 PINNED_SAFETY 而不是 SAFETY(見 text_h)。
    #    ⚠️ 檢查用的列高必須跟**渲染用的**一模一樣。舊版檢查用 _rough、
    #    渲染用 max(_rough, 0.32),於是「量起來放得下」的表格畫出來卻超高。
    def _rough(r):
        return max(0.32, max(
            [text_h(str(r["name"]), name_w - 0.16, 12, safety=PINNED_SAFETY)]
            + [text_h(str(c), cw - 0.16, 12, safety=PINNED_SAFETY)
               for c in r["cells"]]) + 0.06)

    HDR = 0.36
    raw = [dict(r) for r in rows]

    # 🔴 行數上限自動降階。3 行放不下就試 2 行、再試 1 行,直到全部列進得了 h。
    #
    # 🔴🔴 2026-08-14 修掉一個把整份交付檔弄成半成品的 bug:
    #    舊版在迴圈**裡面**就呼叫 fit_items,而 fit_items 會先丟列讓它「放得下」,
    #    於是後面那句 `if 總高 <= h: break` **恆為真** —— 第一圈就 break,
    #    CELL_LINES 2 與 1 從來沒被跑過。降階機制整段是死碼。
    #    後果:第 6 頁標題寫「五條路」只印 2 列、第 7 頁「十條風險」只印 6 列、
    #    第 8 頁「琥珀 7 紅 1」只印 6 列、第 9 頁「①–⑨ 四道護欄」只印 2+2 列。
    #    **投影片自己承諾的條數對不上自己印出來的條數**,比留白難看一百倍。
    #
    # 🔑 正確順序:**先壓字,再考慮丟列**。
    #    丟列是最後手段 —— 一列都不丟能放下,就一列都不丟。
    def _clamped(cell_lines):
        out = [dict(r) for r in raw]
        for r in out:
            r["name"] = clamp(str(r["name"]), name_w - 0.16, 12, cell_lines,
                              f"matrix-name:{r['name']}")
            r["cells"] = [clamp(str(c), cw - 0.16, 12, cell_lines, "matrix-cell")
                          for c in r["cells"]]
        return out

    # ⚠️ clamp() 每截斷一次就往 OVERFLOW 塞一筆。試 3 個級距會塞三份重複的,
    #    所以先記位置,選定後倒回去只留最終那一版。
    mark = len(OVERFLOW)
    rows, chosen = None, 1
    for CELL_LINES in (3, 2, 1):
        cand = _clamped(CELL_LINES)
        rows, chosen = cand, CELL_LINES               # 記住目前最緊的一版
        if HDR + sum(_rough(r) for r in cand) <= h + 1e-6:
            break                                     # 全部列都放得下 → 收工,不丟任何一列
    del OVERFLOW[mark:]
    rows = _clamped(chosen)
    need = HDR + sum(_rough(r) for r in rows)
    if need > h + 1e-6:
        # 連 1 行/格都塞不下,這時才丟列(並且進講者備註)
        rows, _hidden = fit_items(rows, h - HDR, _rough, min_keep=2, who="matrix")
        # 🔴 大聲報出來。上一版就是**默默**丟列,結果標題寫「十條風險」片上只有 6 條,
        #    而我完全不知道 —— 是模擬考官翻頁對照才發現的。
        print(f'   ⚠️ 第 {PAGE} 頁 表格丟了 {len(_hidden)} / {len(raw)} 列'
              f'(需 {need:.2f}" · 只有 {h:.2f}" · 已壓到 {chosen} 行/格):'
              + '、'.join(str(r["name"]).split(chr(10))[0][:14] for r in _hidden))
    n = len(rows)

    _guard(x, y, w, h, "matrix")
    gf = slide.shapes.add_table(n + 1, len(cols) + 1,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False

    tbl.columns[0].width = Inches(name_w)
    for j in range(len(cols)):
        tbl.columns[j + 1].width = Inches(cw)
    # 列高照**各列自己的需要**設,不用平均值 —— 平均值會讓內容多的列被撐開,
    # 而撐開的量正好就是壓到下一個版塊的量。
    # 🔴 分到的高度用不完就會在版面留一塊洞。表格是唯一「給多少都用得完」的版塊 ——
    #    多的平均分到每一列,列距鬆開反而更好讀。
    #    但要設上限:兩列的表格如果把 2" 全吃下去,會變成兩條巨大的色帶。
    #    每列最多加 0.34"(約一行 12pt 的高度),超過的寧可留白。
    bonus = min(max(h - need, 0.0) / max(n, 1), 0.34) if n else 0.0

    tbl.rows[0].height = Inches(HDR)
    for i, r in enumerate(rows):
        tbl.rows[i + 1].height = Inches(_rough(r) + bonus)

    def _cell(c, s, size, color, bold=False, fill=None, align="center"):
        c.text = ""
        c.margin_left = c.margin_right = Emu(45720)      # 0.05"
        c.margin_top = c.margin_bottom = Emu(18288)      # 0.02"
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill:
            c.fill.solid(); c.fill.fore_color.rgb = rgb(fill)
        else:
            c.fill.background()
        para = c.text_frame.paragraphs[0]
        para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}[align]
        run = para.add_run(); run.text = str(s)
        run.font.size = Pt(size); run.font.bold = bold
        run.font.name = FONT; run.font.color.rgb = rgb(color)

    # 表頭
    _cell(tbl.cell(0, 0), spec.get("corner", ""), 12, WHITE, True, NAVY, "left")
    for j, cname in enumerate(cols):
        _cell(tbl.cell(0, j + 1), cname, 12, WHITE, True, NAVY)

    # 內容
    # 🔑 逐列可以指定 tone:"red" / "green",沒指定就沿用舊的 hot 布林。
    #    第 8 頁八條倫理原則要讓「改善」跟「不變」一眼分得開 ——
    #    全部同一個灰,委員會分不出哪幾條是這個案子帶來的好處。
    ROW_TONE = {"red": ("#FBEFEF", RED, RED),
                "green": ("#EDF6EF", GREEN, GREEN)}
    for i, r in enumerate(rows):
        tone = r.get("tone") or ("red" if r.get("hot") else "")
        bg, col, sub = ROW_TONE.get(tone, (WHITE, INK, GREY))
        strong = tone in ROW_TONE
        _cell(tbl.cell(i + 1, 0), r["name"], 12, col, strong, bg, "left")
        for j, c in enumerate(r["cells"]):
            _cell(tbl.cell(i + 1, j + 1), c, 12, sub, False, bg)
    return gf


def blk_callout(slide, x, y, w, h, spec):
    """強調條。

    🔴 callout 是**一句話**,不是段落。規格常常塞進整段,
       而它的高度有天花板(給再多空間也只該用兩三行)——
       所以一定要 clamp,否則字會爆出色塊、壓到隔壁的版塊。
    """
    tone = TONE.get(spec.get("tone", "navy"), NAVY)
    box(slide, x, y, w, h, fill=tone, line=None, who="callout")
    max_lines = max(1, int(h * 72.0 / (17 * 1.25 * SAFETY)))
    body = clamp(spec["text"], w - 0.40, 17, max_lines, "callout")
    text(slide, x + 0.20, y, w - 0.40, h, body, size=17, color=WHITE,
         bold=True, align="center", anchor="middle", who="callout-t")


RENDER = {"cards": blk_cards, "bar": blk_bar, "rows": blk_rows,
          "compare": blk_compare, "matrix": blk_matrix, "callout": blk_callout}
