# -*- coding: utf-8 -*-
"""印出交付檔裡每一個文字框/表格欄的**實際**可用寬度。

🔴 為什麼要有這支:我在改第 6 頁時,用手推的欄寬去量文字 ——
   卡片我算成 5.625",實際只有 4.961"(差 12%)。
   後果是「全客戶一天 13,000 筆 = 129.1–185.9 人時 ÷ 8 小時 = 17–24 人」
   在我的算法裡是一行、實際渲染是兩行,把第三行擠進備註。
   同一輪還發生過標題只超出 0.02" 就折行壓到副標。

🔑 教訓:**版面寬度不要用推的,去量已經畫出來的那個框。**
   這支就是「量已經畫出來的那個框」。

用法:python scripts/widths.py [頁碼]
"""
import os
import sys

from pptx import Presentation

A3 = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(A3, "Huang_26254793_421104_Assessment 3.pptx")
EMU = 914400.0


def main():
    want = {int(a) for a in sys.argv[1:] if a.isdigit()}
    prs = Presentation(DECK)
    for i, s in enumerate(prs.slides, 1):
        if want and i not in want:
            continue
        print(f"── 第 {i} 頁")
        for sh in s.shapes:
            if sh.has_table:
                t = sh.table
                ws = [c.width / EMU for c in t.columns]
                # 儲存格左右內距各 0.05"(見 native._cell)
                print(f'   表格 {len(t.rows)}×{len(t.columns)} 欄寬(扣內距):'
                      + " · ".join(f"{w - 0.10:.2f}" for w in ws))
            elif sh.has_text_frame and sh.text_frame.text.strip():
                head = sh.text_frame.text.split(chr(10))[0][:24]
                print(f'   文字框 {sh.width / EMU:5.2f}" × {sh.height / EMU:4.2f}"  {head}')


if __name__ == "__main__":
    main()
