# -*- coding: utf-8 -*-
"""把交付檔(片上文字 + 講者備註)匯成純文字,給稽核用的 agent 讀。

🔴 這支存在的唯一理由是一個真實事故。
   2026-08-16 我用臨時 python 匯出全片交給答辯題庫的 agent,那段程式只走
   `shape.has_text_frame`。但 blk_matrix 產生的是**真正的 PowerPoint 表格**
   (GraphicFrame),它的儲存格文字不在 text_frame 裡 ——
   於是第 4、5、6、7、8、9、10 頁的表格內容,**55 行片上文字,一行都沒匯出去**。

   後果不是「少一點資訊」,是 agent 反過來斷言那些內容不存在:
     · 「『人推翻機器判定的比例』全份投影片沒有」→ 第 7 頁風險 1 的偵測訊號就是它
     · 「解析度 640×360、1 FPS 片上完全沒有」    → 第 5 頁表格印著
     · 「『沒到 100% 不交回營運』片上沒有」        → 第 9 頁 ⑪ 就是這句
   整份答辯題庫因此帶著三條假禁語交出來,而那是要拿去現場用的東西。

   教訓:**匯出程式漏掉的內容,下游看起來跟「這個內容不存在」一模一樣。**
   所以匯出這件事不能用臨時程式,要有一支會自我檢查的。

用法:
    python scripts/dump_deck.py
輸出 notes/_deck_dump.txt,並印出覆蓋率自檢。
"""
import os
import sys

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
A3 = os.path.dirname(HERE)
DECK = os.path.join(A3, "Huang_26254793_421104_Assessment 3.pptx")
OUT = os.path.join(A3, "notes", "_deck_dump.txt")

LOCKED = {1, 2, 3, 4, 5}          # Kenny 逐頁核可、片上文字不得更動的頁
GROUP = 6                          # MSO_SHAPE_TYPE.GROUP


def walk(shape, out, seen):
    """把一個 shape 的所有文字收進 out。

    🔴 三條路都要走,少一條就是上面那個事故:
       ① group  —— 遞迴,群組本身沒有文字
       ② table  —— 文字在 cell 裡,**不在 text_frame**
       ③ 其餘   —— text_frame
    """
    seen["total"] += 1
    if shape.shape_type == GROUP:
        seen["group"] += 1
        for child in shape.shapes:
            walk(child, out, seen)
        return
    if shape.has_table:
        seen["table"] += 1
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" ｜ ".join(cells))
        return
    if shape.has_text_frame:
        t = shape.text_frame.text
        if t.strip():
            seen["text"] += 1
            out.append(t)
        return
    seen["silent"] += 1            # 線條、色塊 —— 本來就沒有文字


def main():
    prs = Presentation(DECK)
    seen = dict(total=0, group=0, table=0, text=0, silent=0)
    L = []
    for i, s in enumerate(prs.slides, 1):
        L.append("=" * 78)
        L.append(f"### 投影片 {i} / {len(prs.slides)}"
                 + ("  🔒已鎖定(Kenny 核可,不得更動片上文字)" if i in LOCKED else ""))
        L.append("--- 片上文字 ---")
        body = []
        for sh in s.shapes:
            walk(sh, body, seen)
        L.extend(body)
        L.append("--- 講者備註 ---")
        L.append(s.notes_slide.notes_text_frame.text if s.has_notes_slide else "(無)")
    doc = chr(10).join(L)
    open(OUT, "w", encoding="utf-8").write(doc)

    print("✅ " + os.path.relpath(OUT, A3) + f"  ({len(doc)} 字元)")
    print(f"   形狀 {seen['total']} 個 = 文字 {seen['text']} · 表格 {seen['table']}"
          f" · 群組 {seen['group']} · 無字 {seen['silent']}")

    # 🔑 自檢:交付檔裡每一格表格文字,都必須真的出現在匯出檔裡。
    #    這一段就是為了讓上面那個事故不可能再無聲發生。
    missing = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        t = c.text.strip()
                        if len(t) > 4 and t not in doc:
                            missing.append((i, t))
    if missing:
        print(f"   🔴 自檢失敗 —— {len(missing)} 格表格文字沒進匯出檔:")
        for i, t in missing[:10]:
            print(f"      第 {i} 頁「{t[:40]}」")
        sys.exit(1)
    print("   ✅ 自檢:所有表格儲存格文字都在匯出檔裡")


if __name__ == "__main__":
    main()
