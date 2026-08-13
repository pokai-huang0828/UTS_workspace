# -*- coding: utf-8 -*-
"""精簡版投影片 —— 重新生成的一份,圖與字全部重畫,不沿用 fig_v4_*。

為什麼要另做一份(而不是改舊的):
    舊版是「顧問式資訊圖表」——資訊密度高、圖佔滿版面、圖內字被迫壓到 9pt。
    它適合當 leave-behind 文件,不適合一支十分鐘、每頁只停 50 秒的影片。
    作業原文寫明「我們只會對視頻進行評分,幻燈片僅供澄清」——
    所以螢幕的工作不是承載論證(那是口白的事),是**不要跟口白搶注意力**。

本版的設計約束(與舊版相反,刻意的):
    · 每頁 ≤ 3 個要點,一個主張一張圖
    · 圖只畫一個概念,不畫表格
    · **字級下限 14pt**(舊版圖內是 9pt)—— 內容少了才放得大
    · 大量留白;右下 2.05 x 1.15 吋留給人像
    · 口白逐字全部灌進講者備註,錄影時用簡報者檢視

輸出:Huang_26254793_421104_Assessment 3 (slim).pptx —— 與正式版並存,方便對照。
"""
import os
import re
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, FancyArrow
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
A3 = os.path.dirname(HERE)
FIGDIR = os.path.join(A3, "figures_slim")
OUT = os.path.join(A3, "Huang_26254793_421104_Assessment 3 (slim).pptx")
os.makedirs(FIGDIR, exist_ok=True)

# ── 版面 ──────────────────────────────────────────────────────────
W, H = 13.333, 7.5
SAFE_R, SAFE_B = 2.05, 1.15          # 右下人像淨空區
M = 0.75                             # 版心左右邊界
TITLE_Y = 0.52
BODY_Y = 1.62                        # 要點起點
FIG_Y = 1.82                         # 圖起點(緊貼標題線下方)
FIG_H = 4.45                         # 圖高;底緣 6.27 < 淨空區 6.35

NAVY = "#1F4E79"
ORANGE = "#E8833A"
RED = "#C00000"
GREY = "#6B7885"
LIGHT = "#DCE3EA"
INK = "#16212E"
FONT = "Microsoft JhengHei"

plt.rcParams["font.sans-serif"] = [FONT, "Microsoft YaHei", "SimHei"]
plt.rcParams["axes.unicode_minus"] = False      # U+2212 缺字,一律關掉

_P = lambda h: RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def newfig(w=11.6, h=FIG_H):
    fig = plt.figure(figsize=(w, h), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    return fig, ax


def save(fig, name):
    p = os.path.join(FIGDIR, name + ".png")
    fig.savefig(p, dpi=200, bbox_inches="tight", pad_inches=0, facecolor="white")
    plt.close(fig)
    return p


# ══════════════════════════════════════════════════════════════════
# 圖 —— 每張只講一件事
# ══════════════════════════════════════════════════════════════════
def f_pipeline():
    """P1:三階段 + 兩單價的 16 倍落差。"""
    fig, ax = newfig()
    xs = [0.10, 0.34, 0.58]
    labs = ["車上鏡頭\n初判", "雲端\n複判", "5 位專職分析\n人工複核"]
    cols = [LIGHT, LIGHT, ORANGE]
    tcol = [INK, INK, "white"]
    for x, l, c, tc in zip(xs, labs, cols, tcol):
        ax.add_patch(FancyBboxPatch((x, 0.62), 0.18, 0.26,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc=c, ec=c, zorder=2))
        ax.text(x + 0.09, 0.75, l, ha="center", va="center", fontsize=15,
                color=tc, fontweight="bold", zorder=3, linespacing=1.5)
    for x in xs[:2]:
        ax.annotate("", xy=(x + 0.225, 0.75), xytext=(x + 0.185, 0.75),
                    arrowprops=dict(arrowstyle="-|>", lw=2, color=GREY))
    ax.text(0.88, 0.75, "自動化\n在這裡停了", ha="center", va="center",
            fontsize=14, color=RED, fontweight="bold", linespacing=1.5)

    # 兩單價對照
    ax.text(0.10, 0.40, "例行批次複核", fontsize=14, color=GREY, va="center")
    ax.text(0.10, 0.28, "約 20 秒 / 筆", fontsize=22, color=NAVY,
            fontweight="bold", va="center")
    ax.text(0.52, 0.40, "客戶爭議深入判讀", fontsize=14, color=GREY, va="center")
    ax.text(0.52, 0.28, "約 5.25 分鐘 / 筆", fontsize=22, color=ORANGE,
            fontweight="bold", va="center")
    ax.plot([0.46, 0.46], [0.20, 0.46], color=LIGHT, lw=2)
    ax.text(0.5, 0.08, "差 16 倍 —— 差在要不要離開畫面去找證據",
            ha="center", fontsize=17, color=INK, fontweight="bold")
    return save(fig, "s01_pipeline")


def f_timeline():
    """P2:五段時間條,AI 可處理段用深藍。"""
    fig, ax = newfig()
    segs = [("等待排隊", 0.5, LIGHT), ("撈數據", 0.5, NAVY),
            ("判讀", 2.0, NAVY), ("內部覆核", 0.5, LIGHT), ("回信", 0.5, LIGHT)]
    total = sum(s[1] for s in segs)
    x, y, bh = 0.04, 0.52, 0.20
    for name, wgt, c in segs:
        w = wgt / total * 0.92
        ax.add_patch(Rectangle((x, y), w, bh, fc=c, ec="white", lw=2, zorder=2))
        ax.text(x + w / 2, y + bh / 2, name, ha="center", va="center", fontsize=14,
                color=("white" if c == NAVY else INK), fontweight="bold", zorder=3)
        x += w
    ax.text(0.04, 0.82, "一件客戶爭議  2–5 天", fontsize=20, color=INK,
            fontweight="bold", va="center")
    # 括弧標出 AI 可處理段
    x0 = 0.04 + (0.5 / total) * 0.92
    x1 = x0 + (2.5 / total) * 0.92
    ax.plot([x0, x0, x1, x1], [0.46, 0.40, 0.40, 0.46], color=NAVY, lw=2.5)
    ax.text((x0 + x1) / 2, 0.30, "模型幫得上的那兩段", ha="center", fontsize=15,
            color=NAVY, fontweight="bold")
    ax.text((x0 + x1) / 2, 0.17, "43% – 78%", ha="center", fontsize=30,
            color=NAVY, fontweight="bold")
    ax.text(0.96, 0.05, "營運端估計,不是量測", ha="right", fontsize=13, color=GREY)
    return save(fig, "s02_timeline")


def f_twogaps():
    """P3:兩種病,兩個缺口。"""
    fig, ax = newfig()
    for i, (t, sub, big) in enumerate([
            ("行車輔助", "缺的是脈絡", "判定依據\n不在畫面裡"),
            ("駕駛監控", "缺的是判準", "同一段影片\n不同人不同標籤")]):
        x = 0.06 + i * 0.48
        ax.add_patch(FancyBboxPatch((x, 0.30), 0.42, 0.55,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc="white", ec=NAVY, lw=2.2, zorder=2))
        ax.text(x + 0.21, 0.78, t, ha="center", fontsize=19, color=NAVY,
                fontweight="bold", zorder=3)
        ax.text(x + 0.21, 0.67, sub, ha="center", fontsize=15, color=ORANGE,
                fontweight="bold", zorder=3)
        ax.text(x + 0.21, 0.48, big, ha="center", va="center", fontsize=16,
                color=INK, zorder=3, linespacing=1.7)
    ax.text(0.5, 0.13, "交通標誌辨識  =  83% 的誤報,同一個根因",
            ha="center", fontsize=18, color=INK, fontweight="bold")
    ax.text(0.5, 0.03, "而它需要的速限資料,公司本來就有", ha="center",
            fontsize=14, color=GREY)
    return save(fig, "s03_gaps")


def f_blindspot():
    """P4:誤報看得見,漏放看不見。"""
    fig, ax = newfig()
    ax.add_patch(FancyBboxPatch((0.06, 0.34), 0.40, 0.50,
                                boxstyle="round,pad=0,rounding_size=.02",
                                fc=NAVY, ec=NAVY, zorder=2))
    ax.text(0.26, 0.68, "誤報", ha="center", fontsize=22, color="white",
            fontweight="bold", zorder=3)
    ax.text(0.26, 0.52, "量得到", ha="center", fontsize=17, color="white", zorder=3)
    ax.text(0.26, 0.42, "客戶會回頭抱怨", ha="center", fontsize=13,
            color="#C7D6E5", zorder=3)

    ax.add_patch(FancyBboxPatch((0.54, 0.34), 0.40, 0.50,
                                boxstyle="round,pad=0,rounding_size=.02",
                                fc="white", ec=RED, lw=2.5, ls="--", zorder=2))
    ax.text(0.74, 0.68, "漏放", ha="center", fontsize=22, color=RED,
            fontweight="bold", zorder=3)
    ax.text(0.74, 0.52, "目前不存在", ha="center", fontsize=17, color=RED, zorder=3)
    ax.text(0.74, 0.42, "沒有人會來告訴我們", ha="center", fontsize=13,
            color=GREY, zorder=3)
    ax.text(0.5, 0.17, "唯一一份數字:39 筆人工觸發測試,我們自己註明那是下限",
            ha="center", fontsize=15, color=INK)
    ax.text(0.5, 0.05, "沒有人回看的那一流,錯誤會一直留著,還會被當成正確答案餵回模型",
            ha="center", fontsize=14, color=GREY)
    return save(fig, "s04_blind")


def f_options():
    """P5:五條路,一條被選中。"""
    fig, ax = newfig()
    opts = [("加人", "拒絕", RED), ("調班", "待 G1 判", ORANGE),
            ("規則式", "先做", NAVY), ("地端視覺\n語言模型", "採用", NAVY),
            ("既有雲端\n規劃", "前置條件", GREY)]
    wbox, gap = 0.172, 0.035
    x = (1 - (5 * wbox + 4 * gap)) / 2
    for i, (n, badge, c) in enumerate(opts):
        chosen = (badge == "採用")
        ax.add_patch(FancyBboxPatch((x, 0.42), wbox, 0.44,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc=("#EAF1F8" if chosen else "white"),
                                    ec=c, lw=(3.2 if chosen else 1.6), zorder=2))
        ax.text(x + wbox / 2, 0.70, n, ha="center", va="center", fontsize=16,
                color=INK, fontweight="bold", zorder=3, linespacing=1.5)
        ax.add_patch(FancyBboxPatch((x + wbox / 2 - 0.052, 0.475), 0.104, 0.072,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc=c, ec=c, zorder=3))
        ax.text(x + wbox / 2, 0.511, badge, ha="center", va="center", fontsize=13,
                color="white", fontweight="bold", zorder=4)
        x += wbox + gap
    ax.text(0.5, 0.28, "預期 85% 的爭議件,機器有把握、可直接結案",
            ha="center", fontsize=17, color=INK, fontweight="bold")
    ax.text(0.5, 0.185, "推估,非實測 —— 而人現在判得多準,我們沒量過",
            ha="center", fontsize=13, color=GREY)
    ax.add_patch(FancyBboxPatch((0.19, 0.02), 0.62, 0.115,
                                boxstyle="round,pad=0,rounding_size=.02",
                                fc=RED, ec=RED, zorder=2))
    ax.text(0.5, 0.077, "本案六個月:對外自動結案  0%", ha="center", va="center",
            fontsize=19, color="white", fontweight="bold", zorder=3)
    return save(fig, "s05_options")


def f_risks():
    """P6:十條風險,三條紅。"""
    fig, ax = newfig()
    rows = [("R1", "模型不可解釋", False), ("R2", "資料級聯", False),
            ("R3", "自動化偏差", False), ("R4", "駕駛沒有申訴管道", True),
            ("R5", "角色改變抵觸", False), ("R6", "回本難證", False),
            ("R7", "缺高層支持", False), ("R8", "誤報紀錄的舉證地位", True),
            ("R9", "自動結案的漏放", True), ("R10", "自動結案的冤枉", True)]
    for i, (rid, name, hot) in enumerate(rows):
        col, row = i // 5, i % 5
        x = 0.05 + col * 0.49
        y = 0.80 - row * 0.155
        c = RED if hot else GREY
        ax.add_patch(Circle((x + 0.028, y), 0.026, fc=c, ec=c, zorder=3,
                            transform=ax.transData))
        ax.text(x + 0.028, y, rid, ha="center", va="center", fontsize=11,
                color="white", fontweight="bold", zorder=4)
        ax.text(x + 0.075, y, name, ha="left", va="center", fontsize=15,
                color=(RED if hot else INK), fontweight=("bold" if hot else "normal"),
                zorder=3)
    ax.text(0.5, 0.045, "R9 落在安全,R10 落在個人權益 —— 不同量級,不合併成一條",
            ha="center", fontsize=14, color=RED, fontweight="bold")
    return save(fig, "s06_risks")


def f_ethics():
    """P7:八原則,專案期 vs 移交後兩欄。"""
    fig, ax = newfig()
    P = [("人類、社會與環境福祉", "A", "R"), ("以人為本的價值觀", "A", "R"),
         ("公平性", "A", "R"), ("隱私保護與安全", "A", "A"),
         ("可靠性與安全性", "A", "R"), ("透明性與解釋性", "A", "R"),
         ("可申訴性", "R", "R"), ("問責性", "A", "R")]
    CM = {"A": ("#E8A33A", "部分"), "R": (RED, "未達成")}
    ax.text(0.62, 0.94, "專案期", ha="center", fontsize=15, color=GREY, fontweight="bold")
    ax.text(0.83, 0.94, "移交後", ha="center", fontsize=15, color=GREY, fontweight="bold")
    for i, (name, a, b) in enumerate(P):
        y = 0.855 - i * 0.104
        if b == "R":
            ax.add_patch(Rectangle((0.03, y - 0.042), 0.94, 0.084,
                                   fc=RED, ec="none", alpha=0.07, zorder=1))
        ax.text(0.06, y, f"{i+1}", ha="center", va="center", fontsize=13,
                color=NAVY, fontweight="bold", zorder=3)
        ax.text(0.10, y, name, ha="left", va="center", fontsize=15, color=INK, zorder=3)
        for cx, st in ((0.62, a), (0.83, b)):
            c, lab = CM[st]
            ax.add_patch(FancyBboxPatch((cx - 0.058, y - 0.032), 0.116, 0.064,
                                        boxstyle="round,pad=0,rounding_size=.02",
                                        fc=c, ec=c, zorder=3))
            ax.text(cx, y, lab, ha="center", va="center", fontsize=12,
                    color="white", fontweight="bold", zorder=4)
    ax.text(0.5, 0.02, "綠燈  0  條 —— 而移交後那七盞紅燈,全部掛在第 24 週那道門上",
            ha="center", fontsize=14, color=INK, fontweight="bold")
    return save(fig, "s07_ethics")


def f_kpi():
    """P8:三層,護欄層是否決型。"""
    fig, ax = newfig()
    bands = [("業務層", ["客戶爭議回覆時效  2–5 天 → 2 天內"], NAVY),
             ("效率層", ["每筆判讀工時 · 誤報率 · 客戶回頭標記率"], "#4E7FA8"),
             ("護欄層  否決型", ["漏放率 · 複核者一致性 · 訓練資料可追溯率 · 漏濾率"], RED)]
    y = 0.80
    for name, items, c in bands:
        ax.add_patch(FancyBboxPatch((0.05, y - 0.17), 0.90, 0.20,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc="white", ec=c, lw=2.4, zorder=2))
        ax.add_patch(Rectangle((0.05, y - 0.17), 0.015, 0.20, fc=c, ec=c, zorder=3))
        ax.text(0.095, y - 0.03, name, ha="left", va="center", fontsize=16,
                color=c, fontweight="bold", zorder=3)
        ax.text(0.095, y - 0.115, items[0], ha="left", va="center", fontsize=15,
                color=INK, zorder=3)
        y -= 0.255
    ax.text(0.5, 0.14, "任一護欄未達標,上面兩層的成績不採計", ha="center",
            fontsize=16, color=RED, fontweight="bold")
    ax.text(0.5, 0.04, "訂不出目標值的,寫明它由哪一道門、第幾週訂出來",
            ha="center", fontsize=14, color=GREY)
    return save(fig, "s08_kpi")


def f_plan():
    """P9:六個月三道門。"""
    fig, ax = newfig()
    ax.plot([0.06, 0.94], [0.62, 0.62], color=LIGHT, lw=6, solid_capstyle="round",
            zorder=1)
    gates = [(0.28, "G1 · 第 6 週", "因果判定門", "不成立就轉排班"),
             (0.58, "G2 · 第 12 週", "四臂比較門", "含影子模式"),
             (0.88, "G3 · 第 24 週", "部署決定門", "")]
    for x, t, n, sub in gates:
        ax.plot([x], [0.62], marker="D", ms=17, color=NAVY, zorder=3)
        ax.text(x, 0.76, t, ha="center", fontsize=15, color=NAVY, fontweight="bold")
        ax.text(x, 0.50, n, ha="center", fontsize=15, color=INK, fontweight="bold")
        if sub:
            ax.text(x, 0.41, sub, ha="center", fontsize=13, color=GREY)
    ax.text(0.06, 0.62, "W1", ha="right", va="center", fontsize=13, color=GREY)
    ax.add_patch(FancyBboxPatch((0.13, 0.10), 0.74, 0.20,
                                boxstyle="round,pad=0,rounding_size=.02",
                                fc="#EAF1F8", ec=NAVY, lw=2, zorder=2))
    ax.text(0.5, 0.20, "第四臂 · 影子模式:機器自己結案,但不對外送",
            ha="center", va="center", fontsize=17, color=NAVY,
            fontweight="bold", zorder=3)
    return save(fig, "s09_plan")


def f_cost():
    """P10:成本三塊 + 一句請求。"""
    fig, ax = newfig()
    blocks = [("人力", "18.2 人月", "約 NT$158–177 萬", ORANGE),
              ("設備", "一台起步", "NT$16–20 萬", NAVY),
              ("三年可推算 TCO", "專案期", "NT$50–118 萬", GREY)]
    wbox, gap = 0.285, 0.055
    x = (1 - (3 * wbox + 2 * gap)) / 2
    for name, mid, val, c in blocks:
        ax.add_patch(FancyBboxPatch((x, 0.46), wbox, 0.42,
                                    boxstyle="round,pad=0,rounding_size=.02",
                                    fc="white", ec=c, lw=2.4, zorder=2))
        ax.text(x + wbox / 2, 0.80, name, ha="center", fontsize=15, color=c,
                fontweight="bold", zorder=3)
        ax.text(x + wbox / 2, 0.68, mid, ha="center", fontsize=14, color=GREY, zorder=3)
        ax.text(x + wbox / 2, 0.56, val, ha="center", fontsize=19, color=INK,
                fontweight="bold", zorder=3)
        x += wbox + gap
    ax.text(0.5, 0.34, "設備不到一個人一年成本的五分之一 —— 成本主體是人",
            ha="center", fontsize=16, color=INK, fontweight="bold")
    ax.text(0.5, 0.23, "而自動化也不是免費的:移交後稽核會吃掉毛節省的 3%–52%",
            ha="center", fontsize=14, color=ORANGE)
    ax.add_patch(FancyBboxPatch((0.10, 0.02), 0.80, 0.145,
                                boxstyle="round,pad=0,rounding_size=.02",
                                fc=NAVY, ec=NAVY, zorder=2))
    ax.text(0.5, 0.092, "請求:核准第一段 W1–6,並指定一位贊助人", ha="center",
            va="center", fontsize=19, color="white", fontweight="bold", zorder=3)
    return save(fig, "s10_cost")


# ══════════════════════════════════════════════════════════════════
# 逐頁內容 —— 標題 / 要點(≤3) / 圖 / 底線一句
# ══════════════════════════════════════════════════════════════════
PAGES = [
    dict(t=None, hero="五個人。\n三千支影片。\n兩天。", fig=None, foot=None),
    dict(t="用例:兩階段推論,與它下游那道人工防線",
         b=[], fig=f_pipeline, foot=None),
    dict(t="為什麼是現在:工作量沒變多,客戶要的時效變了",
         b=["現況約半數爭議超過兩天", "客戶要的是兩天內、單純案件當天"],
         fig=f_timeline, foot=None),
    dict(t="發現過程:兩種病,兩個缺口", b=[], fig=f_twogaps, foot=None),
    dict(t="我們看不見的那一半:有資料,沒有習慣", b=[], fig=f_blindspot, foot=None),
    dict(t="五條路,含三條不用 AI", b=[], fig=f_options, foot=None),
    dict(t="風險:七類十條", b=[], fig=f_risks,
         foot="R8 曾判誤報的紀錄在事故舉證中的地位 —— 這題我沒有答案,由法務在 G1 前給書面意見"),
    dict(t="倫理:對照澳洲 AI 倫理八原則(2019)", b=[], fig=f_ethics, foot=None),
    dict(t="怎麼算成功:三層 KPI,一條回圈", b=[], fig=f_kpi, foot=None),
    dict(t="六個月 · 六個工作包 · 三道決策門", b=[], fig=f_plan,
         foot="要不要真的對外自動結案,不在這六個月決定"),
    dict(t="資源、當責與請求", b=[], fig=f_cost, foot=None),
    dict(t="參考文獻", b=[], fig=None, foot=None, refs=True),
]


# ══════════════════════════════════════════════════════════════════
def narration():
    """從 notes/_v2_parts 抓每頁的口白逐字,灌進講者備註。"""
    import glob
    H1 = re.compile(r"(?m)^#\s+(.+)$")
    NARR = re.compile(r"(?m)^##\s*口白逐字\s*$")
    NXT = re.compile(r"(?m)^#{1,2}\s+")
    STAGE = re.compile(r"(?ms)^[ \t]*[*_]{0,2}[(（].*?[)）][*_]{0,2}[ \t]*$")
    out = []
    for f in sorted(glob.glob(os.path.join(A3, "notes", "_v2_parts", "*.md"))):
        t = open(f, encoding="utf-8").read()
        hs = list(H1.finditer(t))
        for i, m in enumerate(hs):
            if i == 0:
                continue
            e = hs[i + 1].start() if i + 1 < len(hs) else len(t)
            u = t[m.start():e]
            nm = NARR.search(u)
            if not nm:
                out.append((m.group(1).strip(), ""))
                continue
            rest = u[nm.end():]
            nx = NXT.search(rest)
            blk = rest[:nx.start()] if nx else rest
            blk = STAGE.sub("", blk)
            body = "\n".join(l for l in blk.split("\n")
                             if not l.lstrip().startswith(">")).strip()
            out.append((m.group(1).strip(), re.sub(r"[*_`~]", "", body)))
    return out


def txt(slide, x, y, w, h, s, size, color, bold=False, align="left",
        anchor="top", spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE}[anchor]
    for i, line in enumerate(str(s).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}[align]
        p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = _P(color); r.font.name = FONT
    return tb


REFS = [
    "Australian Human Rights Commission. (2020). Using artificial intelligence to make decisions.",
    "Department of Industry, Science and Resources. (2019). Australia's AI ethics principles.",
    "Ensign, D., Friedler, S. A., Neville, S., Scheidegger, C., & Venkatasubramanian, S. (2018). Runaway feedback loops in predictive policing. PMLR, 81, 160–171.",
    "Goddard, K., Roudsari, A., & Wyatt, J. C. (2012). Automation bias. JAMIA, 19(1), 121–127.",
    "Guo, C., Pleiss, G., Sun, Y., & Weinberger, K. Q. (2017). On calibration of modern neural networks. PMLR, 70, 1321–1330.",
    "National Institute of Standards and Technology. (2023). AI risk management framework (AI RMF 1.0).",
    "Northcutt, C. G., Athalye, A., & Mueller, J. (2021). Pervasive label errors in test sets. NeurIPS Datasets and Benchmarks.",
    "Parasuraman, R., & Manzey, D. H. (2010). Complacency and bias in human use of automation. Human Factors, 52(3), 381–410.",
    "Rudin, C. (2019). Stop explaining black box machine learning models. Nature Machine Intelligence, 1(5), 206–215.",
    "Sambasivan, N., Kapania, S., Highfill, H., Akrong, D., Paritosh, P., & Aroyo, L. (2021). Data cascades in high-stakes AI. CHI '21.",
    "Sculley, D., Holt, G., Golovin, D., et al. (2015). Hidden technical debt in machine learning systems. NeurIPS, 28.",
]


def build():
    narr = narration()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    for i, pg in enumerate(PAGES):
        s = prs.slides.add_slide(prs.slide_layouts[6])

        if pg.get("hero"):
            txt(s, M, 2.1, W - 2 * M, 3.4, pg["hero"], 54, NAVY, bold=True,
                align="center", anchor="middle", spacing=1.5)
        else:
            txt(s, M, TITLE_Y, W - M - 0.6, 0.85, pg["t"], 28, NAVY, bold=True)
            # 標題下的細分隔線。python-pptx 的新圖形預設是 _NoneFill,
            # 必須先 .solid() 才能設 fore_color,否則丟 TypeError。
            rule = s.shapes.add_shape(1, Inches(M), Inches(TITLE_Y + 0.92),
                                      Inches(W - 2 * M), Inches(0.035))
            rule.fill.solid()
            rule.fill.fore_color.rgb = _P(NAVY)
            rule.line.fill.background()
            rule.shadow.inherit = False

            y = BODY_Y
            for b in pg.get("b", []):
                txt(s, M, y, W - 2 * M, 0.45, "・" + b, 18, INK)
                y += 0.48

            if pg.get("fig"):
                p = pg["fig"]()
                from PIL import Image
                iw, ih = Image.open(p).size
                # 有要點的頁,圖往下讓一行;有底線的頁,圖縮 0.45 吋
                fh = FIG_H - (0.60 if pg.get("b") else 0) - (0.45 if pg.get("foot") else 0)
                fy = FIG_Y + (0.60 if pg.get("b") else 0)
                fw = fh * iw / ih
                if fw > W - 2 * M:
                    fw = W - 2 * M; fh = fw * ih / iw
                s.shapes.add_picture(p, Inches((W - fw) / 2), Inches(fy),
                                     Inches(fw), Inches(fh))

            if pg.get("refs"):
                y = BODY_Y
                for r in REFS:
                    txt(s, M, y, W - 2 * M - SAFE_R + 1.0, 0.42, r, 12, INK, spacing=1.15)
                    y += 0.43

            if pg.get("foot"):
                txt(s, M, H - 1.02, W - 2 * M - SAFE_R, 0.5, pg["foot"], 15, GREY)

        if i < len(narr) and narr[i][1]:
            s.notes_slide.notes_text_frame.text = narr[i][1]

    prs.save(OUT)
    print(f"✅ {os.path.basename(OUT)}")
    print(f"   {len(prs.slides)} 頁 · {os.path.getsize(OUT)/2**20:.1f} MB")
    print(f"   圖:{FIGDIR}")
    return OUT


if __name__ == "__main__":
    sys.exit(0 if build() else 1)
