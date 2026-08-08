# -*- coding: utf-8 -*-
"""A3 v4 投影片產生器 —— **骨架**。

結構照抄 v3 的 `Assessment3/build_deck_v3.py`:
本檔只做內容層,版面 helper 一律從上一層的 `build_deck.py` import,不重寫版面。

v4 = **12 張**:封面 + 10 頁內容 + 1 頁參考文獻(v3 是 14 張,超了)。
五段導軌沿用 build_deck.SECTIONS:介紹 / AI 機會 / 業務影響 / 效益與衡量 / 項目計畫。
右下淨空區 SAFE_R 保留(錄影時放人像小窗),picture() 預設會自動避開。

**每頁的文字與配圖都待 v2.1 定稿後才填,現在不得自行編造。**
"""
import os
import sys

# ------------------------------------------------------------------
# 路徑:本檔在 Assessment3/v4/scripts/,版面庫在 Assessment3/build_deck.py
# ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
V4 = os.path.dirname(HERE)                       # .../Assessment3/v4
A3 = os.path.dirname(V4)                         # .../Assessment3
if A3 not in sys.path:
    sys.path.insert(0, A3)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import build_deck
from build_deck import (txt, header, picture, note, card,
                        NAVY, RED, GREY, DARK, WHITE, LIGHT,
                        W, H, TOP, BOT, SAFE_R, SECTIONS)

# 圖檔目錄改指 v4/figures。
# picture() 讀的是 build_deck 模組層的 FIG,所以必須在這裡覆寫模組屬性,
# 只改本檔的區域變數是沒有用的。
FIG = os.path.join(V4, "figures")
build_deck.FIG = FIG

# 交件檔輸出到 v4/ 之下,**不覆蓋上一層 Assessment3/ 裡的 v3 同名檔**。
OUT = os.path.join(V4, "Huang_26254793_421104_Assessment 3.pptx")

N_SLIDES = 12          # 封面 1 + 內容 10 + 參考文獻 1,硬約束
N_CONTENT = 10

# 每頁對應的導軌段(build_deck.SECTIONS 的索引;None = 不顯示導軌)。
# TODO: 十頁對到哪一段,待 v2.1 定稿後填;現在一律 None。
PAGE_SECTION = [None] * N_CONTENT

# 摘要 / 結論頁的四卡橫排參數(沿用 v3 實測值)
CARD_W, CARD_GAP, CARD_X0 = 2.62, 0.16, 0.55


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _todo(slide, tag):
    """骨架期的佔位戳記。內容填進去之後,這行要刪掉。"""
    txt(slide, 0.55, TOP + 0.30, W - 1.1 - SAFE_R, 0.8,
        f"[{tag}] 內容待 v2.1 定稿後填", size=18, color=GREY)


# ==================================================================
# 封面
# ==================================================================
def slide_00_cover(prs):
    """封面。內容待 v2.1 定稿後填。

    v3 樣式參考(定稿後照這個骨架填字,不要現在編):
        txt(s, 1.0, 2.0,  11.3, 1.5, <主標>,  size=38, color=NAVY, bold=True)
        txt(s, 1.0, 3.62, 11.3, 0.6, <副標>,  size=20, color=GREY)
        txt(s, 1.0, 4.62, 11.3, 1.4, <姓名/學號/科目/日期>, size=15, color=DARK)
    """
    s = _blank(prs)
    _todo(s, "封面")
    return s


# ==================================================================
# 10 頁內容
# 版型二選一:
#   A 內容頁    header(...) + picture("fig_v4_XX.png", top=TOP+0.10, maxh=5.05) + note(...)
#   B 摘要/結論 header(...) + 四張 card(CARD_X0 + i*(CARD_W+CARD_GAP), ...) + 一段大字 txt(...)
# ==================================================================
def slide_01(prs):
    """P1。內容待 v2.1 定稿後填(配圖 fig_v4_01)。"""
    s = _blank(prs)
    _todo(s, "P1")
    return s


def slide_02(prs):
    """P2。內容待 v2.1 定稿後填(配圖 fig_v4_02)。"""
    s = _blank(prs)
    _todo(s, "P2")
    return s


def slide_03(prs):
    """P3。內容待 v2.1 定稿後填(配圖 fig_v4_03)。"""
    s = _blank(prs)
    _todo(s, "P3")
    return s


def slide_04(prs):
    """P4。內容待 v2.1 定稿後填(配圖 fig_v4_04)。"""
    s = _blank(prs)
    _todo(s, "P4")
    return s


def slide_05(prs):
    """P5。內容待 v2.1 定稿後填(配圖 fig_v4_05)。"""
    s = _blank(prs)
    _todo(s, "P5")
    return s


def slide_06(prs):
    """P6。內容待 v2.1 定稿後填(配圖 fig_v4_06)。"""
    s = _blank(prs)
    _todo(s, "P6")
    return s


def slide_07(prs):
    """P7。內容待 v2.1 定稿後填(配圖 fig_v4_07)。"""
    s = _blank(prs)
    _todo(s, "P7")
    return s


def slide_08(prs):
    """P8。內容待 v2.1 定稿後填(配圖 fig_v4_08)。"""
    s = _blank(prs)
    _todo(s, "P8")
    return s


def slide_09(prs):
    """P9。內容待 v2.1 定稿後填(配圖 fig_v4_09)。"""
    s = _blank(prs)
    _todo(s, "P9")
    return s


def slide_10(prs):
    """P10。內容待 v2.1 定稿後填(配圖 fig_v4_10)。"""
    s = _blank(prs)
    _todo(s, "P10")
    return s


# ==================================================================
# 參考文獻
# ==================================================================
def slide_11_references(prs):
    """參考文獻頁。內容待 v2.1 定稿後填。

    版型:header(s, "參考文獻", None) + 一大塊 txt(..., size=13, spacing=1.5)。
    三條機械判準(沿用):每筆可點擊深層連結 / 零重複條目 / 零孤兒條目。
    APA 7,字級不小於 12pt。
    """
    s = _blank(prs)
    _todo(s, "參考文獻")
    return s


PAGES = [
    slide_00_cover,
    slide_01, slide_02, slide_03, slide_04, slide_05,
    slide_06, slide_07, slide_08, slide_09, slide_10,
    slide_11_references,
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    for page in PAGES:
        page(prs)

    n = len(prs.slides._sldIdLst)
    assert n == N_SLIDES, f"v4 必須是 {N_SLIDES} 張,目前 {n} 張"

    prs.save(OUT)
    return prs, OUT


if __name__ == "__main__":
    prs, out = build()
    print("投影片頁數:", len(prs.slides._sldIdLst))
    print("圖檔目錄:", build_deck.FIG)
    print("輸出:", out)
    print("大小: %.1f KB" % (os.path.getsize(out) / 1024))

# ==================================================================
# ⚠️ 本檔目前只有結構,12 頁全為空殼 + 佔位戳記。
#    內容待 v2.1 定稿後填;填完後把每頁的 _todo(...) 刪掉。
#    出圖與組版都完成後,必須匯出逐頁圖目視稽核(v3 有三張圖是講稿定稿時才發現破版)。
# ==================================================================
