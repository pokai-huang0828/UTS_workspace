# -*- coding: utf-8 -*-
"""A3 投影片產生器 —— 依 notes/A3_提案內容_v2.md 與 A3_講稿與答辯_v1.md。

版面硬規則(來自 v2 §附):
- 16:9;右下角 15% 面積為淨空區(PowerPoint 錄製的人像小窗固定在那)
- 字體 Microsoft JhengHei;主色深藍 / 強調紅只用於風險與缺口 / 中性灰
- 每頁最多一個 kicker 數字(>=80pt),其餘文字 <=20pt
- 每頁右上五點章節導軌(對應官方五段模板),當前段實心
重跑即覆蓋輸出檔。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "Huang_26254793_421104_Assessment 3.pptx")

NAVY = RGBColor(0x1F, 0x38, 0x64)
RED = RGBColor(0xC0, 0x00, 0x00)
GREY = RGBColor(0x80, 0x80, 0x80)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD9, 0xD9, 0xD9)
FONT = "Microsoft JhengHei"

W, H = 13.333, 7.5                 # 16:9
SAFE_R, SAFE_B = 2.05, 1.15        # 右下淨空區(人像小窗)
TOP = 1.05                         # 標題下緣
BOT = 0.70                         # 底部註記區

SECTIONS = ["介紹", "AI 機會", "業務影響", "效益與衡量", "項目計畫"]


def txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def header(slide, title, section_idx, sub=None):
    txt(slide, 0.55, 0.30, 9.6, 0.7,
        (f"{SECTIONS[section_idx]} · " if section_idx is not None else "") + title,
        size=26, color=NAVY, bold=True)
    if sub:
        txt(slide, 0.55, 0.92, 9.6, 0.4, sub, size=12, color=GREY)
    # 章節導軌
    if section_idx is not None:
        from pptx.enum.shapes import MSO_SHAPE
        for i in range(5):
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(11.05 + i * 0.36), Inches(0.42),
                                       Inches(0.17), Inches(0.17))
            d.fill.solid()
            d.fill.fore_color.rgb = NAVY if i == section_idx else LIGHT
            d.line.fill.background()
            d.shadow.inherit = False


def picture(slide, name, top=TOP + 0.25, avoid_safe=True, maxw=None, maxh=None,
            left=None):
    """等比置入圖片,避開右下淨空區。"""
    path = os.path.join(FIG, name)
    if not os.path.exists(path):
        print("  !! 缺圖:", name)
        return None
    iw, ih = Image.open(path).size
    mw = maxw if maxw else (W - 1.1 - (SAFE_R if avoid_safe else 0))
    mh = maxh if maxh else (H - top - BOT)
    s = min(mw / iw, mh / ih)
    w, h = iw * s, ih * s
    l = left if left is not None else (W - (SAFE_R if avoid_safe else 0) - w) / 2 + 0.15
    slide.shapes.add_picture(path, Inches(l), Inches(top), Inches(w), Inches(h))
    return h


def note(slide, text, color=GREY, size=11):
    txt(slide, 0.55, H - BOT + 0.02, W - 1.1 - SAFE_R, 0.55, text, size=size, color=color)


def card(slide, x, y, w, h, title, body, col=NAVY, tsize=13, bsize=11):
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(0.42))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = col
    hdr.line.fill.background(); hdr.shadow.inherit = False
    txt(slide, x, y + 0.05, w, 0.35, title, size=tsize, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    txt(slide, x + 0.14, y + 0.52, w - 0.28, h - 0.62, body, size=bsize, color=DARK,
        spacing=1.35)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    blank = prs.slide_layouts[6]
    S = lambda: prs.slides.add_slide(blank)

    # ---------- 封面 ----------
    s = S()
    txt(s, 1.0, 2.15, 11.3, 1.3,
        "CBA 房貸固定利率到期留存決策系統", size=40, color=NAVY, bold=True)
    txt(s, 1.0, 3.45, 11.3, 0.6,
        "AI 路線圖提案  ·  提交技術審查委員會", size=20, color=GREY)
    txt(s, 1.0, 4.5, 11.3, 1.4,
        "Po-Kai Huang(26254793)\n421104 Artificial Intelligence for Enterprises\n"
        "Assessment 3  ·  2026 年 8 月", size=15, color=DARK, spacing=1.5)

    # ---------- P1 Executive Summary ----------
    s = S(); header(s, "Executive Summary — Decision Requested", 0)
    cw, gap = 2.62, 0.16
    cards = [
        ("批准什麼", "房貸固定利率到期\n留存決策系統\n\n12 個月\n約 27.5 人月(規劃值)\n\n分兩期撥款,\n今天只請求第一期", NAVY),
        ("為什麼是現在", "CDR 已把轉換成本\n降到接近一次授權點擊\n\n未來 12 個月是\n固定利率到期的\n高峰窗口\n\n窗口過了,\n客戶已經走了", RED),
        ("對到多大的池子", "房貸帳 A$611.5bn\n(CBA, 2025b)\n\n約為 FY25 現金淨利\nA$10,252m 的 60 倍\n\n本案針對每年\n滾動到期的那一段", NAVY),
        ("何時要決定", "本次委員會\n\n下一個決策點:\n第 6 個月\n流程就緒關卡\n\n第二期撥款以\n該關卡通過為條件", NAVY),
    ]
    for i, (t, b, c) in enumerate(cards):
        card(s, 0.55 + i * (cw + gap), TOP + 0.35, cw, 4.35, t, b, c)
    txt(s, 0.55, TOP + 4.95, W - 1.1 - SAFE_R, 0.8,
        "向一家在 Evident AI Index 2025 全球排名第 4、亞太第 1 的機構提案(CBA, 2025a)——\n"
        "本案的問題不是能不能做,是該不該把這筆錢放在這裡。",
        size=14, color=NAVY, bold=True, spacing=1.4)

    # ---------- P2 威脅 ----------
    s = S(); header(s, "威脅:CDR 打開的缺口", 0)
    picture(s, "fig02_market.png", top=TOP + 0.15, maxh=4.05)
    txt(s, 0.55, TOP + 4.35, W - 1.1 - SAFE_R, 1.5,
        "A1 已經證明 CBA 用 AI 贏在拿到新客戶的速度 —— 房貸有條件核准 <10 分鐘、業務信貸年審 14h→2h。\n"
        "但 CDR 打的是既有客戶離開的成本。Acquisition 與 retention 不是同一套能力。",
        size=14, color=NAVY, bold=True, spacing=1.4)
    note(s, "在 CBA 公開揭露中未見針對房貸存量留存的 AI 決策系統(CBA, 2024; 2025c);"
            "若 CEE 內部已涵蓋部分留存場景,本案定位即調整為在既有 CEE 上增設房貸再融資專用模組,投資論證不變。")

    # ---------- P3 選案 ----------
    s = S(); header(s, "選案:三案加權比較與敏感度", 1)
    picture(s, "fig03_matrix.png", top=TOP + 0.10, maxh=2.95, maxw=7.0, left=0.55)
    picture(s, "fig03_tornado.png", top=TOP + 3.20, maxh=2.55, maxw=6.9, left=0.55)
    card(s, 7.85, TOP + 0.10, 3.35, 5.65, "反方意見與回應",
         "「CBA 是 Evident 全球第 4,\n怎麼會有空白?」\n"
         "→ AI 成熟度衡量組織能力與\n部署廣度,不等於每條業務線\n都有對應系統;成熟度高反而\n降低執行風險。\n\n"
         "「CEE 不就是在留客?」\n"
         "→ CEE 是跨產品次佳互動引擎;\n本案要的是固定利率到期事件上\n的流失評分與外撥排序。\n\n"
         "「空白可能是理性的」\n"
         "→ 承認兩項限制:競品訊號\n不可觀測、差別定價敏感。\n兩者都已寫進限制頁與倫理護欄。",
         col=GREY, bsize=10)

    # ---------- P4 策略與閉環 ----------
    s = S(); header(s, "三條留存策略與系統閉環", 1)
    picture(s, "fig04_loop.png", top=TOP + 0.10, maxh=5.0)
    note(s, "會被改變的決策:客戶經理每月的外撥配額排序,從「按貸款餘額高低排」"
            "改為「按流失傾向 × 客戶終身價值排」。", color=NAVY, size=12)

    # ---------- P5 可行性 ----------
    s = S(); header(s, "可行性評估", 1,
                    "本案不需要新建 AI 能力,只需要把既有能力指向新的一軸")
    picture(s, "fig06_feasibility.png", top=TOP + 0.35, maxh=4.35)
    note(s, "go / no-go 門檻:若「資料」維為紅燈,本案順延,先做 6 週資料可行性研究再回委員會。",
         color=RED, size=12)

    # ---------- P6 風險 ----------
    s = S(); header(s, "限制與挑戰:緩解前 → 緩解後", 2)
    picture(s, "fig05_risk_matrix.png", top=TOP + 0.10, maxh=5.05)
    note(s, "每一條殘餘風險,在專案計畫裡都有一個對應的停止或回退關卡。", color=NAVY, size=12)

    # ---------- P7 倫理 ----------
    s = S(); header(s, "社會文化與倫理影響", 2)
    picture(s, "fig07_ethics.png", top=TOP + 0.10, maxh=5.05)
    note(s, "漏掉一個要走的客戶,只是損失利差;把財困客戶當商機,是監管事件。",
         color=RED, size=12)

    # ---------- P8 價值 ----------
    s = S(); header(s, "價值方程式:我不報一個假的 ROI,我報反解", 3)
    picture(s, "fig08_value.png", top=TOP + 0.10, maxh=5.05)
    note(s, "一旦拿到內部到期分佈與產品別 margin,這個式子當場就能算出來。", color=NAVY, size=12)

    # ---------- P9 KPI ----------
    s = S(); header(s, "KPI 樹:證明賺錢,也證明沒有傷到人", 3)
    picture(s, "fig09_kpi.png", top=TOP + 0.10, maxh=5.05)
    note(s, "所有業務層指標以「實驗組 減 對照組」的增量計;絕對值僅供監控,不作為成效證據。",
         color=NAVY, size=12)

    # ---------- P10 時程 ----------
    s = S(); header(s, "實施時程與六道關卡", 4)
    picture(s, "fig09_gantt.png", top=TOP + 0.05, maxh=2.85, maxw=11.1)
    picture(s, "fig10_stagegate.png", top=TOP + 3.05, maxh=2.55, maxw=11.1)
    note(s, "每道關卡四態:Go / Conditional-Go / Hold / Kill。試點裁決未達顯著性即停止,不續投。")

    # ---------- P11 資源 ----------
    s = S(); header(s, "資源配置與當責", 4)
    picture(s, "fig11_resource.png", top=TOP + 0.05, maxh=2.85, maxw=6.6, left=0.55)
    picture(s, "fig11_raci.png", top=TOP + 0.05, maxh=2.85, maxw=4.3, left=7.35)
    txt(s, 0.55, TOP + 3.10, W - 1.1 - SAFE_R, 2.3,
        "資源尖峰不在模型開發,在合規審查(第 5 個月 3.25 FTE)——\n"
        "這個系統真正難的不是把模型做出來,是讓它被允許用。\n"
        "而試點期是全期資源最低(1.42 FTE),因為那三個月是在等實驗結果,不是在建東西。\n\n"
        "第一期(M1–M6)合計 16.0 人月,與「今天只請求批准第一期」精準對應。",
        size=13, color=NAVY, spacing=1.45)
    note(s, "Planning estimate — to be calibrated with CBA internal rates。"
            "成本 = 總人月 × 澳洲 ICT 職類薪資中位數區間 × 1.3 間接費率(間接費率為提案人假設)。")

    # ---------- P12 結論 ----------
    s = S(); header(s, "結論與批准請求", None)
    cards2 = [
        ("批准什麼", "12 個月、約 27.5 人月\n分兩期撥款\n今天只請求第一期\n(M1–M6,16.0 人月)", NAVY),
        ("為什麼值得", "唯一對應到\n尚無 AI 覆蓋的\n戰略缺口的專案", NAVY),
        ("對到多大的池子", "房貸帳 A$611.5bn\n約 FY25 現金淨利的 60 倍\n(同第一頁,同一算式)", NAVY),
        ("已知限制", "方法在公開資料集驗證\n內部基準尚未校準\n與 CEE 的功能重疊待查", RED),
    ]
    cw2 = 2.62
    for i, (t, b, c) in enumerate(cards2):
        card(s, 0.55 + i * (cw2 + gap), TOP + 0.45, cw2, 3.5, t, b, c)
    txt(s, 0.55, TOP + 4.30, W - 1.1 - SAFE_R, 1.2,
        "下一個決策點:第 6 個月的流程就緒關卡 —— 第二期撥款以該關卡通過為條件。",
        size=17, color=NAVY, bold=True)

    # ---------- 參考文獻 ----------
    s = S(); header(s, "參考文獻", None)
    txt(s, 0.55, TOP + 0.30, W - 1.1, 5.6,
        "⚠️ 本頁為佔位。交件前必須逐筆重抓深層連結,不得原樣沿用作業一的書目。\n\n"
        "必修四筆:\n"
        "  1. Macquarie 房貸年增 19% / 市占 5.9% → 改引 APRA Monthly ADI Statistics(可下載、含版本日期)\n"
        "  2. 全球 NIM 3.1%→2.7%(S&P Global 首頁連結)→ 整條棄用\n"
        "  3. CBA (2025b) → 改指 FY25 Results Presentation 的 PDF 直連,並註明 NIM 2.08%、房貸 A$611.5bn 的頁碼\n"
        "  4. NAB → 本提案不用四大行對照表,整條刪除\n\n"
        "需新增 2–3 筆外部文獻(選案論證目前的外部文獻數為 0):\n"
        "  · 零售銀行 churn / retention 建模的同儕審查期刊論文\n"
        "  · 開放銀行 / CDR 對消費者轉換行為影響的官方或監管報告\n"
        "  · uplift modeling 在留存行銷的方法學文獻\n\n"
        "三條機械判準:每筆可點擊深層連結 / 零重複條目 / 零孤兒條目。APA 7,字級不小於 12pt。",
        size=13, color=DARK, spacing=1.5)

    prs.save(OUT)
    return prs, OUT


if __name__ == "__main__":
    prs, out = build()
    print("投影片頁數:", len(prs.slides.__iter__.__self__._sldIdLst))
    print("輸出:", out)
    print("大小: %.1f KB" % (os.path.getsize(out) / 1024))
